from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor

# Colors - AstraZeneca Brand
AZ_MAROON = RGBColor(0x83, 0x00, 0x51)
AZ_PINK = RGBColor(0xE9, 0x1E, 0x63)
AZ_LIGHT_PINK = RGBColor(0xF8, 0xE0, 0xEB)
AZ_GOLD = RGBColor(0xC4, 0xA0, 0x00)
AZ_GREEN = RGBColor(0x2E, 0x7D, 0x32)
AZ_BLUE = RGBColor(0x1E, 0x88, 0xE5)
AZ_PURPLE = RGBColor(0x5E, 0x35, 0xB1)
AZ_DARK_PURPLE = RGBColor(0x4A, 0x14, 0x8C)
TEXT_DARK = RGBColor(0x2D, 0x2D, 0x2D)
TEXT_MEDIUM = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)

def add_text(slide, txt, left, top, width, height, size=12, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return box

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_agent_box(slide, name, left, top, width=Inches(1.1), height=Inches(0.55), fill=AZ_PINK, text_color=WHITE, size=8):
    """Add an agent box"""
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill, AZ_MAROON, 1)
    add_text(slide, name, left, top + Inches(0.1), width, height - Inches(0.1), size, True, text_color, PP_ALIGN.CENTER)
    return box

def add_platform_box(slide, name, left, top, width=Inches(1.3), height=Inches(0.45), fill=AZ_DARK_PURPLE):
    """Add a platform box"""
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill, AZ_MAROON, 1)
    add_text(slide, name, left, top + Inches(0.08), width, height - Inches(0.08), 7, True, WHITE, PP_ALIGN.CENTER)
    return box

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================================================
# SLIDE 1: Test Intelligence Agent - Semantic Flow
# ============================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])

# Title
add_text(slide1, "Test Intelligence Agent – Semantic Flow", Inches(0.3), Inches(0.15), Inches(10), Inches(0.5), 28, False, AZ_MAROON, PP_ALIGN.LEFT)
add_text(slide1, "Multi-Agent Architecture for Autonomous Testing across R&D Platforms", Inches(0.3), Inches(0.55), Inches(12), Inches(0.3), 11, False, TEXT_MEDIUM, PP_ALIGN.LEFT, italic=True)

# ===== LEFT SIDE: R&D Platforms (9 platforms in 2 groups) =====

# KPI Data Products Section
kpi_bg = add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.2), Inches(1.0), Inches(1.8), Inches(2.8), AZ_LIGHT_PINK, AZ_MAROON, 1)
add_text(slide1, "KPI Data Products", Inches(0.2), Inches(1.05), Inches(1.8), Inches(0.35), 9, True, AZ_MAROON, PP_ALIGN.CENTER)

kpi_platforms = [
    "3DP",
    "BIKG",
    "Patient Safety",
    "Clinical Trials",
]
for i, platform in enumerate(kpi_platforms):
    add_platform_box(slide1, platform, Inches(0.35), Inches(1.45 + i * 0.55), Inches(1.5), Inches(0.45))

# Foundation Data Products Section
foundation_bg = add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.2), Inches(3.9), Inches(1.8), Inches(3.3), AZ_LIGHT_PINK, AZ_MAROON, 1)
add_text(slide1, "Foundation Data\nProducts", Inches(0.2), Inches(3.95), Inches(1.8), Inches(0.5), 9, True, AZ_MAROON, PP_ALIGN.CENTER)

foundation_platforms = [
    "CDISC/ADaM",
    "Regulatory Docs",
    "Study PLDB",
    "GxP Systems",
    "HPC Environ.",
]
for i, platform in enumerate(foundation_platforms):
    add_platform_box(slide1, platform, Inches(0.35), Inches(4.5 + i * 0.55), Inches(1.5), Inches(0.45))

# ===== CENTER: Test Intelligence Agent Orchestration =====

# Main orchestrator in the center-top
add_agent_box(slide1, "Orchestrator\nAgent", Inches(5.8), Inches(1.0), Inches(1.4), Inches(0.7), AZ_PINK, WHITE, 10)

# User icon
add_text(slide1, "👤", Inches(8.5), Inches(1.0), Inches(0.5), Inches(0.5), 24, False, TEXT_DARK, PP_ALIGN.CENTER)
add_text(slide1, "User", Inches(8.3), Inches(1.45), Inches(0.9), Inches(0.3), 9, False, TEXT_DARK, PP_ALIGN.CENTER)

# Row 1: Requirement & Strategy Agents
add_agent_box(slide1, "Requirement\nGathering", Inches(3.5), Inches(2.0), Inches(1.2), Inches(0.65), AZ_PINK, WHITE, 8)
add_agent_box(slide1, "Test Strategy\nAgent", Inches(5.0), Inches(2.0), Inches(1.2), Inches(0.65), AZ_PINK, WHITE, 8)
add_agent_box(slide1, "Knowledge &\nMemory Agent", Inches(6.5), Inches(2.0), Inches(1.2), Inches(0.65), AZ_DARK_PURPLE, WHITE, 8)

# Row 2: Generation Agents
add_agent_box(slide1, "Test Case\nGeneration", Inches(3.5), Inches(3.0), Inches(1.2), Inches(0.65), AZ_PINK, WHITE, 8)
add_agent_box(slide1, "Test Script\nGeneration", Inches(5.0), Inches(3.0), Inches(1.2), Inches(0.65), AZ_PINK, WHITE, 8)
add_agent_box(slide1, "Synthetic Data\nGeneration", Inches(6.5), Inches(3.0), Inches(1.2), Inches(0.65), AZ_PINK, WHITE, 8)

# Row 3: Execution & Analysis Agents
add_agent_box(slide1, "Test Execution\nAgent", Inches(3.5), Inches(4.0), Inches(1.2), Inches(0.65), AZ_PINK, WHITE, 8)
add_agent_box(slide1, "Failure Analysis\nAgent", Inches(5.0), Inches(4.0), Inches(1.2), Inches(0.65), AZ_PINK, WHITE, 8)
add_agent_box(slide1, "Code Refactor\nAgent", Inches(6.5), Inches(4.0), Inches(1.2), Inches(0.65), AZ_PINK, WHITE, 8)

# Row 4: Reporting Agent (center bottom)
add_agent_box(slide1, "Reporting\nAgent", Inches(5.0), Inches(5.0), Inches(1.2), Inches(0.65), AZ_PINK, WHITE, 8)

# ===== RIGHT SIDE: Data Infrastructure =====

# GitHub Ontology
add_text(slide1, "GitHub", Inches(8.8), Inches(2.5), Inches(1.2), Inches(0.3), 10, True, TEXT_DARK, PP_ALIGN.CENTER)
add_text(slide1, "Ontology", Inches(8.8), Inches(2.75), Inches(1.2), Inches(0.25), 9, False, TEXT_MEDIUM, PP_ALIGN.CENTER)

# Test Repository
test_repo = add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(3.3), Inches(1.6), Inches(0.6), AZ_DARK_PURPLE, AZ_MAROON, 1)
add_text(slide1, "Test\nRepository", Inches(8.5), Inches(3.35), Inches(1.6), Inches(0.55), 9, True, WHITE, PP_ALIGN.CENTER)

# Evidence Store
evidence = add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(4.1), Inches(1.6), Inches(0.6), AZ_DARK_PURPLE, AZ_MAROON, 1)
add_text(slide1, "Evidence\nStore", Inches(8.5), Inches(4.15), Inches(1.6), Inches(0.55), 9, True, WHITE, PP_ALIGN.CENTER)

# Knowledge Graph (cylinder shape simulated)
kg = add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(4.9), Inches(1.6), Inches(0.7), AZ_DARK_PURPLE, AZ_MAROON, 2)
add_text(slide1, "Knowledge\nGraph", Inches(8.5), Inches(4.95), Inches(1.6), Inches(0.6), 10, True, WHITE, PP_ALIGN.CENTER)

# ===== BOTTOM: Data Flow Components =====

# Data Pipeline
pipeline = add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(5.8), Inches(1.4), Inches(0.55), AZ_DARK_PURPLE, AZ_MAROON, 1)
add_text(slide1, "Data Pipeline", Inches(2.5), Inches(5.88), Inches(1.4), Inches(0.4), 9, True, WHITE, PP_ALIGN.CENTER)

# Test Data Mapping
mapping = add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(5.8), Inches(1.4), Inches(0.55), AZ_DARK_PURPLE, AZ_MAROON, 1)
add_text(slide1, "Test Data\nMapping", Inches(4.2), Inches(5.83), Inches(1.4), Inches(0.5), 8, True, WHITE, PP_ALIGN.CENTER)

# Schema Registry
schema = add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), Inches(5.8), Inches(1.4), Inches(0.55), AZ_DARK_PURPLE, AZ_MAROON, 1)
add_text(slide1, "Schema\nRegistry", Inches(5.9), Inches(5.83), Inches(1.4), Inches(0.5), 8, True, WHITE, PP_ALIGN.CENTER)

# ===== ARROWS / FLOW INDICATORS =====
# Add arrow indicators as text (simplified)
add_text(slide1, "→", Inches(2.0), Inches(2.3), Inches(0.3), Inches(0.3), 16, True, AZ_MAROON, PP_ALIGN.CENTER)
add_text(slide1, "→", Inches(2.0), Inches(3.3), Inches(0.3), Inches(0.3), 16, True, AZ_MAROON, PP_ALIGN.CENTER)
add_text(slide1, "→", Inches(2.0), Inches(4.3), Inches(0.3), Inches(0.3), 16, True, AZ_MAROON, PP_ALIGN.CENTER)
add_text(slide1, "→", Inches(2.0), Inches(5.3), Inches(0.3), Inches(0.3), 16, True, AZ_MAROON, PP_ALIGN.CENTER)

# Horizontal arrows between agents
add_text(slide1, "↔", Inches(4.7), Inches(2.2), Inches(0.3), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.CENTER)
add_text(slide1, "↔", Inches(6.2), Inches(2.2), Inches(0.3), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.CENTER)
add_text(slide1, "↔", Inches(4.7), Inches(3.2), Inches(0.3), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.CENTER)
add_text(slide1, "↔", Inches(6.2), Inches(3.2), Inches(0.3), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.CENTER)
add_text(slide1, "↔", Inches(4.7), Inches(4.2), Inches(0.3), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.CENTER)
add_text(slide1, "↔", Inches(6.2), Inches(4.2), Inches(0.3), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.CENTER)

# Arrows to right side
add_text(slide1, "→", Inches(7.8), Inches(3.5), Inches(0.3), Inches(0.3), 16, True, AZ_MAROON, PP_ALIGN.CENTER)
add_text(slide1, "→", Inches(7.8), Inches(4.3), Inches(0.3), Inches(0.3), 16, True, AZ_MAROON, PP_ALIGN.CENTER)
add_text(slide1, "→", Inches(7.8), Inches(5.1), Inches(0.3), Inches(0.3), 16, True, AZ_MAROON, PP_ALIGN.CENTER)

# Vertical arrows
add_text(slide1, "↓", Inches(6.4), Inches(1.65), Inches(0.3), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.CENTER)
add_text(slide1, "↓", Inches(5.5), Inches(2.7), Inches(0.3), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.CENTER)
add_text(slide1, "↓", Inches(5.5), Inches(3.7), Inches(0.3), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.CENTER)
add_text(slide1, "↓", Inches(5.5), Inches(4.7), Inches(0.3), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.CENTER)

# ===== LEGEND =====
add_text(slide1, "Agent Types:", Inches(10.5), Inches(1.0), Inches(2.5), Inches(0.3), 9, True, TEXT_DARK, PP_ALIGN.LEFT)
add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(1.35), Inches(0.4), Inches(0.25), AZ_PINK, AZ_MAROON, 1)
add_text(slide1, "AI Agent", Inches(11.0), Inches(1.35), Inches(1.5), Inches(0.25), 8, False, TEXT_DARK, PP_ALIGN.LEFT)
add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(1.7), Inches(0.4), Inches(0.25), AZ_DARK_PURPLE, AZ_MAROON, 1)
add_text(slide1, "Data/Storage", Inches(11.0), Inches(1.7), Inches(1.5), Inches(0.25), 8, False, TEXT_DARK, PP_ALIGN.LEFT)
add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(2.05), Inches(0.4), Inches(0.25), AZ_LIGHT_PINK, AZ_MAROON, 1)
add_text(slide1, "Platform Group", Inches(11.0), Inches(2.05), Inches(1.5), Inches(0.25), 8, False, TEXT_DARK, PP_ALIGN.LEFT)

# ============================================================================
# SLIDE 2: Detailed Agent Descriptions
# ============================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

# Header
add_shape(slide2, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.7), AZ_MAROON)
add_text(slide2, "Test Intelligence Agent – Agent Descriptions", Inches(0.3), Inches(0.15), Inches(12), Inches(0.5), 22, True, WHITE, PP_ALIGN.LEFT)

# Agent descriptions table
agents_data = [
    ("Orchestrator Agent", "Central coordinator that manages workflow, delegates tasks to specialized agents, and ensures end-to-end test automation", "All Agents"),
    ("Requirement Gathering", "Analyzes requirements docs, user stories, and specs from R&D platforms to identify testable requirements", "3DP, BIKG, Regulatory"),
    ("Test Strategy Agent", "Determines optimal testing approach based on risk analysis, coverage goals, and platform constraints", "All Platforms"),
    ("Test Case Generation", "Auto-generates test cases from requirements using AI, covering positive, negative, and edge cases", "Schema Registry"),
    ("Test Script Generation", "Converts test cases into executable scripts for target platforms and frameworks", "Test Repository"),
    ("Synthetic Data Generation", "Creates schema-aware test datasets for HPC environments and validation scenarios", "Data Pipeline"),
    ("Test Execution Agent", "Runs tests in parallel across platforms, manages environments, collects results", "HPC, 3DP, BIKG"),
    ("Failure Analysis Agent", "AI-powered root cause analysis of test failures, identifies patterns and flaky tests", "Evidence Store"),
    ("Code Refactor Agent", "Self-healing capability - auto-fixes broken selectors, adapts to UI/API changes", "Test Repository"),
    ("Knowledge & Memory", "Maintains context across sessions, learns from past executions, stores test intelligence", "Knowledge Graph"),
    ("Reporting Agent", "Generates compliance reports, dashboards, and evidence capture for GxP audit trails", "Evidence Store"),
]

# Create table
table = slide2.shapes.add_table(12, 3, Inches(0.3), Inches(0.85), Inches(12.7), Inches(6.4)).table
table.columns[0].width = Inches(2.3)
table.columns[1].width = Inches(7.2)
table.columns[2].width = Inches(3.2)

# Header row
headers = ["Agent", "Function", "Platform Integration"]
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    cell.text_frame.paragraphs[0].font.size = Pt(10)
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.color.rgb = WHITE
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = AZ_MAROON
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

# Data rows
for i, (agent, function, integration) in enumerate(agents_data):
    bg = WHITE if i % 2 == 0 else BG_LIGHT

    # Agent name
    cell0 = table.cell(i+1, 0)
    cell0.text = agent
    cell0.text_frame.paragraphs[0].font.size = Pt(9)
    cell0.text_frame.paragraphs[0].font.bold = True
    cell0.text_frame.paragraphs[0].font.color.rgb = AZ_MAROON
    cell0.fill.solid()
    cell0.fill.fore_color.rgb = bg
    cell0.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Function
    cell1 = table.cell(i+1, 1)
    cell1.text = function
    cell1.text_frame.paragraphs[0].font.size = Pt(8)
    cell1.text_frame.paragraphs[0].font.color.rgb = TEXT_DARK
    cell1.fill.solid()
    cell1.fill.fore_color.rgb = bg
    cell1.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Integration
    cell2 = table.cell(i+1, 2)
    cell2.text = integration
    cell2.text_frame.paragraphs[0].font.size = Pt(8)
    cell2.text_frame.paragraphs[0].font.color.rgb = TEXT_MEDIUM
    cell2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cell2.fill.solid()
    cell2.fill.fore_color.rgb = bg
    cell2.vertical_anchor = MSO_ANCHOR.MIDDLE

# ============================================================================
# SLIDE 3: Platform Integration Matrix
# ============================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

# Header
add_shape(slide3, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.7), AZ_MAROON)
add_text(slide3, "R&D Platform Integration Matrix", Inches(0.3), Inches(0.15), Inches(12), Inches(0.5), 22, True, WHITE, PP_ALIGN.LEFT)

add_text(slide3, "How the Test Intelligence Agent interacts with each R&D platform", Inches(0.3), Inches(0.8), Inches(12), Inches(0.3), 11, False, TEXT_MEDIUM, PP_ALIGN.LEFT, italic=True)

# Platform integration details
platforms_detail = [
    ("3DP", "Drug Development Platform", "Requirements extraction, test execution, API validation", "High"),
    ("BIKG", "Biological Intelligence KG", "Knowledge graph queries, ontology validation, data integrity", "High"),
    ("Patient Safety", "Safety Reporting System", "Compliance testing, regulatory validation, audit trails", "Critical"),
    ("Clinical Trials", "Trial Management System", "Protocol validation, data consistency, CDISC compliance", "Critical"),
    ("CDISC/ADaM", "Clinical Data Standards", "Format validation, transformation testing, standard compliance", "High"),
    ("Regulatory Docs", "Submission Documents", "Module 2/3 validation, cross-reference integrity", "Critical"),
    ("Study PLDB", "Patient Level Database", "Data quality validation, PII protection testing", "High"),
    ("GxP Systems", "Validated Systems", "21 CFR Part 11 compliance, audit trail testing", "Critical"),
    ("HPC Environment", "High Performance Compute", "Environment testing, synthetic data scenarios", "Medium"),
]

# Create table
platform_table = slide3.shapes.add_table(10, 4, Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.5)).table
platform_table.columns[0].width = Inches(1.8)
platform_table.columns[1].width = Inches(2.8)
platform_table.columns[2].width = Inches(5.8)
platform_table.columns[3].width = Inches(2.3)

# Header row
p_headers = ["Platform", "Full Name", "Test Agent Integration", "Priority"]
for i, h in enumerate(p_headers):
    cell = platform_table.cell(0, i)
    cell.text = h
    cell.text_frame.paragraphs[0].font.size = Pt(10)
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.color.rgb = WHITE
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = AZ_MAROON
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

# Data rows
priority_colors = {"Critical": AZ_PINK, "High": AZ_GOLD, "Medium": AZ_BLUE}
for i, (platform, full_name, integration, priority) in enumerate(platforms_detail):
    bg = WHITE if i % 2 == 0 else BG_LIGHT

    cell0 = platform_table.cell(i+1, 0)
    cell0.text = platform
    cell0.text_frame.paragraphs[0].font.size = Pt(9)
    cell0.text_frame.paragraphs[0].font.bold = True
    cell0.text_frame.paragraphs[0].font.color.rgb = AZ_DARK_PURPLE
    cell0.fill.solid()
    cell0.fill.fore_color.rgb = bg
    cell0.vertical_anchor = MSO_ANCHOR.MIDDLE

    cell1 = platform_table.cell(i+1, 1)
    cell1.text = full_name
    cell1.text_frame.paragraphs[0].font.size = Pt(8)
    cell1.text_frame.paragraphs[0].font.color.rgb = TEXT_DARK
    cell1.fill.solid()
    cell1.fill.fore_color.rgb = bg
    cell1.vertical_anchor = MSO_ANCHOR.MIDDLE

    cell2 = platform_table.cell(i+1, 2)
    cell2.text = integration
    cell2.text_frame.paragraphs[0].font.size = Pt(8)
    cell2.text_frame.paragraphs[0].font.color.rgb = TEXT_MEDIUM
    cell2.fill.solid()
    cell2.fill.fore_color.rgb = bg
    cell2.vertical_anchor = MSO_ANCHOR.MIDDLE

    cell3 = platform_table.cell(i+1, 3)
    cell3.text = priority
    cell3.text_frame.paragraphs[0].font.size = Pt(9)
    cell3.text_frame.paragraphs[0].font.bold = True
    cell3.text_frame.paragraphs[0].font.color.rgb = WHITE
    cell3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cell3.fill.solid()
    cell3.fill.fore_color.rgb = priority_colors.get(priority, AZ_BLUE)
    cell3.vertical_anchor = MSO_ANCHOR.MIDDLE

# Bottom note
add_text(slide3, "Priority Legend:  ", Inches(0.3), Inches(6.85), Inches(1.2), Inches(0.3), 9, True, TEXT_DARK, PP_ALIGN.LEFT)
add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(6.9), Inches(0.8), Inches(0.25), AZ_PINK, AZ_MAROON, 1)
add_text(slide3, "Critical", Inches(1.55), Inches(6.9), Inches(0.7), Inches(0.25), 8, True, WHITE, PP_ALIGN.CENTER)
add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(6.9), Inches(0.6), Inches(0.25), AZ_GOLD, AZ_MAROON, 1)
add_text(slide3, "High", Inches(2.55), Inches(6.9), Inches(0.5), Inches(0.25), 8, True, WHITE, PP_ALIGN.CENTER)
add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.3), Inches(6.9), Inches(0.7), Inches(0.25), AZ_BLUE, AZ_MAROON, 1)
add_text(slide3, "Medium", Inches(3.35), Inches(6.9), Inches(0.6), Inches(0.25), 8, True, WHITE, PP_ALIGN.CENTER)

# Save
output_path = "/Users/devesh.b.sharma/Astra Zeneca/R&D IT/Test-Intelligence-Agent-Architecture.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
