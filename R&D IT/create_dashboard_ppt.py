from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Colors
AZ_MAROON = RGBColor(0x83, 0x00, 0x51)
AZ_GOLD = RGBColor(0xC4, 0xA0, 0x00)
AZ_GREEN = RGBColor(0x2E, 0x7D, 0x32)
TEXT_DARK = RGBColor(0x2D, 0x2D, 0x2D)
TEXT_MEDIUM = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_CREAM = RGBColor(0xFA, 0xF9, 0xF7)

def set_cell(cell, text, size=10, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, fill=None):
    cell.text = text
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill

def text(slide, txt, left, top, width, height, size=12, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    p = box.text_frame.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# === TITLE ===
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.7))
p = title.text_frame.paragraphs[0]
r1 = p.add_run()
r1.text = "Example of "
r1.font.size = Pt(28)
r1.font.italic = True
r1.font.color.rgb = TEXT_DARK
r2 = p.add_run()
r2.text = "Monitoring & Benefits Dashboard"
r2.font.size = Pt(28)
r2.font.italic = True
r2.font.bold = True
r2.font.color.rgb = AZ_MAROON

text(slide, "Test Intelligence Agent — Phase 1", Inches(0.5), Inches(0.7), Inches(10), Inches(0.3), 12, color=TEXT_MEDIUM)

# === LEFT SIDE: BENEFITS TABLE ===
text(slide, "BENEFITS", Inches(0.5), Inches(1.1), Inches(3), Inches(0.3), 11, True, AZ_MAROON)

benefits_table = slide.shapes.add_table(4, 1, Inches(0.5), Inches(1.4), Inches(5.5), Inches(5.4)).table
benefits_table.columns[0].width = Inches(5.5)

benefits = [
    ("Reduced Manual Testing Effort", "AI generates test cases automatically, freeing engineers to focus on development"),
    ("Increased Test Coverage", "Discovers edge cases and scenarios that manual testing often misses"),
    ("Faster Release Cycles", "Automated test generation accelerates validation and deployment timelines"),
    ("Regulatory Compliance", "Consistent, auditable test coverage for FDA/EMA submission requirements"),
]

for i, (title_txt, desc) in enumerate(benefits):
    cell = benefits_table.cell(i, 0)
    cell.text = f"{title_txt}\n{desc}"
    p1 = cell.text_frame.paragraphs[0]
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_DARK
    if len(cell.text_frame.paragraphs) > 1:
        p2 = cell.text_frame.paragraphs[1]
        p2.font.size = Pt(9)
        p2.font.bold = False
        p2.font.color.rgb = TEXT_MEDIUM
    cell.fill.solid()
    cell.fill.fore_color.rgb = WHITE

# === RIGHT SIDE: MONITORING ===
text(slide, "MONITORING EXAMPLE", Inches(6.3), Inches(1.1), Inches(6), Inches(0.3), 11, True, AZ_MAROON)

# KPI Table
text(slide, "Key Performance Indicators", Inches(6.3), Inches(1.45), Inches(6), Inches(0.25), 9, True, TEXT_DARK)
kpi_table = slide.shapes.add_table(2, 4, Inches(6.3), Inches(1.7), Inches(6.5), Inches(0.9)).table

kpis = [("Test Effort", "Reduction"), ("Coverage", "Increase"), ("Data Quality", "Score"), ("Platforms", "Integrated")]
placeholders = ["[%]", "[%]", "[%]", "[#]"]
kpi_colors = [AZ_GREEN, AZ_GREEN, AZ_GOLD, AZ_MAROON]

for i in range(4):
    kpi_table.columns[i].width = Inches(1.6)
    set_cell(kpi_table.cell(0, i), placeholders[i], 16, True, kpi_colors[i], PP_ALIGN.CENTER, WHITE)
    set_cell(kpi_table.cell(1, i), f"{kpis[i][0]}\n{kpis[i][1]}", 8, False, TEXT_DARK, PP_ALIGN.CENTER, BG_CREAM)

# Usage Metrics Table
text(slide, "Usage Metrics", Inches(6.3), Inches(2.75), Inches(3), Inches(0.25), 9, True, TEXT_DARK)
usage_table = slide.shapes.add_table(4, 2, Inches(6.3), Inches(3.0), Inches(3.0), Inches(1.5)).table
usage_table.columns[0].width = Inches(1.8)
usage_table.columns[1].width = Inches(1.2)

usage = [("Tests Generated", "[count]"), ("Active Users", "[users]"), ("Avg Response Time", "[seconds]"), ("Success Rate", "[percent]")]
for i, (name, val) in enumerate(usage):
    set_cell(usage_table.cell(i, 0), name, 9, False, TEXT_DARK, PP_ALIGN.LEFT, WHITE)
    set_cell(usage_table.cell(i, 1), val, 9, True, AZ_GREEN, PP_ALIGN.CENTER, WHITE)

# Platform Integration Table
text(slide, "Platform Integration", Inches(9.5), Inches(2.75), Inches(3), Inches(0.25), 9, True, TEXT_DARK)
platform_table = slide.shapes.add_table(2, 2, Inches(9.5), Inches(3.0), Inches(3.2), Inches(1.5)).table
platform_table.columns[0].width = Inches(2.4)
platform_table.columns[1].width = Inches(0.8)

platforms = [("3DP - Drug Development", "READY"), ("BIKG - Knowledge Graph", "READY")]
for i, (name, status) in enumerate(platforms):
    set_cell(platform_table.cell(i, 0), name, 9, False, TEXT_DARK, PP_ALIGN.LEFT, WHITE)
    set_cell(platform_table.cell(i, 1), status, 9, True, WHITE, PP_ALIGN.CENTER, AZ_GREEN)

# Cost Tracking Table
text(slide, "Cost Tracking", Inches(6.3), Inches(4.65), Inches(3), Inches(0.25), 9, True, TEXT_DARK)
cost_table = slide.shapes.add_table(3, 2, Inches(6.3), Inches(4.9), Inches(3.0), Inches(1.1)).table
cost_table.columns[0].width = Inches(1.8)
cost_table.columns[1].width = Inches(1.2)

costs = [("LLM API (Bedrock)", "[$/month]"), ("Embedding Service", "[$/month]"), ("Infrastructure", "[$/month]")]
for i, (name, val) in enumerate(costs):
    set_cell(cost_table.cell(i, 0), name, 9, False, TEXT_DARK, PP_ALIGN.LEFT, WHITE)
    set_cell(cost_table.cell(i, 1), val, 9, True, AZ_GOLD, PP_ALIGN.CENTER, WHITE)

# Model Performance Table
text(slide, "Model Performance", Inches(9.5), Inches(4.65), Inches(3), Inches(0.25), 9, True, TEXT_DARK)
perf_table = slide.shapes.add_table(3, 2, Inches(9.5), Inches(4.9), Inches(3.2), Inches(1.1)).table
perf_table.columns[0].width = Inches(1.8)
perf_table.columns[1].width = Inches(1.4)

perfs = [("Test Accuracy", "[score]"), ("Data Validity", "[score]"), ("User Satisfaction", "[score]")]
for i, (name, val) in enumerate(perfs):
    set_cell(perf_table.cell(i, 0), name, 9, False, TEXT_DARK, PP_ALIGN.LEFT, WHITE)
    set_cell(perf_table.cell(i, 1), val, 9, True, AZ_GREEN, PP_ALIGN.CENTER, AZ_GREEN)
    perf_table.cell(i, 1).text_frame.paragraphs[0].font.color.rgb = WHITE

# Status Row Table
status_table = slide.shapes.add_table(1, 2, Inches(6.3), Inches(6.15), Inches(6.4), Inches(0.45)).table
status_table.columns[0].width = Inches(3.1)
status_table.columns[1].width = Inches(3.3)

set_cell(status_table.cell(0, 0), "Budget Status: On Track", 10, True, AZ_GREEN, PP_ALIGN.CENTER, BG_CREAM)
set_cell(status_table.cell(0, 1), "ROI Tracking: Positive", 10, True, WHITE, PP_ALIGN.CENTER, AZ_MAROON)

prs.save("/Users/devesh.b.sharma/Astra Zeneca/R&D IT/Test-Intelligence-Agent-Dashboard.pptx")
print("Done!")
