from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
MAROON = RGBColor(0x83, 0x00, 0x51)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
DARK_PURPLE = RGBColor(0x5C, 0x3D, 0x6E)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
BLUE = RGBColor(0x15, 0x65, 0xC0)
ORANGE = RGBColor(0xFF, 0x6F, 0x00)
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)
RED = RGBColor(0xC6, 0x28, 0x28)
DARK_GRAY = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(255, 255, 255)
GRAY_66 = RGBColor(0x66, 0x66, 0x66)
GRAY_55 = RGBColor(0x55, 0x55, 0x55)

def add_text_box(slide, left, top, width, height, text, font_size=12, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return shape

def add_rectangle(slide, left, top, width, height, fill_color, text="", font_size=10, font_color=WHITE, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
    return shape

# ============ SLIDE 1: Monitoring & Benefits Dashboard ============
slide_layout = prs.slide_layouts[6]  # Blank
slide = prs.slides.add_slide(slide_layout)

# Title
add_text_box(slide, 0.5, 0.3, 12, 0.5, "Example of Monitoring and Benefits Dashboard", 24, True, DARK_GRAY)
add_text_box(slide, 0.5, 0.7, 12, 0.3, "Use an example for one of our four agents", 14, False, MAROON)

# Subtitle
add_text_box(slide, 0.5, 1.1, 8, 0.4, "Monitoring & Benefits Dashboard", 22, True, DARK_GRAY)
add_text_box(slide, 0.5, 1.5, 6, 0.3, "Test Intelligence Agent — Phase 1", 12, False, GRAY_66)

# Platform badge
add_rectangle(slide, 4.5, 1.5, 1, 0.25, MAROON, "3DP & BIKG", 9, WHITE, True)

# Benefits Section
add_text_box(slide, 0.5, 2.0, 2, 0.3, "BENEFITS", 11, True, MAROON)

benefits = [
    ("Reduced Manual Testing Effort", "AI generates test cases automatically for regulatory validation, freeing engineers to focus on development work"),
    ("Increased Test Coverage", "Discovers edge cases and validation scenarios that manual testing often misses, improving quality"),
    ("Faster Release Cycles", "Automated test generation accelerates validation timelines for submissions and data pipelines"),
    ("Regulatory Compliance", "Consistent, auditable test coverage for FDA/EMA submission requirements and GxP validation")
]

y_pos = 2.4
for title, desc in benefits:
    # Yellow bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y_pos), Inches(0.08), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

    add_text_box(slide, 0.7, y_pos, 5, 0.25, title, 12, True, DARK_GRAY)
    add_text_box(slide, 0.7, y_pos + 0.25, 5, 0.4, desc, 9, False, GRAY_66)
    y_pos += 0.75

# Monitoring Section
add_text_box(slide, 6.5, 2.0, 2, 0.3, "MONITORING", 11, True, MAROON)

# Metrics boxes
metrics = [("25%", "Test Effort\nReduction"), ("15%", "Coverage\nIncrease"), ("23", "Defects\nCaught"), ("2", "Platforms\nIntegrated")]
x_pos = 6.5
for value, label in metrics:
    add_rectangle(slide, x_pos, 2.4, 1.4, 0.9, LIGHT_GRAY, "", 10)
    add_text_box(slide, x_pos, 2.5, 1.4, 0.4, value, 28, True, MAROON, PP_ALIGN.CENTER)
    add_text_box(slide, x_pos, 2.95, 1.4, 0.35, label, 8, False, GRAY_66, PP_ALIGN.CENTER)
    x_pos += 1.5

# Usage Metrics
add_text_box(slide, 6.5, 3.5, 2, 0.25, "Usage Metrics", 10, True, MAROON)
usage = [("Tests Generated", "1,247"), ("Test Suites Created", "34"), ("Success Rate", "98.2%"), ("Avg Generation Time", "2.3s")]
y_pos = 3.8
for label, value in usage:
    add_text_box(slide, 6.5, y_pos, 2.2, 0.2, label, 9, False, GRAY_55)
    add_text_box(slide, 8.7, y_pos, 0.8, 0.2, value, 9, True, MAROON, PP_ALIGN.RIGHT)
    y_pos += 0.25

# Cost Tracking
add_text_box(slide, 6.5, 5.0, 2, 0.25, "Cost Tracking", 10, True, MAROON)
costs = [("LLM API (Bedrock)", "$4,200"), ("Embedding Service", "$850"), ("Infrastructure", "$500")]
y_pos = 5.3
for label, value in costs:
    add_text_box(slide, 6.5, y_pos, 2.2, 0.2, label, 9, False, GRAY_55)
    add_text_box(slide, 8.7, y_pos, 0.8, 0.2, value, 9, True, MAROON, PP_ALIGN.RIGHT)
    y_pos += 0.25

# Platform Integration
add_text_box(slide, 10, 3.5, 2.5, 0.25, "Platform Integration", 10, True, MAROON)
add_rectangle(slide, 10, 3.8, 2.8, 0.5, LIGHT_GRAY, "")
add_text_box(slide, 10.1, 3.85, 1.5, 0.2, "3DP", 11, True, DARK_GRAY)
add_text_box(slide, 10.1, 4.05, 2, 0.2, "Drug Development Platform", 8, False, GRAY_66)
add_rectangle(slide, 12.1, 3.9, 0.5, 0.25, GREEN, "LIVE", 8, WHITE, True)

add_rectangle(slide, 10, 4.4, 2.8, 0.5, LIGHT_GRAY, "")
add_text_box(slide, 10.1, 4.45, 1.5, 0.2, "BIKG", 11, True, DARK_GRAY)
add_text_box(slide, 10.1, 4.65, 2, 0.2, "Knowledge Graph", 8, False, GRAY_66)
add_rectangle(slide, 12.1, 4.5, 0.5, 0.25, GREEN, "LIVE", 8, WHITE, True)

# Model Performance
add_text_box(slide, 10, 5.1, 2.5, 0.25, "Model Performance", 10, True, MAROON)
perf = [("Test Accuracy", "94%"), ("Data Validity", "91%"), ("First Review Approval", "89%")]
y_pos = 5.4
for label, value in perf:
    add_text_box(slide, 10, y_pos, 1.8, 0.2, label, 9, False, GRAY_55)
    add_rectangle(slide, 11.8, y_pos, 0.9, 0.2, GREEN, value, 8, WHITE, True)
    y_pos += 0.3

# Budget and ROI bars
add_rectangle(slide, 6.5, 6.3, 3, 0.4, RGBColor(0xE8, 0xF5, 0xE9), "Budget: $5,550/mo (Under $7.5K target)", 10, GREEN, True)
add_rectangle(slide, 9.7, 6.3, 3, 0.4, GREEN, "ROI: $48K Annual Savings", 11, WHITE, True)

# ============ SLIDE 2: Key Performance Indicators ============
slide = prs.slides.add_slide(slide_layout)

add_text_box(slide, 0.5, 0.4, 10, 0.5, "Key Performance Indicators", 28, True, DARK_GRAY)
add_text_box(slide, 0.5, 0.85, 10, 0.3, "Definition and measurement methodology for each KPI", 12, False, GRAY_66)

kpis = [
    ("Test Effort Reduction", "25%", "Percentage decrease in manual hours spent writing test cases compared to baseline.",
     "(Baseline Hours − Current Hours) / Baseline Hours × 100",
     "Before: 40 hrs/sprint for Module 2/3 validation tests → After: 30 hrs/sprint = 25% reduction"),
    ("Coverage Increase", "15%", "Increase in test scenarios covered, including edge cases previously untested.",
     "(New Scenarios − Original) / Original × 100",
     "AI discovered 47 new edge cases for GxP validation and CDISC format compliance"),
    ("Data Quality Score", "90%", "Validity rate of AI-generated synthetic test data passing all schema validations.",
     "Valid Records / Total Generated × 100",
     "Synthetic patient data for submissions — 90% passes CDISC and regulatory checks"),
    ("Platforms Integrated", "2", "Number of R&D IT platforms where the agent is deployed and generating tests.",
     "3DP (Drug Development Data Platform) and BIKG (Biological Intelligence KG)",
     "2 primary platforms in Phase 1, extensible to additional platforms")
]

positions = [(0.5, 1.4), (6.7, 1.4), (0.5, 4.2), (6.7, 4.2)]
for i, (title, value, desc, measurement, example) in enumerate(kpis):
    x, y = positions[i]

    # Card border
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.8), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    add_text_box(slide, x + 0.2, y + 0.2, 4, 0.3, title, 16, True, DARK_GRAY)
    add_rectangle(slide, x + 4.8, y + 0.2, 0.8, 0.35, RGBColor(0xE8, 0xF5, 0xE9), value, 14, GREEN, True)
    add_text_box(slide, x + 0.2, y + 0.6, 5.4, 0.5, desc, 10, False, GRAY_66)

    # Measurement box
    add_rectangle(slide, x + 0.2, y + 1.2, 5.4, 0.5, LIGHT_GRAY, "", 9)
    add_text_box(slide, x + 0.3, y + 1.25, 1.2, 0.2, "Measurement:", 9, True, MAROON)
    add_text_box(slide, x + 0.3, y + 1.45, 5.2, 0.3, measurement, 9, False, GRAY_55)

    # Example box
    add_rectangle(slide, x + 0.2, y + 1.8, 5.4, 0.5, LIGHT_GRAY, "", 9)
    add_text_box(slide, x + 0.3, y + 1.85, 1, 0.2, "Example:", 9, True, MAROON)
    add_text_box(slide, x + 0.3, y + 2.05, 5.2, 0.3, example, 9, False, GRAY_55)

# ============ SLIDE 3: Platform Use Cases ============
slide = prs.slides.add_slide(slide_layout)

add_text_box(slide, 0.5, 0.4, 10, 0.5, "Platform Use Cases", 28, True, DARK_GRAY)
add_text_box(slide, 0.5, 0.85, 10, 0.3, "How Test Intelligence Agent delivers value on 3DP and BIKG", 12, False, GRAY_66)

use_cases = [
    ("3DP", MAROON, "Regulatory Validation Testing", "Auto-generate validation tests for FDA/EMA submission modules and GxP compliance.",
     ["Module 2/3 submission validation", "Patient Safety data checks", "GxP compliance testing", "Audit trail validation"]),
    ("3DP", MAROON, "Data Pipeline Testing", "Validate data transformations and integrations across the platform.",
     ["Portfolio Tower ingestion tests", "Study Cycle Time benchmarking", "CRO operational data validation", "Cross-reference integrity checks"]),
    ("3DP", MAROON, "Synthetic Test Data", "Generate compliant synthetic data for testing without using real patient data.",
     ["CDISC-compliant patient records", "Adverse event test data", "Lab results within valid ranges", "Clinical trial mock data"]),
    ("BIKG", GOLD, "Knowledge Graph Query Testing", "Generate test cases for graph traversal and relationship queries.",
     ["Gene-disease relationship queries", "Protein-drug interactions", "Pathway traversal validation", "Query performance testing"]),
    ("BIKG", GOLD, "Entity Validation", "Validate biological entity data integrity and ontology consistency.",
     ["Entity attribute completeness", "Ontology term consistency", "Cross-reference accuracy", "Relationship cardinality checks"]),
    ("BIKG", GOLD, "Synthetic Biological Data", "Generate test data for biological entities and relationships.",
     ["Mock gene expression data", "Synthetic pathway data", "Test entity relationships", "Sample KG traversal results"])
]

positions = [(0.5, 1.3), (4.5, 1.3), (8.5, 1.3), (0.5, 4.1), (4.5, 4.1), (8.5, 4.1)]
for i, (platform, color, title, desc, items) in enumerate(use_cases):
    x, y = positions[i]

    # Top color bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    # Card
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.08), Inches(3.8), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    add_rectangle(slide, x + 0.1, y + 0.2, 0.5, 0.2, color, platform, 8, WHITE, True)
    add_text_box(slide, x + 0.1, y + 0.5, 3.6, 0.3, title, 12, True, DARK_GRAY)
    add_text_box(slide, x + 0.1, y + 0.8, 3.6, 0.4, desc, 9, False, GRAY_66)

    y_item = y + 1.3
    for item in items:
        add_text_box(slide, x + 0.1, y_item, 3.6, 0.2, "• " + item, 9, False, GRAY_55)
        y_item += 0.25

# ============ SLIDE 4: Architecture ============
slide = prs.slides.add_slide(slide_layout)

add_text_box(slide, 0.5, 0.4, 12, 0.6, "Scalable, Agent Mesh, future proof architecture for AZ's R&D Agentic ambition", 24, True, DARK_GRAY)

# Architecture placeholder box
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(8.5), Inches(5.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
shape.line.color.rgb = RED
shape.line.width = Pt(2)

add_text_box(slide, 1, 2.5, 7, 0.3, "[Architecture Diagram]", 14, False, RGBColor(0x99, 0x99, 0x99), PP_ALIGN.CENTER)
add_text_box(slide, 1, 3, 7, 0.3, "AWS Cloud | VPC | ECS Fargate | Agent Runtime | Mesh Control Plane", 10, False, RGBColor(0x99, 0x99, 0x99), PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.4, 7, 0.3, "Azure AD → Route 53 → ALB → Frontend ECS → Agent Runtime ECS → Bedrock", 10, False, RGBColor(0x99, 0x99, 0x99), PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.8, 7, 0.3, "4 Agents: Test Intelligence | Dependency | Deployment | Regulatory", 10, False, RGBColor(0x99, 0x99, 0x99), PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.2, 7, 0.3, "CI/CD: CodeCommit → CodeBuild → CodeDeploy → ECR", 10, False, RGBColor(0x99, 0x99, 0x99), PP_ALIGN.CENTER)

# Right side notes
add_text_box(slide, 9.3, 1.5, 3.5, 1.5,
    "1. These four agents collectively form a modular R&D intelligence fabric — enabling knowledge retrieval, experiment validation, data reasoning, and governed orchestration — all operating within a secure, federated Agent Mesh architecture aligned to AZ's Enterprise Agentic AI blueprint.",
    10, False, GRAY_55)

add_text_box(slide, 9.3, 3.5, 3.5, 1.5,
    "2. To accelerate implementation, we will leverage AI-assisted engineering (Claude Code) to generate validated scaffolding, typed contracts, test harnesses, and infrastructure templates — while maintaining full human governance, security review, and CI/CD controls.",
    10, False, GRAY_55)

# ============ SLIDE 5: AI for AI Delivery (Table) ============
slide = prs.slides.add_slide(slide_layout)

add_text_box(slide, 0.3, 0.25, 12, 0.5, "From Manual Processes to Intelligent Automation: AI-Powered R&D IT Platform Delivery", 22, True, DARK_GRAY)
add_text_box(slide, 0.3, 0.65, 12, 0.3, "AI Agents reduce cycle time by 40-60% while improving compliance, quality, and cross-team coordination for 3DP & BIKG Platforms", 11, False, GRAY_66)

# Banner
add_rectangle(slide, 0.3, 1.0, 12.7, 0.4, MAROON, "Accelerate delivery by 40-60% | AI-Orchestrated intelligent workflow across Testing, Dependencies, Deployments & Compliance", 11, WHITE, True)

# Table headers
headers = [("Phase\nDuration", DARK_PURPLE, 1.0),
           ("Test Intelligence\nWave 1", GREEN, 2.7),
           ("Dependency Coordinator\nWave 1", BLUE, 2.7),
           ("Deployment & Change\nWave 2", ORANGE, 2.7),
           ("Regulatory Readiness\nWave 2", PURPLE, 2.7)]

x_pos = 0.3
for header, color, width in headers:
    add_rectangle(slide, x_pos, 1.5, width, 0.55, MAROON if not header.startswith("Phase") else color, header, 10, WHITE, True)
    if not header.startswith("Phase"):
        # Color bar on top
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_pos), Inches(1.45), Inches(width), Inches(0.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
    x_pos += width + 0.05

# Today row
add_rectangle(slide, 0.3, 2.1, 1.0, 1.1, DARK_PURPLE, "Today\n(Human-Led)", 9, WHITE, True)
today_content = [
    "3-4 weeks\nManual test authoring. Long regression cycles, defect leakage.",
    "2-3 weeks\nManual tracking via Jira/ADO. Cross-team friction, approval delays.",
    "1-2 weeks\nManual IaC. High change failure rate, drift incidents.",
    "2-4 weeks\nManual evidence gathering. Long audit prep, missing artefacts."
]
x_pos = 1.35
for content in today_content:
    add_rectangle(slide, x_pos, 2.1, 2.7, 1.1, RGBColor(0xFA, 0xFA, 0xFA), "", 8)
    add_text_box(slide, x_pos + 0.1, 2.15, 2.5, 1.0, content, 9, False, DARK_GRAY)
    x_pos += 2.75

# Tomorrow row
add_rectangle(slide, 0.3, 3.25, 1.0, 1.1, GOLD, "Tomorrow\n(AI-Powered)", 9, WHITE, True)
tomorrow_content = [
    "1 week\nAutomated test case generation. 25% reduction in test effort.",
    "Days\nAutomated dependency tracking. Early conflict detection.",
    "Hours\nCommon IaC patterns. Reduced change failure rate, faster MTTR.",
    "Hours\nAutomated evidence generation. 100% complete evidence."
]
x_pos = 1.35
for content in tomorrow_content:
    add_rectangle(slide, x_pos, 3.25, 2.7, 1.1, RGBColor(0xFF, 0xFD, 0xF5), "", 8)
    add_text_box(slide, x_pos + 0.1, 3.3, 2.5, 1.0, content, 9, False, GREEN)
    x_pos += 2.75

# Tools row
add_rectangle(slide, 0.3, 4.4, 1.0, 0.9, LIGHT_GRAY, "AI Tools\nUsed", 9, DARK_GRAY, True)
tools_content = ["Claude\nTest Intelligence Agent", "Claude\nDependency Coordinator", "Claude + Copilot\nDeploy & Change Agent", "Claude\nRegulatory Readiness Agent"]
x_pos = 1.35
for content in tools_content:
    add_rectangle(slide, x_pos, 4.4, 2.7, 0.9, WHITE, "", 8)
    add_text_box(slide, x_pos + 0.1, 4.5, 2.5, 0.8, content, 9, False, DARK_GRAY, PP_ALIGN.CENTER)
    x_pos += 2.75

# Value row
add_rectangle(slide, 0.3, 5.35, 1.0, 0.9, GREEN, "Value\nDelivered", 9, WHITE, True)
value_content = ["25% faster\nReduced test effort\nHigher coverage", "50% faster\nReduced approval cycle\nEarly conflict detection", "70% faster\nNear-instant deployments\nReduced change failure", "100% compliant\nComplete evidence packs\nFaster audit readiness"]
x_pos = 1.35
for content in value_content:
    add_rectangle(slide, x_pos, 5.35, 2.7, 0.9, RGBColor(0xE8, 0xF5, 0xE9), "", 8)
    add_text_box(slide, x_pos + 0.1, 5.4, 2.5, 0.85, content, 9, True, GREEN, PP_ALIGN.CENTER)
    x_pos += 2.75

# Impact footer
add_rectangle(slide, 0.3, 6.4, 12.7, 0.5, DARK_GRAY, "", 10)
add_text_box(slide, 0.5, 6.5, 2, 0.3, "AI for AI Delivery Impact:", 11, True, GOLD)
add_text_box(slide, 2.8, 6.5, 10, 0.3, "✓ Launch readiness accelerated by 40-60%    ✓ 70% reduction in manual effort    ✓ 100% regulatory compliance    ✓ Enterprise-grade quality", 9, False, WHITE)

# ============ SLIDE 6: AI for AI Delivery (3-Column) ============
slide = prs.slides.add_slide(slide_layout)

add_text_box(slide, 0.3, 0.25, 10, 0.4, "'AI for AI' approach for delivery with R&D IT Platforms", 22, True, DARK_GRAY)
add_text_box(slide, 0.3, 0.6, 12, 0.3, "Leveraging an AI-first & Agentic approach to establish the AI foundation in AZ's environment for 3DP & BIKG Platforms", 11, False, MAROON)

# TODAY Column
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.0), Inches(4.0), Inches(5.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xF9, 0xF9, 0xF9)
shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

add_text_box(slide, 0.5, 1.1, 2, 0.3, "Today:", 16, True, DARK_GRAY)
add_text_box(slide, 0.5, 1.4, 3.6, 0.3, "Manual, error-prone processes", 9, False, MAROON)

add_text_box(slide, 0.5, 1.8, 3, 0.2, "WAVE 1 (MONTHS 0-3)", 8, True, GREEN)

# Today phases
phases_today = [
    (GREEN, "Testing", "3-4 wks", "Manual test authoring. Error-prone testing."),
    (BLUE, "Dependencies", "2-3 wks", "Cross-team friction. Approval delays."),
]
y_pos = 2.1
for color, name, duration, desc in phases_today:
    add_rectangle(slide, 0.5, y_pos, 0.7, 0.6, color, name, 8, WHITE, True)
    add_rectangle(slide, 1.25, y_pos, 0.5, 0.6, LIGHT_GRAY, duration, 8, GRAY_66, False)
    add_text_box(slide, 1.8, y_pos + 0.1, 2.3, 0.5, desc, 8, False, GRAY_55)
    y_pos += 0.7

add_text_box(slide, 0.5, 3.6, 3, 0.2, "WAVE 2 (MONTHS 3-6)", 8, True, PURPLE)

phases_today2 = [
    (ORANGE, "Deployment", "1-2 wks", "Manual IaC. High change failure rate."),
    (PURPLE, "Compliance", "2-4 wks", "Manual evidence. Missing artefacts."),
]
y_pos = 3.9
for color, name, duration, desc in phases_today2:
    add_rectangle(slide, 0.5, y_pos, 0.7, 0.6, color, name, 8, WHITE, True)
    add_rectangle(slide, 1.25, y_pos, 0.5, 0.6, LIGHT_GRAY, duration, 8, GRAY_66, False)
    add_text_box(slide, 1.8, y_pos + 0.1, 2.3, 0.5, desc, 8, False, GRAY_55)
    y_pos += 0.7

add_rectangle(slide, 0.5, 6.0, 3.6, 0.4, DARK_PURPLE, "Siloed processes with limited visibility", 8, WHITE, True)

# TOMORROW Column
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.0), Inches(4.3), Inches(5.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFD, 0xF5)
shape.line.color.rgb = GOLD

add_text_box(slide, 4.7, 1.1, 2, 0.3, "Tomorrow:", 16, True, DARK_GRAY)
add_text_box(slide, 4.7, 1.4, 4, 0.3, "Agentic AI-powered automation", 9, False, MAROON)

add_text_box(slide, 4.7, 1.8, 3, 0.2, "WAVE 1 (MONTHS 0-3)", 8, True, GREEN)

phases_tomorrow = [
    (GREEN, "Testing", "1 week", "Automated test generation. 25% reduction."),
    (BLUE, "Dependencies", "Days", "Automated tracking. Early conflict detection."),
]
y_pos = 2.1
for color, name, duration, desc in phases_tomorrow:
    add_rectangle(slide, 4.7, y_pos, 0.7, 0.6, color, name, 8, WHITE, True)
    add_rectangle(slide, 5.45, y_pos, 0.5, 0.6, LIGHT_GRAY, duration, 8, GRAY_66, False)
    add_text_box(slide, 6.0, y_pos + 0.1, 2.6, 0.5, desc, 8, False, GREEN)
    y_pos += 0.7

add_text_box(slide, 4.7, 3.6, 3, 0.2, "WAVE 2 (MONTHS 3-6)", 8, True, PURPLE)

phases_tomorrow2 = [
    (ORANGE, "Deployment", "Hours", "Common IaC patterns. Faster MTTR."),
    (PURPLE, "Compliance", "Hours", "Automated evidence. 100% complete."),
]
y_pos = 3.9
for color, name, duration, desc in phases_tomorrow2:
    add_rectangle(slide, 4.7, y_pos, 0.7, 0.6, color, name, 8, WHITE, True)
    add_rectangle(slide, 5.45, y_pos, 0.5, 0.6, LIGHT_GRAY, duration, 8, GRAY_66, False)
    add_text_box(slide, 6.0, y_pos + 0.1, 2.6, 0.5, desc, 8, False, GREEN)
    y_pos += 0.7

add_rectangle(slide, 4.7, 6.0, 3.9, 0.4, GOLD, "Unified Agent Mesh with integrated workflows", 8, WHITE, True)

# AI TOOLS Column
add_text_box(slide, 9.2, 1.0, 3.8, 0.3, "AI for AI Delivery", 14, True, MAROON, PP_ALIGN.RIGHT)

tools = [
    (RGBColor(0x63, 0x66, 0xF1), "Design Acceleration", "Accelerate front-end design through prompting"),
    (DARK_GRAY, "Code Generation", "Fast track code generation with AI assistants"),
    (GREEN, "Test Intelligence Agent", "Automates testing across 3DP & BIKG"),
    (BLUE, "Dependency Coordinator", "Reduces rework and approval delays"),
    (ORANGE, "Deployment & Change Agent", "IaC maturity with high compliance value"),
    (PURPLE, "Regulatory Readiness Agent", "Compliance-heavy, evidence-focused")
]

y_pos = 1.4
for color, title, desc in tools:
    # Left border
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.1), Inches(y_pos), Inches(0.05), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    add_rectangle(slide, 9.2, y_pos, 3.8, 0.7, RGBColor(0xF8, 0xF8, 0xF8), "", 8)
    add_text_box(slide, 9.3, y_pos + 0.08, 3.5, 0.2, title, 10, True, DARK_GRAY)
    add_text_box(slide, 9.3, y_pos + 0.32, 3.5, 0.35, desc, 8, False, GRAY_66)

    y_pos += 0.78

# Save presentation
output_path = '/Users/devesh.b.sharma/Astra Zeneca/R&D IT/RnD-IT-Agentic-AI-Presentation.pptx'
prs.save(output_path)
print(f"Presentation saved: {output_path}")
