from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Colors - AstraZeneca Brand
AZ_MAROON = RGBColor(0x83, 0x00, 0x51)
AZ_GOLD = RGBColor(0xC4, 0xA0, 0x00)
AZ_GREEN = RGBColor(0x2E, 0x7D, 0x32)
AZ_BLUE = RGBColor(0x1E, 0x88, 0xE5)
AZ_ORANGE = RGBColor(0xF5, 0x7C, 0x00)
AZ_RED = RGBColor(0xC6, 0x28, 0x28)
TEXT_DARK = RGBColor(0x2D, 0x2D, 0x2D)
TEXT_MEDIUM = RGBColor(0x55, 0x55, 0x55)
TEXT_LIGHT = RGBColor(0x75, 0x75, 0x75)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_CREAM = RGBColor(0xFA, 0xF9, 0xF7)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)

def set_cell(cell, text, size=10, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, fill=None):
    cell.text = text
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill

def add_text(slide, txt, left, top, width, height, size=12, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return box

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_agent_box(slide, name, status, status_color, left, top):
    """Add an agent box with status indicator"""
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.8), WHITE, AZ_MAROON, 2)
    add_text(slide, name, left + Inches(0.1), top + Inches(0.1), Inches(1.1), Inches(0.35), 9, True, AZ_MAROON, PP_ALIGN.CENTER)
    add_text(slide, status, left + Inches(0.1), top + Inches(0.45), Inches(1.1), Inches(0.25), 8, False, status_color, PP_ALIGN.CENTER)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])

# Background gradient effect (maroon bar)
add_shape(slide1, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(2.5), AZ_MAROON)

# Title
add_text(slide1, "Test Intelligence Agent", Inches(0.5), Inches(0.8), Inches(12), Inches(0.7), 44, True, WHITE, PP_ALIGN.LEFT)
add_text(slide1, "Agentic AI Dashboard for Autonomous QA", Inches(0.5), Inches(1.5), Inches(12), Inches(0.5), 24, False, WHITE, PP_ALIGN.LEFT, italic=True)

# Subtitle info
add_text(slide1, "Phase 1  |  3DP & BIKG Platforms  |  R&D IT", Inches(0.5), Inches(2.8), Inches(12), Inches(0.4), 14, False, TEXT_MEDIUM, PP_ALIGN.LEFT)

# Key stats boxes
stats = [
    ("5", "AI Agents", AZ_MAROON),
    ("78%", "Autonomy Rate", AZ_GREEN),
    ("94%", "Test Accuracy", AZ_GREEN),
    ("$48K", "Annual Savings", AZ_GOLD),
]
for i, (value, label, color) in enumerate(stats):
    left = Inches(0.5 + i * 3.2)
    box = add_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.8), Inches(2.8), Inches(1.4), WHITE, color, 3)
    add_text(slide1, value, left, Inches(3.95), Inches(2.8), Inches(0.7), 32, True, color, PP_ALIGN.CENTER)
    add_text(slide1, label, left, Inches(4.65), Inches(2.8), Inches(0.4), 12, False, TEXT_MEDIUM, PP_ALIGN.CENTER)

# What is Agentic AI
add_text(slide1, "What is Agentic AI Testing?", Inches(0.5), Inches(5.5), Inches(12), Inches(0.4), 16, True, AZ_MAROON, PP_ALIGN.LEFT)
add_text(slide1, "Intelligent AI agents that autonomously generate, execute, analyze, and heal tests - reducing manual effort by 40% while increasing coverage and catching 98% of critical bugs before production.",
         Inches(0.5), Inches(5.95), Inches(12), Inches(0.8), 12, False, TEXT_MEDIUM, PP_ALIGN.LEFT)

# ============================================================================
# SLIDE 2: BEFORE vs AFTER - The Transformation
# ============================================================================
slide_before_after = prs.slides.add_slide(prs.slide_layouts[6])

# Header
add_shape(slide_before_after, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.8), AZ_MAROON)
add_text(slide_before_after, "Before vs After: The Agentic AI Transformation", Inches(0.5), Inches(0.2), Inches(12), Inches(0.5), 24, True, WHITE, PP_ALIGN.LEFT)

# Divider line in the middle
add_shape(slide_before_after, MSO_SHAPE.RECTANGLE, Inches(6.55), Inches(1.0), Inches(0.05), Inches(6.2), AZ_MAROON)

# Process Bottleneck callout
add_text(slide_before_after, "Process Bottleneck: Slow Validation and Release Cycles", Inches(0.3), Inches(1.0), Inches(12.7), Inches(0.35), 12, True, AZ_MAROON, PP_ALIGN.CENTER, italic=True)

# ===== LEFT SIDE: BEFORE (Traditional Testing) =====
add_shape(slide_before_after, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.4), Inches(6.1), Inches(0.5), AZ_RED)
add_text(slide_before_after, "BEFORE: Manual Testing", Inches(0.3), Inches(1.5), Inches(6.1), Inches(0.4), 16, True, WHITE, PP_ALIGN.CENTER)

# Before metrics - aligned with Test Intelligence Agent context
before_metrics = [
    ("Test Design", "Manual", "Engineers write every test case"),
    ("Test Execution", "Manual", "Run tests one by one"),
    ("Failure Interpretation", "Manual", "Hours spent debugging failures"),
    ("Evidence Capture", "Manual", "Screenshots, logs collected by hand"),
    ("Synthetic Test Data", "None", "No schema-aware data generation"),
    ("HPC Environment Tests", "Limited", "Manual selection of test scenarios"),
    ("Validation Cycle", "2 weeks", "Slow release due to QA bottleneck"),
    ("Manual Testing Effort", "100%", "Heavy reliance on human testers"),
]

before_table = slide_before_after.shapes.add_table(9, 3, Inches(0.3), Inches(2.0), Inches(6.1), Inches(3.2)).table
before_table.columns[0].width = Inches(2.2)
before_table.columns[1].width = Inches(1.2)
before_table.columns[2].width = Inches(2.7)

set_cell(before_table.cell(0, 0), "Capability", 9, True, WHITE, PP_ALIGN.LEFT, AZ_RED)
set_cell(before_table.cell(0, 1), "Status", 9, True, WHITE, PP_ALIGN.CENTER, AZ_RED)
set_cell(before_table.cell(0, 2), "Pain Point", 9, True, WHITE, PP_ALIGN.LEFT, AZ_RED)

for i, (metric, value, pain) in enumerate(before_metrics):
    set_cell(before_table.cell(i+1, 0), metric, 9, False, TEXT_DARK, PP_ALIGN.LEFT, WHITE)
    set_cell(before_table.cell(i+1, 1), value, 9, True, AZ_RED, PP_ALIGN.CENTER, BG_LIGHT)
    set_cell(before_table.cell(i+1, 2), pain, 8, False, TEXT_MEDIUM, PP_ALIGN.LEFT, WHITE)

# Before challenges box
add_text(slide_before_after, "Key Challenges", Inches(0.3), Inches(5.35), Inches(6.1), Inches(0.3), 11, True, AZ_RED, PP_ALIGN.LEFT)
challenges_box = add_shape(slide_before_after, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(5.7), Inches(6.1), Inches(1.4), BG_LIGHT, AZ_RED, 2)

challenges = [
    "Slow Validation: QA bottleneck delays releases",
    "Rework: Failed tests require manual investigation",
    "Limited Coverage: Can't test all metadata scenarios",
    "Not Reusable: Tests tied to specific platforms",
]
for i, challenge in enumerate(challenges):
    add_text(slide_before_after, "✗ " + challenge, Inches(0.5), Inches(5.8 + i * 0.32), Inches(5.8), Inches(0.3), 9, False, AZ_RED, PP_ALIGN.LEFT)

# ===== RIGHT SIDE: AFTER (Agentic AI Testing) =====
add_shape(slide_before_after, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.4), Inches(6.1), Inches(0.5), AZ_GREEN)
add_text(slide_before_after, "AFTER: Test Intelligence Agent", Inches(6.9), Inches(1.5), Inches(6.1), Inches(0.4), 16, True, WHITE, PP_ALIGN.CENTER)

# After metrics - aligned with Test Intelligence Agent capabilities
after_metrics = [
    ("Test Design", "Automated", "AI generates test cases from specs"),
    ("Test Execution", "Automated", "Parallel execution across platforms"),
    ("Failure Interpretation", "Automated", "AI analyzes & explains failures"),
    ("Evidence Capture", "Automated", "Auto-capture logs, screenshots, traces"),
    ("Synthetic Test Data", "AI-Generated", "Schema-aware test datasets"),
    ("HPC Environment Tests", "Smart Selection", "AI selects relevant test scenarios"),
    ("Validation Cycle", "3 days", "78% faster release cycles"),
    ("Manual Testing Effort", "25-30% reduction", "Reusable across teams"),
]

after_table = slide_before_after.shapes.add_table(9, 3, Inches(6.9), Inches(2.0), Inches(6.1), Inches(3.2)).table
after_table.columns[0].width = Inches(2.2)
after_table.columns[1].width = Inches(1.4)
after_table.columns[2].width = Inches(2.5)

set_cell(after_table.cell(0, 0), "Capability", 9, True, WHITE, PP_ALIGN.LEFT, AZ_GREEN)
set_cell(after_table.cell(0, 1), "Status", 9, True, WHITE, PP_ALIGN.CENTER, AZ_GREEN)
set_cell(after_table.cell(0, 2), "Benefit", 9, True, WHITE, PP_ALIGN.LEFT, AZ_GREEN)

for i, (metric, value, benefit) in enumerate(after_metrics):
    set_cell(after_table.cell(i+1, 0), metric, 9, False, TEXT_DARK, PP_ALIGN.LEFT, WHITE)
    set_cell(after_table.cell(i+1, 1), value, 9, True, AZ_GREEN, PP_ALIGN.CENTER, BG_LIGHT)
    set_cell(after_table.cell(i+1, 2), benefit, 8, False, TEXT_MEDIUM, PP_ALIGN.LEFT, WHITE)

# After benefits box
add_text(slide_before_after, "Key Benefits", Inches(6.9), Inches(5.35), Inches(6.1), Inches(0.3), 11, True, AZ_GREEN, PP_ALIGN.LEFT)
benefits_box = add_shape(slide_before_after, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(5.7), Inches(6.1), Inches(1.4), BG_LIGHT, AZ_GREEN, 2)

benefits = [
    "Faster Validation: Accelerated release cycles",
    "Reduced Rework: AI interprets failures automatically",
    "Metadata-Scenario Testing: Improved reliability",
    "Reusable: Works across engineering & data teams",
]
for i, benefit in enumerate(benefits):
    add_text(slide_before_after, "✓ " + benefit, Inches(7.1), Inches(5.8 + i * 0.32), Inches(5.8), Inches(0.3), 9, False, AZ_GREEN, PP_ALIGN.LEFT)

# ============================================================================
# SLIDE 2B: Detailed Improvement Metrics
# ============================================================================
slide_improvements = prs.slides.add_slide(prs.slide_layouts[6])

# Header
add_shape(slide_improvements, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.8), AZ_MAROON)
add_text(slide_improvements, "Detailed Improvement Metrics", Inches(0.5), Inches(0.2), Inches(12), Inches(0.5), 24, True, WHITE, PP_ALIGN.LEFT)

# Score card from context
add_text(slide_improvements, "Test Intelligence Agent Scorecard", Inches(0.4), Inches(1.0), Inches(12), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.LEFT)

score_table = slide_improvements.shapes.add_table(2, 5, Inches(0.4), Inches(1.35), Inches(12.5), Inches(0.8)).table
score_table.columns[0].width = Inches(2.5)
score_table.columns[1].width = Inches(2.5)
score_table.columns[2].width = Inches(2.5)
score_table.columns[3].width = Inches(2.5)
score_table.columns[4].width = Inches(2.5)

scores = [("Efficiency", "5"), ("Scalability", "4"), ("Feasibility", "4"), ("Adoption", "4"), ("Weighted Score", "4.30")]
for i, (label, score) in enumerate(scores):
    set_cell(score_table.cell(0, i), label, 10, True, WHITE if i < 4 else TEXT_DARK, PP_ALIGN.CENTER, AZ_MAROON if i < 4 else AZ_GOLD)
    set_cell(score_table.cell(1, i), score, 14, True, AZ_MAROON if i < 4 else AZ_GOLD, PP_ALIGN.CENTER, WHITE)

# Main improvement KPIs
improvement_kpis = [
    ("25-30%", "Effort Reduction", "Manual testing reduced", AZ_GREEN),
    ("78%", "Faster Releases", "2 weeks → 3 days", AZ_BLUE),
    ("Auto", "Test Generation", "Schema-aware datasets", AZ_MAROON),
    ("Reusable", "Cross-Platform", "Engineering & data teams", AZ_GOLD),
]

for i, (value, label, desc, color) in enumerate(improvement_kpis):
    left = Inches(0.4 + i * 3.2)
    box = add_shape(slide_improvements, MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.3), Inches(3.0), Inches(1.2), WHITE, color, 3)
    add_text(slide_improvements, value, left, Inches(2.4), Inches(3.0), Inches(0.5), 28, True, color, PP_ALIGN.CENTER)
    add_text(slide_improvements, label, left, Inches(2.9), Inches(3.0), Inches(0.25), 10, True, TEXT_DARK, PP_ALIGN.CENTER)
    add_text(slide_improvements, desc, left, Inches(3.15), Inches(3.0), Inches(0.25), 8, False, TEXT_MEDIUM, PP_ALIGN.CENTER)

# Detailed comparison table
add_text(slide_improvements, "Detailed Capability Comparison", Inches(0.4), Inches(3.65), Inches(12), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.LEFT)

comparison_table = slide_improvements.shapes.add_table(9, 4, Inches(0.4), Inches(4.0), Inches(12.5), Inches(2.8)).table
comparison_table.columns[0].width = Inches(3.5)
comparison_table.columns[1].width = Inches(2.5)
comparison_table.columns[2].width = Inches(2.5)
comparison_table.columns[3].width = Inches(4.0)

set_cell(comparison_table.cell(0, 0), "Capability", 9, True, WHITE, PP_ALIGN.LEFT, AZ_MAROON)
set_cell(comparison_table.cell(0, 1), "Before", 9, True, WHITE, PP_ALIGN.CENTER, AZ_RED)
set_cell(comparison_table.cell(0, 2), "After", 9, True, WHITE, PP_ALIGN.CENTER, AZ_GREEN)
set_cell(comparison_table.cell(0, 3), "Impact", 9, True, WHITE, PP_ALIGN.LEFT, AZ_MAROON)

detailed_comparison = [
    ("Test Design & Creation", "Manual (40+ hrs/wk)", "Automated (AI)", "Engineers focus on development"),
    ("Test Execution", "Sequential, Manual", "Parallel, Automated", "Faster feedback loops"),
    ("Failure Interpretation", "Manual (hours)", "AI-Powered (seconds)", "Instant root cause analysis"),
    ("Evidence Capture", "Manual collection", "Auto-captured", "Complete audit trail for GxP"),
    ("Synthetic Test Data", "None / Manual", "Schema-aware generation", "Better edge case coverage"),
    ("HPC Environment Testing", "Limited selection", "AI-selected scenarios", "Relevant tests automatically"),
    ("Validation Cycle Time", "2 weeks", "3 days", "78% faster releases"),
    ("Manual Testing Effort", "100%", "70-75%", "25-30% reduction achieved"),
]

for i, (metric, before, after, improvement) in enumerate(detailed_comparison):
    bg = WHITE if i % 2 == 0 else BG_LIGHT
    set_cell(comparison_table.cell(i+1, 0), metric, 8, False, TEXT_DARK, PP_ALIGN.LEFT, bg)
    set_cell(comparison_table.cell(i+1, 1), before, 8, True, AZ_RED, PP_ALIGN.CENTER, bg)
    set_cell(comparison_table.cell(i+1, 2), after, 8, True, AZ_GREEN, PP_ALIGN.CENTER, bg)
    set_cell(comparison_table.cell(i+1, 3), improvement, 8, False, AZ_GREEN, PP_ALIGN.LEFT, bg)

# Bottom summary - Potential Value
add_shape(slide_improvements, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.5), AZ_GREEN)
add_text(slide_improvements, "Potential Value: 25-30% manual effort reduction  |  Reusable across engineering & data teams  |  Improved reliability",
         Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.4), 11, True, WHITE, PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 3: Agent Orchestration Architecture
# ============================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

# Header
add_shape(slide2, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.8), AZ_MAROON)
add_text(slide2, "Agent Orchestration Architecture", Inches(0.5), Inches(0.2), Inches(12), Inches(0.5), 24, True, WHITE, PP_ALIGN.LEFT)

# Section title
add_text(slide2, "Multi-Agent System for Autonomous Testing", Inches(0.5), Inches(1.0), Inches(12), Inches(0.4), 14, True, AZ_MAROON, PP_ALIGN.LEFT)

# Agent boxes - Main flow
agents = [
    ("GENERATOR", "Creates Tests", AZ_GREEN),
    ("EXECUTOR", "Runs Tests", AZ_GREEN),
    ("ANALYZER", "Finds Issues", AZ_BLUE),
    ("HEALER", "Auto-Fixes", AZ_GOLD),
    ("OPTIMIZER", "Improves", AZ_MAROON),
]

# Draw agent flow
y_pos = Inches(1.6)
for i, (name, desc, color) in enumerate(agents[:4]):
    left = Inches(0.8 + i * 3.0)
    # Agent box
    box = add_shape(slide2, MSO_SHAPE.ROUNDED_RECTANGLE, left, y_pos, Inches(2.2), Inches(1.2), WHITE, color, 3)
    add_text(slide2, name, left, y_pos + Inches(0.15), Inches(2.2), Inches(0.4), 14, True, color, PP_ALIGN.CENTER)
    add_text(slide2, desc, left, y_pos + Inches(0.55), Inches(2.2), Inches(0.3), 10, False, TEXT_MEDIUM, PP_ALIGN.CENTER)
    add_text(slide2, "Agent", left, y_pos + Inches(0.85), Inches(2.2), Inches(0.25), 9, False, TEXT_LIGHT, PP_ALIGN.CENTER)

    # Arrow
    if i < 3:
        add_shape(slide2, MSO_SHAPE.RIGHT_ARROW, left + Inches(2.3), y_pos + Inches(0.4), Inches(0.6), Inches(0.4), AZ_MAROON)

# Optimizer box (below, connecting to analyzer)
opt_left = Inches(5.3)
opt_top = Inches(3.2)
box = add_shape(slide2, MSO_SHAPE.ROUNDED_RECTANGLE, opt_left, opt_top, Inches(2.2), Inches(1.0), WHITE, AZ_MAROON, 3)
add_text(slide2, "OPTIMIZER", opt_left, opt_top + Inches(0.15), Inches(2.2), Inches(0.35), 14, True, AZ_MAROON, PP_ALIGN.CENTER)
add_text(slide2, "Learns & Improves", opt_left, opt_top + Inches(0.5), Inches(2.2), Inches(0.3), 10, False, TEXT_MEDIUM, PP_ALIGN.CENTER)

# Agent capabilities table
add_text(slide2, "Agent Capabilities", Inches(0.5), Inches(4.5), Inches(12), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.LEFT)

cap_table = slide2.shapes.add_table(6, 4, Inches(0.5), Inches(4.85), Inches(12.3), Inches(2.4)).table
cap_table.columns[0].width = Inches(1.8)
cap_table.columns[1].width = Inches(4.0)
cap_table.columns[2].width = Inches(3.5)
cap_table.columns[3].width = Inches(3.0)

# Header row
headers = ["Agent", "Function", "Input Sources", "Output"]
for i, h in enumerate(headers):
    set_cell(cap_table.cell(0, i), h, 10, True, WHITE, PP_ALIGN.CENTER, AZ_MAROON)

# Data rows
cap_data = [
    ("Generator", "Auto-generates test cases from requirements, code changes, and user behavior", "Requirements, Code Commits, User Journeys", "Test Cases, Test Suites"),
    ("Executor", "Runs tests across 3,000+ device/browser/OS combinations in parallel", "Test Suites, CI/CD Triggers", "Execution Results, Logs"),
    ("Analyzer", "Root cause analysis, pattern detection, flaky test identification", "Test Results, Historical Data", "Defect Reports, Risk Scores"),
    ("Healer", "Auto-fixes broken selectors, adapts to UI changes, self-healing scripts", "Failed Tests, DOM Changes", "Fixed Tests, Change Logs"),
    ("Optimizer", "Learns from results, prioritizes high-risk areas, reduces redundancy", "All Agent Outputs, Metrics", "Optimized Test Suite, Recommendations"),
]

for i, (agent, func, inputs, outputs) in enumerate(cap_data):
    colors = [AZ_GREEN, AZ_GREEN, AZ_BLUE, AZ_GOLD, AZ_MAROON]
    set_cell(cap_table.cell(i+1, 0), agent, 9, True, colors[i], PP_ALIGN.LEFT, WHITE)
    set_cell(cap_table.cell(i+1, 1), func, 8, False, TEXT_DARK, PP_ALIGN.LEFT, WHITE)
    set_cell(cap_table.cell(i+1, 2), inputs, 8, False, TEXT_MEDIUM, PP_ALIGN.LEFT, BG_LIGHT)
    set_cell(cap_table.cell(i+1, 3), outputs, 8, False, TEXT_MEDIUM, PP_ALIGN.LEFT, BG_LIGHT)

# ============================================================================
# SLIDE 3: Dashboard Overview
# ============================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

# Header
add_shape(slide3, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.8), AZ_MAROON)
add_text(slide3, "Agentic AI Testing Dashboard", Inches(0.5), Inches(0.2), Inches(12), Inches(0.5), 24, True, WHITE, PP_ALIGN.LEFT)

# KPI Cards Row
kpi_data = [
    ("AUTONOMY", "78%", "+12%", AZ_MAROON, "Tests without human intervention"),
    ("ACCURACY", "94%", "+3%", AZ_GREEN, "Correct test generation rate"),
    ("SELF-HEAL", "67%", "+15%", AZ_GOLD, "Auto-fixed broken tests"),
    ("COVERAGE", "92%", "+15%", AZ_BLUE, "Code paths tested"),
]

for i, (title, value, trend, color, desc) in enumerate(kpi_data):
    left = Inches(0.4 + i * 3.2)
    top = Inches(1.0)
    box = add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.0), Inches(1.5), WHITE, color, 2)
    add_text(slide3, title, left, top + Inches(0.1), Inches(3.0), Inches(0.3), 10, True, color, PP_ALIGN.CENTER)
    add_text(slide3, value, left, top + Inches(0.35), Inches(3.0), Inches(0.6), 36, True, color, PP_ALIGN.CENTER)
    add_text(slide3, trend, left + Inches(2.2), top + Inches(0.5), Inches(0.6), Inches(0.3), 10, True, AZ_GREEN, PP_ALIGN.CENTER)
    add_text(slide3, desc, left, top + Inches(1.1), Inches(3.0), Inches(0.3), 8, False, TEXT_LIGHT, PP_ALIGN.CENTER)

# Left section: Agent Status
add_text(slide3, "Agent Status (Live)", Inches(0.4), Inches(2.7), Inches(6), Inches(0.3), 11, True, AZ_MAROON, PP_ALIGN.LEFT)

status_table = slide3.shapes.add_table(5, 4, Inches(0.4), Inches(3.05), Inches(6.0), Inches(1.8)).table
status_table.columns[0].width = Inches(1.5)
status_table.columns[1].width = Inches(1.5)
status_table.columns[2].width = Inches(1.5)
status_table.columns[3].width = Inches(1.5)

status_headers = ["Agent", "Status", "Current Task", "Actions/hr"]
for i, h in enumerate(status_headers):
    set_cell(status_table.cell(0, i), h, 9, True, WHITE, PP_ALIGN.CENTER, AZ_MAROON)

status_data = [
    ("Generator", "Active", "GxP Module 3.2", "23"),
    ("Executor", "Running", "1,247 tests", "156"),
    ("Analyzer", "Thinking", "Root cause #2847", "12"),
    ("Healer", "Standby", "Waiting", "8"),
]

for i, (agent, status, task, actions) in enumerate(status_data):
    status_colors = {"Active": AZ_GREEN, "Running": AZ_GREEN, "Thinking": AZ_BLUE, "Standby": TEXT_LIGHT}
    set_cell(status_table.cell(i+1, 0), agent, 9, True, TEXT_DARK, PP_ALIGN.LEFT, WHITE)
    set_cell(status_table.cell(i+1, 1), status, 9, True, status_colors.get(status, TEXT_DARK), PP_ALIGN.CENTER, WHITE)
    set_cell(status_table.cell(i+1, 2), task, 8, False, TEXT_MEDIUM, PP_ALIGN.LEFT, BG_LIGHT)
    set_cell(status_table.cell(i+1, 3), actions, 9, True, AZ_GREEN, PP_ALIGN.CENTER, WHITE)

# Right section: Reasoning Log
add_text(slide3, "Agent Reasoning Log", Inches(6.8), Inches(2.7), Inches(6), Inches(0.3), 11, True, AZ_MAROON, PP_ALIGN.LEFT)

log_box = add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(3.05), Inches(6.0), Inches(1.8), BG_LIGHT, TEXT_LIGHT, 1)
log_entries = [
    "14:32  GENERATOR  Generated 23 GxP tests (94% confidence)",
    "14:28  ANALYZER   Found flaky test root cause: async race",
    "14:15  HEALER     Auto-fixed selector #submit → #submitBtn",
    "14:02  EXECUTOR   Completed Module 2.3 suite (156 tests)",
]
for i, entry in enumerate(log_entries):
    add_text(slide3, entry, Inches(6.95), Inches(3.15 + i * 0.4), Inches(5.7), Inches(0.35), 9, False, TEXT_DARK, PP_ALIGN.LEFT)

# Bottom Left: Human Controls
add_text(slide3, "Human-in-the-Loop Controls", Inches(0.4), Inches(5.0), Inches(6), Inches(0.3), 11, True, AZ_MAROON, PP_ALIGN.LEFT)

control_box = add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(5.35), Inches(6.0), Inches(1.8), WHITE, AZ_MAROON, 2)
add_text(slide3, "Execution Mode:   [Autonomous]  [Supervised]  [Manual]", Inches(0.6), Inches(5.5), Inches(5.6), Inches(0.35), 10, False, TEXT_DARK, PP_ALIGN.LEFT)
add_text(slide3, "Controls:   [ Pause All ]   [ Stop All ]   [ Rollback ]", Inches(0.6), Inches(5.85), Inches(5.6), Inches(0.35), 10, False, TEXT_DARK, PP_ALIGN.LEFT)
add_text(slide3, "Approval Queue: 3 actions pending human review", Inches(0.6), Inches(6.3), Inches(5.6), Inches(0.3), 10, True, AZ_GOLD, PP_ALIGN.LEFT)
add_text(slide3, "   Delete 47 obsolete tests | Modify prod API tests", Inches(0.6), Inches(6.6), Inches(5.6), Inches(0.3), 9, False, TEXT_MEDIUM, PP_ALIGN.LEFT)

# Bottom Right: Risk Analysis
add_text(slide3, "Predictive Risk Analysis", Inches(6.8), Inches(5.0), Inches(6), Inches(0.3), 11, True, AZ_MAROON, PP_ALIGN.LEFT)

risk_box = add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(5.35), Inches(6.0), Inches(1.8), WHITE, AZ_ORANGE, 2)
risk_items = [
    ("HIGH", "Module 3.2.S.4 (Drug Substance)", "87%", AZ_RED),
    ("MED", "Patient Safety Reporting API", "62%", AZ_ORANGE),
    ("LOW", "CDISC ADaM Transformation", "41%", AZ_GOLD),
]
for i, (level, module, risk, color) in enumerate(risk_items):
    add_text(slide3, level, Inches(7.0), Inches(5.5 + i * 0.5), Inches(0.6), Inches(0.35), 9, True, color, PP_ALIGN.CENTER)
    add_text(slide3, module, Inches(7.7), Inches(5.5 + i * 0.5), Inches(3.5), Inches(0.35), 9, False, TEXT_DARK, PP_ALIGN.LEFT)
    add_text(slide3, risk, Inches(11.5), Inches(5.5 + i * 0.5), Inches(1.0), Inches(0.35), 9, True, color, PP_ALIGN.RIGHT)

# ============================================================================
# SLIDE 4: Test Generation Intelligence
# ============================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])

# Header
add_shape(slide4, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.8), AZ_MAROON)
add_text(slide4, "AI Test Generation Intelligence", Inches(0.5), Inches(0.2), Inches(12), Inches(0.5), 24, True, WHITE, PP_ALIGN.LEFT)

# KPI Row
gen_kpis = [
    ("1,247", "Tests Generated", "This Month"),
    ("156", "New This Week", "+23% vs last"),
    ("50%", "Time Saved", "vs Manual"),
    ("98%", "Bug Detection", "Pre-Production"),
]

for i, (value, label, sub) in enumerate(gen_kpis):
    left = Inches(0.4 + i * 3.2)
    box = add_shape(slide4, MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.0), Inches(3.0), Inches(1.2), WHITE, AZ_GREEN, 2)
    add_text(slide4, value, left, Inches(1.1), Inches(3.0), Inches(0.55), 32, True, AZ_GREEN, PP_ALIGN.CENTER)
    add_text(slide4, label, left, Inches(1.65), Inches(3.0), Inches(0.25), 10, True, TEXT_DARK, PP_ALIGN.CENTER)
    add_text(slide4, sub, left, Inches(1.9), Inches(3.0), Inches(0.2), 8, False, TEXT_LIGHT, PP_ALIGN.CENTER)

# Left: Generation Sources
add_text(slide4, "Test Generation Sources", Inches(0.4), Inches(2.4), Inches(6), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.LEFT)

sources_table = slide4.shapes.add_table(5, 3, Inches(0.4), Inches(2.8), Inches(6.0), Inches(2.0)).table
sources_table.columns[0].width = Inches(2.5)
sources_table.columns[1].width = Inches(1.0)
sources_table.columns[2].width = Inches(2.5)

set_cell(sources_table.cell(0, 0), "Source", 9, True, WHITE, PP_ALIGN.LEFT, AZ_MAROON)
set_cell(sources_table.cell(0, 1), "%", 9, True, WHITE, PP_ALIGN.CENTER, AZ_MAROON)
set_cell(sources_table.cell(0, 2), "Example", 9, True, WHITE, PP_ALIGN.LEFT, AZ_MAROON)

sources_data = [
    ("Requirements Documents", "34%", "Module 2/3 submission specs"),
    ("User Behavior Analytics", "28%", "Click patterns, user journeys"),
    ("Code Changes (Commits)", "22%", "New API endpoints detected"),
    ("Historical Bug Patterns", "16%", "Previous defect regression"),
]

for i, (source, pct, example) in enumerate(sources_data):
    set_cell(sources_table.cell(i+1, 0), source, 9, False, TEXT_DARK, PP_ALIGN.LEFT, WHITE)
    set_cell(sources_table.cell(i+1, 1), pct, 10, True, AZ_GREEN, PP_ALIGN.CENTER, BG_LIGHT)
    set_cell(sources_table.cell(i+1, 2), example, 8, False, TEXT_MEDIUM, PP_ALIGN.LEFT, WHITE)

# Right: Test Types Generated
add_text(slide4, "Test Types Generated", Inches(6.8), Inches(2.4), Inches(6), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.LEFT)

types_table = slide4.shapes.add_table(5, 3, Inches(6.8), Inches(2.8), Inches(6.0), Inches(2.0)).table
types_table.columns[0].width = Inches(2.5)
types_table.columns[1].width = Inches(1.5)
types_table.columns[2].width = Inches(2.0)

set_cell(types_table.cell(0, 0), "Test Type", 9, True, WHITE, PP_ALIGN.LEFT, AZ_MAROON)
set_cell(types_table.cell(0, 1), "Count", 9, True, WHITE, PP_ALIGN.CENTER, AZ_MAROON)
set_cell(types_table.cell(0, 2), "Coverage", 9, True, WHITE, PP_ALIGN.CENTER, AZ_MAROON)

types_data = [
    ("Positive Path Tests", "512 (41%)", "Happy flows"),
    ("Negative Path Tests", "387 (31%)", "Error handling"),
    ("Edge Case Tests", "223 (18%)", "Boundary conditions"),
    ("Security Tests", "125 (10%)", "Injection, auth"),
]

for i, (ttype, count, coverage) in enumerate(types_data):
    colors = [AZ_GREEN, AZ_ORANGE, AZ_GOLD, AZ_RED]
    set_cell(types_table.cell(i+1, 0), ttype, 9, False, TEXT_DARK, PP_ALIGN.LEFT, WHITE)
    set_cell(types_table.cell(i+1, 1), count, 9, True, colors[i], PP_ALIGN.CENTER, BG_LIGHT)
    set_cell(types_table.cell(i+1, 2), coverage, 8, False, TEXT_MEDIUM, PP_ALIGN.CENTER, WHITE)

# Bottom: Example Generated Test
add_text(slide4, "Example: AI-Generated Test Case", Inches(0.4), Inches(5.0), Inches(12), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.LEFT)

example_box = add_shape(slide4, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(5.4), Inches(12.5), Inches(1.8), BG_LIGHT, AZ_MAROON, 2)

add_text(slide4, "Test: GxP_Module3_DrugSubstance_NullPointer_EdgeCase", Inches(0.6), Inches(5.55), Inches(8), Inches(0.3), 11, True, TEXT_DARK, PP_ALIGN.LEFT)
add_text(slide4, "AUTO-GENERATED", Inches(10.5), Inches(5.55), Inches(2.2), Inches(0.3), 9, True, WHITE, PP_ALIGN.CENTER)
add_shape(slide4, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.4), Inches(5.5), Inches(2.3), Inches(0.35), AZ_GREEN)
add_text(slide4, "AUTO-GENERATED", Inches(10.5), Inches(5.55), Inches(2.2), Inches(0.3), 9, True, WHITE, PP_ALIGN.CENTER)

details = [
    "Source:      Code commit #a3f2d91 - New null check path in DrugSubstanceValidator.java",
    "Rationale:   Agent detected unguarded null pointer dereference with no existing test coverage",
    "Priority:    HIGH - Regulatory impact (GxP compliance requirement)",
    "Status:      Generated → Executed → PASSED",
]
for i, detail in enumerate(details):
    add_text(slide4, detail, Inches(0.6), Inches(5.95 + i * 0.35), Inches(12), Inches(0.3), 9, False, TEXT_DARK, PP_ALIGN.LEFT)

# ============================================================================
# SLIDE 5: Self-Healing & Maintenance
# ============================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])

# Header
add_shape(slide5, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.8), AZ_MAROON)
add_text(slide5, "Self-Healing Test Automation", Inches(0.5), Inches(0.2), Inches(12), Inches(0.5), 24, True, WHITE, PP_ALIGN.LEFT)

# Stats row
heal_stats = [
    ("67%", "Self-Heal Rate", "Tests auto-fixed"),
    ("40%", "Maintenance Reduction", "vs traditional automation"),
    ("2.3s", "Avg Fix Time", "Instant remediation"),
    ("156", "Auto-Healed", "This month"),
]

for i, (value, label, sub) in enumerate(heal_stats):
    left = Inches(0.4 + i * 3.2)
    color = [AZ_GOLD, AZ_GREEN, AZ_BLUE, AZ_GOLD][i]
    box = add_shape(slide5, MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.0), Inches(3.0), Inches(1.2), WHITE, color, 2)
    add_text(slide5, value, left, Inches(1.1), Inches(3.0), Inches(0.55), 32, True, color, PP_ALIGN.CENTER)
    add_text(slide5, label, left, Inches(1.65), Inches(3.0), Inches(0.25), 10, True, TEXT_DARK, PP_ALIGN.CENTER)
    add_text(slide5, sub, left, Inches(1.9), Inches(3.0), Inches(0.2), 8, False, TEXT_LIGHT, PP_ALIGN.CENTER)

# How Self-Healing Works
add_text(slide5, "How the Healer Agent Works", Inches(0.4), Inches(2.4), Inches(12), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.LEFT)

heal_table = slide5.shapes.add_table(5, 4, Inches(0.4), Inches(2.8), Inches(12.5), Inches(1.8)).table
heal_table.columns[0].width = Inches(0.8)
heal_table.columns[1].width = Inches(3.0)
heal_table.columns[2].width = Inches(4.5)
heal_table.columns[3].width = Inches(4.2)

set_cell(heal_table.cell(0, 0), "Step", 9, True, WHITE, PP_ALIGN.CENTER, AZ_MAROON)
set_cell(heal_table.cell(0, 1), "Action", 9, True, WHITE, PP_ALIGN.LEFT, AZ_MAROON)
set_cell(heal_table.cell(0, 2), "Description", 9, True, WHITE, PP_ALIGN.LEFT, AZ_MAROON)
set_cell(heal_table.cell(0, 3), "Example", 9, True, WHITE, PP_ALIGN.LEFT, AZ_MAROON)

heal_data = [
    ("1", "Detect Failure", "Agent monitors test execution for failures", "Test #2831 failed: element not found"),
    ("2", "Analyze Cause", "AI analyzes DOM changes, screenshots, logs", "Button ID changed: #submit → #submitBtn"),
    ("3", "Generate Fix", "Agent creates updated selector/script", "Updated XPath with visual context"),
    ("4", "Verify & Deploy", "Re-run test, confirm pass, deploy fix", "Test passed Healed selector deployed"),
]

for i, (step, action, desc, example) in enumerate(heal_data):
    set_cell(heal_table.cell(i+1, 0), step, 10, True, AZ_MAROON, PP_ALIGN.CENTER, BG_LIGHT)
    set_cell(heal_table.cell(i+1, 1), action, 9, True, TEXT_DARK, PP_ALIGN.LEFT, WHITE)
    set_cell(heal_table.cell(i+1, 2), desc, 8, False, TEXT_MEDIUM, PP_ALIGN.LEFT, WHITE)
    set_cell(heal_table.cell(i+1, 3), example, 8, False, TEXT_MEDIUM, PP_ALIGN.LEFT, BG_LIGHT)

# Healing Types
add_text(slide5, "Types of Auto-Healing", Inches(0.4), Inches(4.8), Inches(6), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.LEFT)

healing_types_table = slide5.shapes.add_table(5, 2, Inches(0.4), Inches(5.15), Inches(6.0), Inches(2.0)).table
healing_types_table.columns[0].width = Inches(2.5)
healing_types_table.columns[1].width = Inches(3.5)

set_cell(healing_types_table.cell(0, 0), "Healing Type", 9, True, WHITE, PP_ALIGN.LEFT, AZ_GOLD)
set_cell(healing_types_table.cell(0, 1), "What It Fixes", 9, True, WHITE, PP_ALIGN.LEFT, AZ_GOLD)

healing_data = [
    ("Selector Healing", "Element IDs, classes, XPaths that changed"),
    ("Visual Healing", "UI layout changes using image recognition"),
    ("Timing Healing", "Async issues, race conditions, waits"),
    ("Data Healing", "Test data that became stale or invalid"),
]

for i, (htype, fixes) in enumerate(healing_data):
    set_cell(healing_types_table.cell(i+1, 0), htype, 9, True, TEXT_DARK, PP_ALIGN.LEFT, WHITE)
    set_cell(healing_types_table.cell(i+1, 1), fixes, 8, False, TEXT_MEDIUM, PP_ALIGN.LEFT, BG_LIGHT)

# Before/After
add_text(slide5, "Impact: Before vs After", Inches(6.8), Inches(4.8), Inches(6), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.LEFT)

impact_table = slide5.shapes.add_table(4, 3, Inches(6.8), Inches(5.15), Inches(6.0), Inches(2.0)).table
impact_table.columns[0].width = Inches(2.5)
impact_table.columns[1].width = Inches(1.75)
impact_table.columns[2].width = Inches(1.75)

set_cell(impact_table.cell(0, 0), "Metric", 9, True, WHITE, PP_ALIGN.LEFT, AZ_MAROON)
set_cell(impact_table.cell(0, 1), "Before", 9, True, WHITE, PP_ALIGN.CENTER, AZ_RED)
set_cell(impact_table.cell(0, 2), "After", 9, True, WHITE, PP_ALIGN.CENTER, AZ_GREEN)

impact_data = [
    ("Test Maintenance Time", "40 hrs/week", "24 hrs/week"),
    ("Flaky Test Rate", "15%", "3%"),
    ("Time to Fix Broken Test", "4 hours avg", "2.3 seconds"),
]

for i, (metric, before, after) in enumerate(impact_data):
    set_cell(impact_table.cell(i+1, 0), metric, 9, False, TEXT_DARK, PP_ALIGN.LEFT, WHITE)
    set_cell(impact_table.cell(i+1, 1), before, 9, True, AZ_RED, PP_ALIGN.CENTER, BG_LIGHT)
    set_cell(impact_table.cell(i+1, 2), after, 9, True, AZ_GREEN, PP_ALIGN.CENTER, BG_LIGHT)

# ============================================================================
# SLIDE 6: Cost & ROI
# ============================================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])

# Header
add_shape(slide6, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.8), AZ_MAROON)
add_text(slide6, "Cost Tracking & ROI Analysis", Inches(0.5), Inches(0.2), Inches(12), Inches(0.5), 24, True, WHITE, PP_ALIGN.LEFT)

# Top KPIs
roi_kpis = [
    ("$5,550", "Monthly Cost", "56% of budget"),
    ("$7,500", "Budget Target", "On Track"),
    ("$48K", "Annual Savings", "ROI: 3.2x"),
    ("$3.37", "Cost per Test", "-18% vs last month"),
]

for i, (value, label, sub) in enumerate(roi_kpis):
    left = Inches(0.4 + i * 3.2)
    color = [AZ_GOLD, AZ_MAROON, AZ_GREEN, AZ_BLUE][i]
    box = add_shape(slide6, MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.0), Inches(3.0), Inches(1.2), WHITE, color, 2)
    add_text(slide6, value, left, Inches(1.1), Inches(3.0), Inches(0.55), 32, True, color, PP_ALIGN.CENTER)
    add_text(slide6, label, left, Inches(1.65), Inches(3.0), Inches(0.25), 10, True, TEXT_DARK, PP_ALIGN.CENTER)
    add_text(slide6, sub, left, Inches(1.9), Inches(3.0), Inches(0.2), 8, False, TEXT_LIGHT, PP_ALIGN.CENTER)

# Cost Breakdown by Agent
add_text(slide6, "Cost Breakdown by Agent", Inches(0.4), Inches(2.4), Inches(6), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.LEFT)

cost_table = slide6.shapes.add_table(6, 4, Inches(0.4), Inches(2.8), Inches(6.0), Inches(2.2)).table
cost_table.columns[0].width = Inches(1.5)
cost_table.columns[1].width = Inches(1.5)
cost_table.columns[2].width = Inches(1.5)
cost_table.columns[3].width = Inches(1.5)

set_cell(cost_table.cell(0, 0), "Agent", 9, True, WHITE, PP_ALIGN.LEFT, AZ_MAROON)
set_cell(cost_table.cell(0, 1), "Tokens (M)", 9, True, WHITE, PP_ALIGN.CENTER, AZ_MAROON)
set_cell(cost_table.cell(0, 2), "Cost", 9, True, WHITE, PP_ALIGN.CENTER, AZ_MAROON)
set_cell(cost_table.cell(0, 3), "% of Total", 9, True, WHITE, PP_ALIGN.CENTER, AZ_MAROON)

cost_data = [
    ("Generator", "1.4M", "$1,840", "44%"),
    ("Analyzer", "0.8M", "$1,050", "25%"),
    ("Executor", "0.5M", "$680", "16%"),
    ("Healer", "0.3M", "$420", "10%"),
    ("Optimizer", "0.2M", "$210", "5%"),
]

for i, (agent, tokens, cost, pct) in enumerate(cost_data):
    colors = [AZ_GREEN, AZ_BLUE, AZ_GREEN, AZ_GOLD, AZ_MAROON]
    set_cell(cost_table.cell(i+1, 0), agent, 9, True, colors[i], PP_ALIGN.LEFT, WHITE)
    set_cell(cost_table.cell(i+1, 1), tokens, 9, False, TEXT_DARK, PP_ALIGN.CENTER, BG_LIGHT)
    set_cell(cost_table.cell(i+1, 2), cost, 9, True, AZ_GOLD, PP_ALIGN.CENTER, WHITE)
    set_cell(cost_table.cell(i+1, 3), pct, 9, False, TEXT_MEDIUM, PP_ALIGN.CENTER, BG_LIGHT)

# Infrastructure Costs
add_text(slide6, "Infrastructure Costs", Inches(6.8), Inches(2.4), Inches(6), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.LEFT)

infra_table = slide6.shapes.add_table(4, 2, Inches(6.8), Inches(2.8), Inches(6.0), Inches(1.5)).table
infra_table.columns[0].width = Inches(4.0)
infra_table.columns[1].width = Inches(2.0)

set_cell(infra_table.cell(0, 0), "Service", 9, True, WHITE, PP_ALIGN.LEFT, AZ_MAROON)
set_cell(infra_table.cell(0, 1), "Monthly Cost", 9, True, WHITE, PP_ALIGN.CENTER, AZ_MAROON)

infra_data = [
    ("LLM API (AWS Bedrock - Claude)", "$4,200"),
    ("Embedding Service (Titan)", "$850"),
    ("Infrastructure (Lambda, S3, etc.)", "$500"),
]

for i, (service, cost) in enumerate(infra_data):
    set_cell(infra_table.cell(i+1, 0), service, 9, False, TEXT_DARK, PP_ALIGN.LEFT, WHITE)
    set_cell(infra_table.cell(i+1, 1), cost, 10, True, AZ_GOLD, PP_ALIGN.CENTER, BG_LIGHT)

# ROI Breakdown
add_text(slide6, "ROI Analysis", Inches(6.8), Inches(4.5), Inches(6), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.LEFT)

roi_table = slide6.shapes.add_table(5, 2, Inches(6.8), Inches(4.9), Inches(6.0), Inches(2.0)).table
roi_table.columns[0].width = Inches(4.0)
roi_table.columns[1].width = Inches(2.0)

set_cell(roi_table.cell(0, 0), "Savings Category", 9, True, WHITE, PP_ALIGN.LEFT, AZ_GREEN)
set_cell(roi_table.cell(0, 1), "Annual Value", 9, True, WHITE, PP_ALIGN.CENTER, AZ_GREEN)

roi_data = [
    ("Reduced Manual Test Creation", "$24,000"),
    ("Faster Release Cycles", "$12,000"),
    ("Self-Healing (Maintenance)", "$8,000"),
    ("Defect Prevention", "$4,000"),
]

for i, (category, value) in enumerate(roi_data):
    set_cell(roi_table.cell(i+1, 0), category, 9, False, TEXT_DARK, PP_ALIGN.LEFT, WHITE)
    set_cell(roi_table.cell(i+1, 1), value, 10, True, AZ_GREEN, PP_ALIGN.CENTER, BG_LIGHT)

# Total ROI Box
add_shape(slide6, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(5.2), Inches(6.0), Inches(1.6), AZ_GREEN)
add_text(slide6, "TOTAL ANNUAL ROI", Inches(0.4), Inches(5.4), Inches(6.0), Inches(0.4), 14, True, WHITE, PP_ALIGN.CENTER)
add_text(slide6, "$48,000", Inches(0.4), Inches(5.8), Inches(6.0), Inches(0.6), 40, True, WHITE, PP_ALIGN.CENTER)
add_text(slide6, "3.2x Return on Investment", Inches(0.4), Inches(6.4), Inches(6.0), Inches(0.3), 12, False, WHITE, PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 7: Platform Integration & Status
# ============================================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])

# Header
add_shape(slide7, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.8), AZ_MAROON)
add_text(slide7, "Platform Integration & Observability", Inches(0.5), Inches(0.2), Inches(12), Inches(0.5), 24, True, WHITE, PP_ALIGN.LEFT)

# Platform Status
add_text(slide7, "Integrated Platforms", Inches(0.4), Inches(1.0), Inches(6), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.LEFT)

platform_table = slide7.shapes.add_table(4, 4, Inches(0.4), Inches(1.4), Inches(6.0), Inches(1.8)).table
platform_table.columns[0].width = Inches(1.2)
platform_table.columns[1].width = Inches(2.3)
platform_table.columns[2].width = Inches(1.0)
platform_table.columns[3].width = Inches(1.5)

set_cell(platform_table.cell(0, 0), "Platform", 9, True, WHITE, PP_ALIGN.LEFT, AZ_MAROON)
set_cell(platform_table.cell(0, 1), "Description", 9, True, WHITE, PP_ALIGN.LEFT, AZ_MAROON)
set_cell(platform_table.cell(0, 2), "Status", 9, True, WHITE, PP_ALIGN.CENTER, AZ_MAROON)
set_cell(platform_table.cell(0, 3), "Pass Rate", 9, True, WHITE, PP_ALIGN.CENTER, AZ_MAROON)

platforms = [
    ("3DP", "Drug Development Platform", "LIVE", "98.2%"),
    ("BIKG", "Biological Intelligence KG", "LIVE", "97.8%"),
    ("CDISC", "Clinical Data Standards", "STAGING", "94.1%"),
]

for i, (name, desc, status, rate) in enumerate(platforms):
    status_color = AZ_GREEN if status == "LIVE" else AZ_GOLD
    set_cell(platform_table.cell(i+1, 0), name, 10, True, AZ_MAROON, PP_ALIGN.LEFT, WHITE)
    set_cell(platform_table.cell(i+1, 1), desc, 8, False, TEXT_MEDIUM, PP_ALIGN.LEFT, WHITE)
    set_cell(platform_table.cell(i+1, 2), status, 9, True, WHITE, PP_ALIGN.CENTER, status_color)
    set_cell(platform_table.cell(i+1, 3), rate, 10, True, AZ_GREEN, PP_ALIGN.CENTER, BG_LIGHT)

# CI/CD Pipeline
add_text(slide7, "CI/CD Pipeline Integration", Inches(6.8), Inches(1.0), Inches(6), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.LEFT)

pipeline_box = add_shape(slide7, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4), Inches(6.0), Inches(1.8), WHITE, AZ_BLUE, 2)
add_text(slide7, "main  →  build  →  test  →  deploy", Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.5), 16, True, TEXT_DARK, PP_ALIGN.CENTER)
add_text(slide7, "   ✓         ✓          ●          ○", Inches(7.0), Inches(2.1), Inches(5.6), Inches(0.4), 20, True, AZ_GREEN, PP_ALIGN.CENTER)
add_text(slide7, "Last successful deploy: 2 hours ago", Inches(7.0), Inches(2.6), Inches(5.6), Inches(0.3), 10, False, TEXT_MEDIUM, PP_ALIGN.CENTER)
add_text(slide7, "Current: Running validation suite (1,247 tests)", Inches(7.0), Inches(2.9), Inches(5.6), Inches(0.3), 9, False, AZ_BLUE, PP_ALIGN.CENTER)

# Observability Features
add_text(slide7, "Agent Observability Features", Inches(0.4), Inches(3.4), Inches(12), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.LEFT)

obs_table = slide7.shapes.add_table(5, 3, Inches(0.4), Inches(3.8), Inches(12.5), Inches(1.8)).table
obs_table.columns[0].width = Inches(2.5)
obs_table.columns[1].width = Inches(5.5)
obs_table.columns[2].width = Inches(4.5)

set_cell(obs_table.cell(0, 0), "Feature", 9, True, WHITE, PP_ALIGN.LEFT, AZ_MAROON)
set_cell(obs_table.cell(0, 1), "Description", 9, True, WHITE, PP_ALIGN.LEFT, AZ_MAROON)
set_cell(obs_table.cell(0, 2), "Benefit", 9, True, WHITE, PP_ALIGN.LEFT, AZ_MAROON)

obs_data = [
    ("Reasoning Traces", "Full visibility into agent decision-making process", "Debug AI behavior, explain outcomes"),
    ("Real-time Dashboards", "Live metrics on agent status, performance, costs", "Immediate issue detection"),
    ("Alerting & Notifications", "Configurable alerts for anomalies and failures", "Proactive incident response"),
    ("Audit Logs", "Complete history of agent actions for compliance", "Regulatory traceability (GxP)"),
]

for i, (feature, desc, benefit) in enumerate(obs_data):
    set_cell(obs_table.cell(i+1, 0), feature, 9, True, TEXT_DARK, PP_ALIGN.LEFT, WHITE)
    set_cell(obs_table.cell(i+1, 1), desc, 8, False, TEXT_MEDIUM, PP_ALIGN.LEFT, BG_LIGHT)
    set_cell(obs_table.cell(i+1, 2), benefit, 8, False, TEXT_MEDIUM, PP_ALIGN.LEFT, WHITE)

# Model Performance
add_text(slide7, "Model Performance Metrics", Inches(0.4), Inches(5.8), Inches(6), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.LEFT)

perf_table = slide7.shapes.add_table(4, 2, Inches(0.4), Inches(6.2), Inches(6.0), Inches(1.1)).table
perf_table.columns[0].width = Inches(3.0)
perf_table.columns[1].width = Inches(3.0)

perfs = [
    ("Test Accuracy", "94%"),
    ("Data Validity", "91%"),
    ("User Satisfaction", "4.2/5 (87%)"),
]

set_cell(perf_table.cell(0, 0), "Metric", 9, True, WHITE, PP_ALIGN.LEFT, AZ_GREEN)
set_cell(perf_table.cell(0, 1), "Score", 9, True, WHITE, PP_ALIGN.CENTER, AZ_GREEN)

for i, (metric, score) in enumerate(perfs):
    set_cell(perf_table.cell(i+1, 0), metric, 9, False, TEXT_DARK, PP_ALIGN.LEFT, WHITE)
    set_cell(perf_table.cell(i+1, 1), score, 10, True, AZ_GREEN, PP_ALIGN.CENTER, BG_LIGHT)

# Governance
add_text(slide7, "Governance & Compliance", Inches(6.8), Inches(5.8), Inches(6), Inches(0.3), 12, True, AZ_MAROON, PP_ALIGN.LEFT)

gov_box = add_shape(slide7, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(6.2), Inches(6.0), Inches(1.1), BG_LIGHT, AZ_MAROON, 1)
gov_items = [
    "GxP Compliance: All agent actions logged and auditable",
    "Human Oversight: Critical actions require approval",
    "Data Privacy: No PII in training or inference",
]
for i, item in enumerate(gov_items):
    add_text(slide7, item, Inches(7.0), Inches(6.35 + i * 0.3), Inches(5.6), Inches(0.25), 9, False, TEXT_DARK, PP_ALIGN.LEFT)

# ============================================================================
# SLIDE 8: Summary & Next Steps
# ============================================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])

# Header
add_shape(slide8, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.8), AZ_MAROON)
add_text(slide8, "Summary & Next Steps", Inches(0.5), Inches(0.2), Inches(12), Inches(0.5), 24, True, WHITE, PP_ALIGN.LEFT)

# Key Achievements
add_text(slide8, "Key Achievements - Phase 1", Inches(0.4), Inches(1.0), Inches(6), Inches(0.3), 14, True, AZ_MAROON, PP_ALIGN.LEFT)

achievements = [
    ("5 AI Agents deployed", "Generator, Executor, Analyzer, Healer, Optimizer"),
    ("78% Autonomy Rate", "Tests running without human intervention"),
    ("94% Test Accuracy", "AI-generated tests catching real bugs"),
    ("67% Self-Healing", "Broken tests auto-fixed instantly"),
    ("$48K Annual Savings", "3.2x ROI on investment"),
    ("2 Platforms Live", "3DP and BIKG fully integrated"),
]

for i, (title, desc) in enumerate(achievements):
    y = Inches(1.4 + i * 0.5)
    add_shape(slide8, MSO_SHAPE.OVAL, Inches(0.5), y + Inches(0.05), Inches(0.15), Inches(0.15), AZ_GREEN)
    add_text(slide8, title, Inches(0.8), y, Inches(5.5), Inches(0.25), 11, True, TEXT_DARK, PP_ALIGN.LEFT)
    add_text(slide8, desc, Inches(0.8), y + Inches(0.22), Inches(5.5), Inches(0.25), 9, False, TEXT_MEDIUM, PP_ALIGN.LEFT)

# Next Steps - Phase 2
add_text(slide8, "Next Steps - Phase 2", Inches(6.8), Inches(1.0), Inches(6), Inches(0.3), 14, True, AZ_MAROON, PP_ALIGN.LEFT)

next_steps = [
    ("Expand to CDISC Platform", "Clinical data standards integration"),
    ("Increase Autonomy to 90%", "Reduce human intervention further"),
    ("Add Predictive Testing", "AI predicts where bugs will occur"),
    ("Multi-Agent Collaboration", "Agents working together on complex tasks"),
    ("Natural Language Interface", "Talk to agents to generate tests"),
    ("Production Monitoring Agent", "Detect issues in live systems"),
]

for i, (title, desc) in enumerate(next_steps):
    y = Inches(1.4 + i * 0.5)
    add_shape(slide8, MSO_SHAPE.OVAL, Inches(6.9), y + Inches(0.05), Inches(0.15), Inches(0.15), AZ_GOLD)
    add_text(slide8, title, Inches(7.2), y, Inches(5.5), Inches(0.25), 11, True, TEXT_DARK, PP_ALIGN.LEFT)
    add_text(slide8, desc, Inches(7.2), y + Inches(0.22), Inches(5.5), Inches(0.25), 9, False, TEXT_MEDIUM, PP_ALIGN.LEFT)

# Bottom banner
add_shape(slide8, MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(4.8), Inches(12.5), Inches(2.4), BG_CREAM, AZ_MAROON, 2)

add_text(slide8, "The Future of Testing is Agentic", Inches(0.4), Inches(5.0), Inches(12.5), Inches(0.5), 20, True, AZ_MAROON, PP_ALIGN.CENTER)
add_text(slide8, "AI agents that think, act, and learn autonomously - transforming QA from a bottleneck into a competitive advantage",
         Inches(0.6), Inches(5.5), Inches(12.1), Inches(0.6), 12, False, TEXT_MEDIUM, PP_ALIGN.CENTER, italic=True)

# Stats in bottom banner
final_stats = [
    ("70%", "of enterprises adopting\nAgentic AI by 2025"),
    ("15%", "of work decisions made\nby AI agents by 2028"),
    ("40%", "reduction in testing\ncosts with AI agents"),
]

for i, (stat, label) in enumerate(final_stats):
    left = Inches(1.5 + i * 4.0)
    add_text(slide8, stat, left, Inches(6.1), Inches(2.5), Inches(0.5), 28, True, AZ_GREEN, PP_ALIGN.CENTER)
    add_text(slide8, label, left, Inches(6.55), Inches(2.5), Inches(0.5), 9, False, TEXT_MEDIUM, PP_ALIGN.CENTER)

# Save the presentation
output_path = "/Users/devesh.b.sharma/Astra Zeneca/R&D IT/Agentic-AI-Testing-Dashboard.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
