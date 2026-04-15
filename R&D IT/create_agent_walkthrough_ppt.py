from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ============================================================================
# COLORS
# ============================================================================
AZ_MAROON = RGBColor(0x83, 0x00, 0x51)
AZ_GREEN = RGBColor(0x2E, 0x7D, 0x32)
AZ_LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
AZ_GOLD = RGBColor(0xC4, 0xA0, 0x00)
AZ_LIGHT_GOLD = RGBColor(0xFF, 0xF8, 0xE1)
AZ_BLUE = RGBColor(0x1E, 0x88, 0xE5)
AZ_RED = RGBColor(0xD3, 0x2F, 0x2F)
AZ_LIGHT_RED = RGBColor(0xFF, 0xEB, 0xEE)
AZ_PINK = RGBColor(0xE9, 0x1E, 0x63)
AZ_PURPLE = RGBColor(0x5E, 0x35, 0xB1)
AZ_ORANGE = RGBColor(0xFF, 0x98, 0x00)
AZ_DARK = RGBColor(0x2D, 0x2D, 0x2D)
AZ_MEDIUM = RGBColor(0x55, 0x55, 0x55)
AZ_LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_CREAM = RGBColor(0xFA, 0xF9, 0xF7)

# ============================================================================
# HELPERS
# ============================================================================
def add_text(slide, txt, left, top, width, height, size=12, bold=False, color=AZ_DARK, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return box

def add_multiline_text(slide, lines, left, top, width, height, size=10, color=AZ_DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(1)
        p.space_before = Pt(0)
        if isinstance(line, tuple):
            txt, b, c = line
            p.text = txt
            p.font.size = Pt(size)
            p.font.bold = b
            p.font.color.rgb = c
        else:
            p.text = line
            p.font.size = Pt(size)
            p.font.color.rgb = color
    return box

def add_shape_bg(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def rbox(slide, left, top, width, height, fill, line=None, lw=1):
    return add_shape_bg(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill, line, lw)

def set_cell(cell, text, size=10, bold=False, color=AZ_DARK, align=PP_ALIGN.LEFT, fill=None):
    cell.text = text
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill

# ============================================================================
# CREATE PRESENTATION
# ============================================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================================================
# SLIDE 1: TITLE — The Problem & The Promise (kept as is)
# ============================================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])

add_shape_bg(s1, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), BG_CREAM)
add_shape_bg(s1, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08), AZ_MAROON)
add_shape_bg(s1, MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.5), Inches(13.333), Inches(2.5), WHITE)

title_box = s1.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.8))
tf = title_box.text_frame
tf.word_wrap = True
r1 = tf.paragraphs[0].add_run()
r1.text = "Test Intelligence Agent"
r1.font.size = Pt(36)
r1.font.bold = True
r1.font.color.rgb = AZ_MAROON

add_text(s1, "How AI Agents Automate Testing for R&D Platforms", Inches(0.8), Inches(2.5), Inches(11), Inches(0.5), 18, False, AZ_MEDIUM, PP_ALIGN.LEFT, italic=True)
add_text(s1, "A Real-Life Walkthrough", Inches(0.8), Inches(3.0), Inches(11), Inches(0.4), 14, False, AZ_GOLD, PP_ALIGN.LEFT)

# LEFT: Problem
rbox(s1, Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.5), AZ_LIGHT_RED, AZ_RED, 1.5)
add_text(s1, "TODAY — Manual Testing", Inches(1.0), Inches(4.4), Inches(5.2), Inches(0.35), 14, True, AZ_RED, PP_ALIGN.LEFT)
add_multiline_text(s1, [
    ("QA engineer reads requirement docs", False, AZ_DARK),
    ("Manually writes test cases in Excel", False, AZ_DARK),
    ("Creates test data by hand", False, AZ_DARK),
    ("Runs tests, documents results", False, AZ_DARK),
    ("Writes compliance report for audit", False, AZ_DARK),
    ("", False, AZ_DARK),
    ("~2 weeks per release cycle", True, AZ_RED),
], Inches(1.0), Inches(4.8), Inches(5.2), Inches(2.0), 11)

# RIGHT: Promise
rbox(s1, Inches(6.9), Inches(4.3), Inches(5.6), Inches(2.5), AZ_LIGHT_GREEN, AZ_GREEN, 1.5)
add_text(s1, "WITH AGENT — AI-Powered Testing", Inches(7.1), Inches(4.4), Inches(5.2), Inches(0.35), 14, True, AZ_GREEN, PP_ALIGN.LEFT)
add_multiline_text(s1, [
    ("QA engineer types one sentence", False, AZ_DARK),
    ("Agent reads specs automatically", False, AZ_DARK),
    ("Agent generates test cases + data", False, AZ_DARK),
    ("Agent runs tests, finds bugs", False, AZ_DARK),
    ("Agent writes compliance report", False, AZ_DARK),
    ("", False, AZ_DARK),
    ("~5 minutes, same quality", True, AZ_GREEN),
], Inches(7.1), Inches(4.8), Inches(5.2), Inches(2.0), 11)

# VS circle
add_shape_bg(s1, MSO_SHAPE.OVAL, Inches(6.1), Inches(5.15), Inches(0.6), Inches(0.6), AZ_MAROON)
add_text(s1, "vs", Inches(6.1), Inches(5.25), Inches(0.6), Inches(0.4), 14, True, WHITE, PP_ALIGN.CENTER)

add_text(s1, "1 / 3", Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3), 9, False, AZ_MEDIUM, PP_ALIGN.RIGHT)

# ============================================================================
# SLIDE 2: FULL FLOWCHART — The Entire Agent Pipeline on One Page
# ============================================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])

# Header bar
add_shape_bg(s2, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.6), AZ_MAROON)
add_text(s2, "Test Intelligence Agent — Complete Flow", Inches(0.4), Inches(0.08), Inches(8), Inches(0.4), 20, True, WHITE, PP_ALIGN.LEFT)
add_text(s2, "Real-Life Example: Patient Safety — Adverse Event Reporting", Inches(8.5), Inches(0.12), Inches(4.5), Inches(0.35), 10, False, RGBColor(0xFF, 0xCC, 0xDD), PP_ALIGN.RIGHT, italic=True)

# ── LAYOUT: Vertical flow down the LEFT/CENTER, with annotations on RIGHT ──
# Using compact sizing to fit everything on one page

# ── ROW 1: TRIGGER (QA Engineer) ──
CX = Inches(3.0)   # center-x for the main flow column
BW = Inches(4.5)    # box width for main flow
BH = Inches(0.42)   # standard box height

rbox(s2, CX, Inches(0.72), BW, BH, WHITE, AZ_MAROON, 1.5)
add_text(s2, 'QA Engineer: "Test adverse event reporting for Patient Safety"', CX + Inches(0.1), Inches(0.75), BW - Inches(0.2), BH, 8, True, AZ_DARK, PP_ALIGN.CENTER)

# Arrow
add_text(s2, "▼", CX + Inches(2.1), Inches(1.12), Inches(0.3), Inches(0.22), 10, True, AZ_MAROON, PP_ALIGN.CENTER)

# ── ROW 2: ORCHESTRATOR ──
rbox(s2, CX, Inches(1.32), BW, Inches(0.50), AZ_MAROON)
add_text(s2, "ORCHESTRATOR", CX + Inches(0.1), Inches(1.33), BW - Inches(0.2), Inches(0.22), 10, True, WHITE, PP_ALIGN.CENTER)
add_text(s2, "Platform: Patient Safety  |  Feature: Adverse Event  |  Action: Test", CX + Inches(0.1), Inches(1.53), BW - Inches(0.2), Inches(0.22), 7, False, RGBColor(0xFF, 0xCC, 0xDD), PP_ALIGN.CENTER)

# Arrow
add_text(s2, "▼", CX + Inches(2.1), Inches(1.82), Inches(0.3), Inches(0.2), 10, True, AZ_MAROON, PP_ALIGN.CENTER)

# ── ROW 3: STEP 1 — REQUIREMENT AGENT ──
Y3 = Inches(2.0)
rbox(s2, CX, Y3, BW, Inches(0.65), AZ_BLUE)
add_text(s2, "STEP 1  REQUIREMENT AGENT", CX + Inches(0.1), Y3 + Inches(0.02), BW - Inches(0.2), Inches(0.2), 9, True, WHITE, PP_ALIGN.LEFT)
add_text(s2, "Reads specs automatically. Output: 4 mandatory fields\n(Patient ID, Drug, Severity, Description). Severity = enum.", CX + Inches(0.1), Y3 + Inches(0.22), BW - Inches(0.2), Inches(0.4), 7, False, WHITE, PP_ALIGN.LEFT)
# Time badge
rbox(s2, Inches(7.7), Y3 + Inches(0.1), Inches(1.3), Inches(0.25), AZ_GREEN)
add_text(s2, "30s vs 3 days", Inches(7.7), Y3 + Inches(0.1), Inches(1.3), Inches(0.25), 7, True, WHITE, PP_ALIGN.CENTER)

# Arrow
add_text(s2, "▼", CX + Inches(2.1), Inches(2.65), Inches(0.3), Inches(0.2), 10, True, AZ_MAROON, PP_ALIGN.CENTER)

# ── ROW 4: STEP 2 — TEST CASE GEN ──
Y4 = Inches(2.82)
rbox(s2, CX, Y4, BW, Inches(0.65), AZ_PURPLE)
add_text(s2, "STEP 2  TEST CASE GENERATION", CX + Inches(0.1), Y4 + Inches(0.02), BW - Inches(0.2), Inches(0.2), 9, True, WHITE, PP_ALIGN.LEFT)
add_text(s2, "TC1: Valid event  TC2: Missing ID  TC3: Invalid severity\nTC4: Duplicate  TC5: Special chars  TC6: Bad patient ID", CX + Inches(0.1), Y4 + Inches(0.22), BW - Inches(0.2), Inches(0.4), 7, False, WHITE, PP_ALIGN.LEFT)
# Time badge
rbox(s2, Inches(7.7), Y4 + Inches(0.1), Inches(1.3), Inches(0.25), AZ_GREEN)
add_text(s2, "30s vs 3 days", Inches(7.7), Y4 + Inches(0.1), Inches(1.3), Inches(0.25), 7, True, WHITE, PP_ALIGN.CENTER)

# Arrow
add_text(s2, "▼", CX + Inches(2.1), Inches(3.47), Inches(0.3), Inches(0.2), 10, True, AZ_MAROON, PP_ALIGN.CENTER)

# ── ROW 5: DECISION — Need data? ──
Y5 = Inches(3.65)
rbox(s2, CX + Inches(0.75), Y5, Inches(3.0), Inches(0.35), AZ_LIGHT_GOLD, AZ_GOLD, 1.5)
add_text(s2, "DECISION: Need test data?", CX + Inches(0.85), Y5 + Inches(0.03), Inches(2.8), Inches(0.3), 8, True, AZ_GOLD, PP_ALIGN.CENTER)

# YES arrow going down
add_text(s2, "YES ▼", CX + Inches(1.0), Inches(3.98), Inches(0.6), Inches(0.2), 7, True, AZ_GOLD, PP_ALIGN.CENTER)
# NO arrow going right
add_text(s2, "NO ►  Skip to Step 4", CX + Inches(3.85), Y5 + Inches(0.05), Inches(1.8), Inches(0.25), 7, True, AZ_MEDIUM, PP_ALIGN.LEFT)

# ── ROW 6: STEP 3 — SYNTHETIC DATA (conditional) ──
Y6 = Inches(4.18)
rbox(s2, CX, Y6, BW, Inches(0.50), AZ_GOLD)
add_text(s2, "STEP 3  SYNTHETIC DATA AGENT  (if needed)", CX + Inches(0.1), Y6 + Inches(0.02), BW - Inches(0.2), Inches(0.2), 9, True, WHITE, PP_ALIGN.LEFT)
add_text(s2, "Patient: PAT-2847  |  Drug: Tagrisso  |  Severity: Serious  (fake, schema-valid)", CX + Inches(0.1), Y6 + Inches(0.23), BW - Inches(0.2), Inches(0.25), 7, False, WHITE, PP_ALIGN.LEFT)
# Time badge
rbox(s2, Inches(7.7), Y6 + Inches(0.1), Inches(1.3), Inches(0.25), AZ_GREEN)
add_text(s2, "15s vs 1 day", Inches(7.7), Y6 + Inches(0.1), Inches(1.3), Inches(0.25), 7, True, WHITE, PP_ALIGN.CENTER)

# Arrow
add_text(s2, "▼", CX + Inches(2.1), Inches(4.68), Inches(0.3), Inches(0.2), 10, True, AZ_MAROON, PP_ALIGN.CENTER)

# ── ROW 7: STEP 4 — EXECUTION ──
Y7 = Inches(4.86)
rbox(s2, CX, Y7, BW, Inches(0.50), AZ_PINK)
add_text(s2, "STEP 4  EXECUTION AGENT", CX + Inches(0.1), Y7 + Inches(0.02), BW - Inches(0.2), Inches(0.2), 9, True, WHITE, PP_ALIGN.LEFT)
add_text(s2, "TC1: PASS  TC2: PASS  TC3: FAIL  TC4: PASS  TC5: PASS  TC6: PASS", CX + Inches(0.1), Y7 + Inches(0.23), BW - Inches(0.2), Inches(0.25), 7, False, WHITE, PP_ALIGN.LEFT)
# Time badge
rbox(s2, Inches(7.7), Y7 + Inches(0.1), Inches(1.3), Inches(0.25), AZ_GREEN)
add_text(s2, "2 min vs 2 days", Inches(7.7), Y7 + Inches(0.1), Inches(1.3), Inches(0.25), 7, True, WHITE, PP_ALIGN.CENTER)

# Arrow
add_text(s2, "▼", CX + Inches(2.1), Inches(5.36), Inches(0.3), Inches(0.2), 10, True, AZ_MAROON, PP_ALIGN.CENTER)

# ── ROW 8: DECISION — All pass? ──
Y8 = Inches(5.52)
rbox(s2, CX + Inches(0.75), Y8, Inches(3.0), Inches(0.35), AZ_LIGHT_GOLD, AZ_GOLD, 1.5)
add_text(s2, "DECISION: All tests pass?", CX + Inches(0.85), Y8 + Inches(0.03), Inches(2.8), Inches(0.3), 8, True, AZ_GOLD, PP_ALIGN.CENTER)

# YES going right
add_text(s2, "YES ►  Skip to Step 7", CX + Inches(3.85), Y8 + Inches(0.05), Inches(1.8), Inches(0.25), 7, True, AZ_GREEN, PP_ALIGN.LEFT)
# NO going down
add_text(s2, "NO ▼", CX + Inches(1.0), Inches(5.86), Inches(0.6), Inches(0.2), 7, True, AZ_RED, PP_ALIGN.CENTER)

# ── ROW 9: STEP 5 — FAILURE ANALYSIS (conditional) ──
Y9 = Inches(6.04)
rbox(s2, CX, Y9, BW, Inches(0.50), AZ_RED)
add_text(s2, "STEP 5  FAILURE ANALYSIS  (if tests fail)", CX + Inches(0.1), Y9 + Inches(0.02), BW - Inches(0.2), Inches(0.2), 9, True, WHITE, PP_ALIGN.LEFT)
add_text(s2, "TC3: No severity validation — accepts any string. ROOT CAUSE: Code bug. RISK: Regulatory!", CX + Inches(0.1), Y9 + Inches(0.23), BW - Inches(0.2), Inches(0.25), 7, False, WHITE, PP_ALIGN.LEFT)

# Arrow to decision
add_text(s2, "▼", CX + Inches(2.1), Inches(6.54), Inches(0.3), Inches(0.18), 10, True, AZ_MAROON, PP_ALIGN.CENTER)

# ── ROW 10: DECISION — Code bug or test bug? + STEP 6 CODE REFACTOR ──
Y10 = Inches(6.7)
rbox(s2, CX + Inches(0.75), Y10, Inches(1.8), Inches(0.32), AZ_LIGHT_GOLD, AZ_GOLD, 1.5)
add_text(s2, "Code or test bug?", CX + Inches(0.85), Y10 + Inches(0.03), Inches(1.6), Inches(0.28), 7, True, AZ_GOLD, PP_ALIGN.CENTER)

# Code bug arrow to Step 6
add_text(s2, "CODE BUG ►", CX + Inches(2.65), Y10 + Inches(0.04), Inches(1.0), Inches(0.25), 7, True, AZ_RED, PP_ALIGN.LEFT)

# Step 6 box
rbox(s2, CX + Inches(3.5), Y10 - Inches(0.02), Inches(1.8), Inches(0.36), AZ_ORANGE)
add_text(s2, "STEP 6  CODE REFACTOR", CX + Inches(3.55), Y10, Inches(1.7), Inches(0.16), 7, True, WHITE, PP_ALIGN.CENTER)
add_text(s2, "Suggest fix → re-run", CX + Inches(3.55), Y10 + Inches(0.16), Inches(1.7), Inches(0.16), 6, False, WHITE, PP_ALIGN.CENTER)

# Test bug path label
add_text(s2, "TEST BUG → Fix test, re-run", CX + Inches(0.05), Y10 + Inches(0.05), Inches(1.5), Inches(0.25), 6, False, AZ_BLUE, PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# RIGHT COLUMN: STEP 7 REPORTING + BOTTOM LINE TABLE
# ══════════════════════════════════════════════════════════════════════════════

RX = Inches(9.3)   # right column x
RW = Inches(3.8)   # right column width

# ── STEP 7: REPORTING AGENT ──
rbox(s2, RX, Inches(0.72), RW, Inches(2.4), AZ_LIGHT_GREEN, AZ_GREEN, 1.5)
rbox(s2, RX, Inches(0.72), RW, Inches(0.32), AZ_GREEN)
add_text(s2, "STEP 7  REPORTING AGENT", RX + Inches(0.1), Inches(0.74), RW - Inches(0.2), Inches(0.28), 9, True, WHITE, PP_ALIGN.LEFT)
add_text(s2, "(always runs at end)", RX + Inches(2.2), Inches(0.74), Inches(1.4), Inches(0.28), 7, False, WHITE, PP_ALIGN.RIGHT)

# Mini report inside
report_lines = [
    ("TEST EVIDENCE REPORT", True, AZ_GREEN),
    ("", False, AZ_DARK),
    ("Platform:    Patient Safety", False, AZ_DARK),
    ("Tests:         6 generated, 6/6 passed", False, AZ_GREEN),
    ("Bug Found:  Severity validation missing", False, AZ_RED),
    ("Risk:          HIGH — Regulatory", False, AZ_RED),
    ("Fix:            Enum validation added", False, AZ_GREEN),
    ("Compliance: GxP audit trail ready", False, AZ_GREEN),
    ("", False, AZ_DARK),
    ("Time: 4 min 32 sec (zero human effort)", True, AZ_MAROON),
]
add_multiline_text(s2, report_lines, RX + Inches(0.15), Inches(1.08), RW - Inches(0.3), Inches(1.9), 7, AZ_DARK)

# ── ARROW from reporting to bottom line ──
add_text(s2, "▼", RX + Inches(1.7), Inches(3.12), Inches(0.3), Inches(0.2), 10, True, AZ_MAROON, PP_ALIGN.CENTER)

# ── BOTTOM LINE TABLE ──
rbox(s2, RX, Inches(3.3), RW, Inches(0.28), AZ_MAROON)
add_text(s2, "THE BOTTOM LINE", RX + Inches(0.1), Inches(3.3), RW - Inches(0.2), Inches(0.28), 9, True, WHITE, PP_ALIGN.CENTER)

bt = s2.shapes.add_table(8, 3, RX, Inches(3.58), RW, Inches(2.1)).table
bt.columns[0].width = Inches(1.5)
bt.columns[1].width = Inches(1.0)
bt.columns[2].width = Inches(1.3)

set_cell(bt.cell(0, 0), "Task", 7, True, WHITE, PP_ALIGN.LEFT, AZ_MAROON)
set_cell(bt.cell(0, 1), "Manual", 7, True, WHITE, PP_ALIGN.CENTER, AZ_RED)
set_cell(bt.cell(0, 2), "With Agent", 7, True, WHITE, PP_ALIGN.CENTER, AZ_GREEN)

rows = [
    ("Read specs", "3 days", "30 sec"),
    ("Write test cases", "3 days", "30 sec"),
    ("Create test data", "1 day", "15 sec"),
    ("Run tests", "2 days", "2 min"),
    ("Investigate failures", "1 day", "30 sec"),
    ("Write report", "1 day", "15 sec"),
    ("TOTAL", "~2 weeks", "~5 min"),
]

for i, (task, manual, agent) in enumerate(rows):
    is_total = i == len(rows) - 1
    bg = AZ_LIGHT_GRAY if is_total else (WHITE if i % 2 == 0 else AZ_LIGHT_GRAY)
    set_cell(bt.cell(i+1, 0), task, 7 if not is_total else 8, is_total, AZ_DARK, PP_ALIGN.LEFT, bg)
    set_cell(bt.cell(i+1, 1), manual, 7 if not is_total else 9, is_total, AZ_RED, PP_ALIGN.CENTER, AZ_LIGHT_RED if is_total else bg)
    set_cell(bt.cell(i+1, 2), agent, 7 if not is_total else 9, is_total, AZ_GREEN, PP_ALIGN.CENTER, AZ_LIGHT_GREEN if is_total else bg)

# ── PUNCHLINE ──
rbox(s2, RX, Inches(5.78), RW, Inches(0.8), AZ_LIGHT_RED, AZ_RED, 1)
add_text(s2, "WITHOUT AGENT", RX + Inches(0.1), Inches(5.8), RW - Inches(0.2), Inches(0.2), 7, True, AZ_RED, PP_ALIGN.LEFT)
add_multiline_text(s2, [
    ("Bug found during FDA audit", False, AZ_DARK),
    ("Submission delayed by months", True, AZ_RED),
], RX + Inches(0.1), Inches(6.0), RW - Inches(0.2), Inches(0.5), 7)

rbox(s2, RX, Inches(6.62), RW, Inches(0.8), AZ_LIGHT_GREEN, AZ_GREEN, 1)
add_text(s2, "WITH AGENT", RX + Inches(0.1), Inches(6.64), RW - Inches(0.2), Inches(0.2), 7, True, AZ_GREEN, PP_ALIGN.LEFT)
add_multiline_text(s2, [
    ("Bug caught in 5 minutes, fix suggested", False, AZ_DARK),
    ("Zero delay, audit-ready from day one", True, AZ_GREEN),
], RX + Inches(0.1), Inches(6.84), RW - Inches(0.2), Inches(0.5), 7)

# ── LEFT SIDE LABELS: "ALWAYS" vs "IF NEEDED" ──
# Green bar for always-runs steps
add_shape_bg(s2, MSO_SHAPE.RECTANGLE, Inches(2.7), Inches(2.0), Inches(0.06), Inches(1.47), AZ_GREEN)   # Step 1-2
add_shape_bg(s2, MSO_SHAPE.RECTANGLE, Inches(2.7), Inches(4.86), Inches(0.06), Inches(0.50), AZ_GREEN)   # Step 4
add_text(s2, "ALWAYS", Inches(1.8), Inches(2.35), Inches(0.85), Inches(0.2), 6, True, AZ_GREEN, PP_ALIGN.RIGHT)
add_text(s2, "ALWAYS", Inches(1.8), Inches(4.98), Inches(0.85), Inches(0.2), 6, True, AZ_GREEN, PP_ALIGN.RIGHT)

# Gold bar for conditional steps
add_shape_bg(s2, MSO_SHAPE.RECTANGLE, Inches(2.7), Inches(4.18), Inches(0.06), Inches(0.50), AZ_GOLD)   # Step 3
add_shape_bg(s2, MSO_SHAPE.RECTANGLE, Inches(2.7), Inches(6.04), Inches(0.06), Inches(0.98), AZ_GOLD)   # Step 5-6
add_text(s2, "IF\nNEEDED", Inches(1.8), Inches(4.22), Inches(0.85), Inches(0.35), 6, True, AZ_GOLD, PP_ALIGN.RIGHT)
add_text(s2, "IF\nNEEDED", Inches(1.8), Inches(6.2), Inches(0.85), Inches(0.35), 6, True, AZ_GOLD, PP_ALIGN.RIGHT)

# Step 7 label pointing right
add_text(s2, "ALWAYS ►", Inches(7.9), Inches(0.78), Inches(1.2), Inches(0.2), 7, True, AZ_GREEN, PP_ALIGN.RIGHT)

add_text(s2, "2 / 3", Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3), 9, False, AZ_MEDIUM, PP_ALIGN.RIGHT)


# ============================================================================
# SLIDE 3: BOTTOM LINE — Full comparison table + business value (bigger)
# ============================================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])

add_shape_bg(s3, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.6), AZ_MAROON)
add_text(s3, "The Bottom Line", Inches(0.4), Inches(0.08), Inches(8), Inches(0.4), 20, True, WHITE, PP_ALIGN.LEFT)

# Time comparison table (full size)
add_text(s3, "TIME COMPARISON — Per Release Cycle", Inches(0.5), Inches(0.85), Inches(12), Inches(0.35), 14, True, AZ_MAROON, PP_ALIGN.LEFT)

tt = s3.shapes.add_table(8, 3, Inches(0.5), Inches(1.25), Inches(7.5), Inches(3.0)).table
tt.columns[0].width = Inches(3.2)
tt.columns[1].width = Inches(1.9)
tt.columns[2].width = Inches(2.4)

set_cell(tt.cell(0, 0), "Task", 11, True, WHITE, PP_ALIGN.LEFT, AZ_MAROON)
set_cell(tt.cell(0, 1), "Without Agent", 11, True, WHITE, PP_ALIGN.CENTER, AZ_RED)
set_cell(tt.cell(0, 2), "With Agent", 11, True, WHITE, PP_ALIGN.CENTER, AZ_GREEN)

comparisons = [
    ("Read specs & understand requirements", "3 days", "30 seconds"),
    ("Write test cases", "3 days", "30 seconds"),
    ("Create test data", "1 day", "15 seconds"),
    ("Run tests & document results", "2 days", "2 minutes"),
    ("Investigate failures", "1 day", "30 seconds"),
    ("Write compliance report", "1 day", "15 seconds"),
    ("TOTAL", "~2 weeks", "~5 minutes"),
]

for i, (task, without, with_agent) in enumerate(comparisons):
    is_total = i == len(comparisons) - 1
    bg = AZ_LIGHT_GRAY if is_total else (WHITE if i % 2 == 0 else AZ_LIGHT_GRAY)
    set_cell(tt.cell(i+1, 0), task, 10 if not is_total else 12, is_total, AZ_DARK, PP_ALIGN.LEFT, bg)
    set_cell(tt.cell(i+1, 1), without, 10 if not is_total else 14, is_total, AZ_RED, PP_ALIGN.CENTER, AZ_LIGHT_RED if is_total else bg)
    set_cell(tt.cell(i+1, 2), with_agent, 10 if not is_total else 14, is_total, AZ_GREEN, PP_ALIGN.CENTER, AZ_LIGHT_GREEN if is_total else bg)

# RIGHT: Business value cards
add_text(s3, "BUSINESS VALUE", Inches(8.5), Inches(0.85), Inches(4), Inches(0.35), 14, True, AZ_MAROON, PP_ALIGN.LEFT)

rbox(s3, Inches(8.5), Inches(1.25), Inches(4.3), Inches(1.0), AZ_MAROON)
add_text(s3, "REGULATORY RISK PREVENTION", Inches(8.7), Inches(1.3), Inches(3.9), Inches(0.25), 11, True, RGBColor(0xFF, 0xCC, 0xDD), PP_ALIGN.LEFT)
add_text(s3, "Catches compliance gaps that humans miss.\nOne finding during FDA audit can delay\nsubmission by months.", Inches(8.7), Inches(1.6), Inches(3.9), Inches(0.6), 9, False, WHITE, PP_ALIGN.LEFT)

rbox(s3, Inches(8.5), Inches(2.35), Inches(4.3), Inches(0.9), AZ_GREEN)
add_text(s3, "TESTING EFFORT REDUCTION", Inches(8.7), Inches(2.4), Inches(3.9), Inches(0.25), 11, True, WHITE, PP_ALIGN.LEFT)
add_text(s3, "QA engineers focus on strategy & review\ninstead of repetitive manual execution.", Inches(8.7), Inches(2.7), Inches(3.9), Inches(0.5), 9, False, WHITE, PP_ALIGN.LEFT)

rbox(s3, Inches(8.5), Inches(3.35), Inches(4.3), Inches(0.9), AZ_BLUE)
add_text(s3, "FASTER RELEASE CYCLES", Inches(8.7), Inches(3.4), Inches(3.9), Inches(0.25), 11, True, WHITE, PP_ALIGN.LEFT)
add_text(s3, "2 weeks → 5 minutes per cycle.\nMultiply across 9 R&D platforms.", Inches(8.7), Inches(3.7), Inches(3.9), Inches(0.5), 9, False, WHITE, PP_ALIGN.LEFT)

# ── Divider ──
add_shape_bg(s3, MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.55), Inches(13.333), Inches(0.03), AZ_MAROON)

# ── WITHOUT vs WITH (punchline) ──
rbox(s3, Inches(0.5), Inches(4.85), Inches(6.0), Inches(2.3), AZ_LIGHT_RED, AZ_RED, 2)
add_text(s3, "WITHOUT THE AGENT", Inches(0.7), Inches(4.95), Inches(5.6), Inches(0.35), 14, True, AZ_RED, PP_ALIGN.LEFT)
add_multiline_text(s3, [
    ("Missing severity validation goes undetected", False, AZ_DARK),
    ("Found 3 weeks later during manual QA review", False, AZ_DARK),
    ("Or worse — found during FDA audit", False, AZ_DARK),
    ("", False, AZ_DARK),
    ("Result: Submission delayed by months", True, AZ_RED),
], Inches(0.7), Inches(5.35), Inches(5.6), Inches(1.7), 11)

rbox(s3, Inches(6.8), Inches(4.85), Inches(6.0), Inches(2.3), AZ_LIGHT_GREEN, AZ_GREEN, 2)
add_text(s3, "WITH THE AGENT", Inches(7.0), Inches(4.95), Inches(5.6), Inches(0.35), 14, True, AZ_GREEN, PP_ALIGN.LEFT)
add_multiline_text(s3, [
    ("Agent catches missing validation in 5 minutes", False, AZ_DARK),
    ("Fix is suggested immediately", False, AZ_DARK),
    ("Compliance evidence auto-generated for audit", False, AZ_DARK),
    ("", False, AZ_DARK),
    ("Result: Zero delay, audit-ready from day one", True, AZ_GREEN),
], Inches(7.0), Inches(5.35), Inches(5.6), Inches(1.7), 11)

# VS circle
add_shape_bg(s3, MSO_SHAPE.OVAL, Inches(6.25), Inches(5.55), Inches(0.55), Inches(0.55), AZ_MAROON)
add_text(s3, "vs", Inches(6.25), Inches(5.63), Inches(0.55), Inches(0.4), 13, True, WHITE, PP_ALIGN.CENTER)

add_text(s3, "3 / 3", Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3), 9, False, AZ_MEDIUM, PP_ALIGN.RIGHT)

# ============================================================================
# SAVE
# ============================================================================
output_path = "/Users/devesh.b.sharma/Astra Zeneca/R&D IT/Test-Intelligence-Agent-Walkthrough.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
