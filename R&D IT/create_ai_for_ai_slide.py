from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ============================================================================
# COLORS
# ============================================================================
AZ_MAROON = RGBColor(0x83, 0x00, 0x51)
AZ_GREEN = RGBColor(0x2E, 0x7D, 0x32)
AZ_LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
AZ_GOLD = RGBColor(0xC4, 0xA0, 0x00)
AZ_BLUE = RGBColor(0x1E, 0x88, 0xE5)
AZ_LIGHT_BLUE = RGBColor(0xE3, 0xF2, 0xFD)
AZ_ORANGE = RGBColor(0xFF, 0x98, 0x00)
AZ_DARK = RGBColor(0x2D, 0x2D, 0x2D)
AZ_MEDIUM = RGBColor(0x55, 0x55, 0x55)
AZ_LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xFA, 0xF9, 0xF7)
LIGHT_MAROON = RGBColor(0xF8, 0xE0, 0xEB)
GRAY_LINE = RGBColor(0xCC, 0xCC, 0xCC)

# ============================================================================
# HELPERS
# ============================================================================
def add_text(slide, txt, left, top, width, height, size=12, bold=False, color=AZ_DARK, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return box

def add_rich_text(slide, runs_list, left, top, width, height, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for txt, sz, b, c, it in runs_list:
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = b
        r.font.color.rgb = c
        r.font.italic = it
    return box

def add_multiline(slide, lines, left, top, width, height, size=10, color=AZ_DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(1)
        p.space_before = Pt(0)
        if isinstance(line, tuple):
            txt, b, c = line
            p.text = txt
            p.font.size = Pt(size)
            p.font.bold = b
            p.font.color.rgb = c
        else:
            p.text = line
            p.font.size = Pt(size)
            p.font.color.rgb = color
    return box

def rect(slide, left, top, width, height, fill=None, line=None, lw=1):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def rrbox(slide, left, top, width, height, fill, line=None, lw=1):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

# ============================================================================
# PRESENTATION
# ============================================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
sl = prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================================
# TITLE BAR
# ============================================================================
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(1.1), AZ_MAROON)

add_rich_text(sl, [
    ("AI for AI ", 28, True, WHITE, False),
    ("approach for delivery with AI Engineering in Ops", 26, False, WHITE, False),
], Inches(0.4), Inches(0.05), Inches(12.5), Inches(0.55))

add_text(sl, "Leveraging an AI-first & Agentic approach to build the Test Intelligence Agent across End-to-End Lifecycle",
         Inches(0.4), Inches(0.62), Inches(12.5), Inches(0.4), 12, False, RGBColor(0xFF, 0xCC, 0xDD), PP_ALIGN.LEFT, italic=True)

# ============================================================================
# LEFT SIDE — "Today:"
# ============================================================================
LX = Inches(0.25)
LW = Inches(5.9)
LY = Inches(1.25)
LH = Inches(5.95)

# Container
rect(sl, LX, LY, LW, LH, None, GRAY_LINE, 1)

# Header
add_text(sl, "Today:", LX + Inches(0.15), LY + Inches(0.05), Inches(1.5), Inches(0.35), 18, True, AZ_DARK)
add_rich_text(sl, [
    ("Traditional delivery: ", 9, True, AZ_MAROON, True),
    ("precision through Process, progress through People i.e. ", 9, False, AZ_MAROON, True),
    ("'Human centric'", 9, True, AZ_MAROON, True),
    (" delivery approach", 9, False, AZ_MAROON, True),
], LX + Inches(0.15), LY + Inches(0.38), LW - Inches(0.3), Inches(0.35))

# Timeline bar
TBAR_X = LX + Inches(0.15)
TBAR_Y = LY + Inches(0.85)
TBAR_H = Inches(4.55)
rect(sl, TBAR_X, TBAR_Y, Inches(0.07), TBAR_H, AZ_GOLD)
add_text(sl, "16-20 weeks", TBAR_X - Inches(0.55), TBAR_Y + Inches(1.5), Inches(0.55), Inches(1.5), 9, True, AZ_GOLD, PP_ALIGN.CENTER)

# Phase columns
PX = LX + Inches(0.45)    # phase box x
PW = Inches(1.0)           # phase box width
TX = PX + PW + Inches(0.05)  # time x
TW = Inches(0.55)          # time width
DX = TX + TW + Inches(0.15)  # desc x
DW = Inches(3.5)           # desc width
PH = Inches(1.0)           # phase height
PG = Inches(0.12)          # phase gap

today_phases = [
    {
        "name": "Design",
        "color": AZ_GOLD,
        "time": "4-6\nweeks",
        "lines": [
            ("Manual wireframes in PowerPoint", True, AZ_DARK),
            ("or static screenshots emailed around for feedback.", False, AZ_DARK),
            ("Requirements hand-typed into Confluence —", False, AZ_DARK),
            ("40-60 page BRDs that nobody reads end-to-end.", False, AZ_DARK),
            ("Jira stories written manually, 3-hour grooming.", False, AZ_MEDIUM),
        ],
    },
    {
        "name": "Build",
        "color": AZ_GOLD,
        "time": "8-9\nweeks",
        "lines": [
            ("Manual coding", True, AZ_DARK),
            (", line by line. Developers spend 40%", False, AZ_DARK),
            ("of time searching Stack Overflow.", False, AZ_DARK),
            ("PRs sit for 5-10 days waiting for review.", False, AZ_DARK),
            ("New devs take 3-6 months to onboard.", False, AZ_MEDIUM),
        ],
    },
    {
        "name": "Test",
        "color": AZ_GOLD,
        "time": "3-4\nweeks",
        "lines": [
            ("Test cases in Excel spreadsheets", True, AZ_DARK),
            (" — 500-2,000", False, AZ_DARK),
            ("rows manually written and updated.", False, AZ_DARK),
            ("Manual regression: 100+ scenarios by hand.", False, AZ_DARK),
            ("GxP validation docs take months.", False, AZ_MEDIUM),
        ],
    },
    {
        "name": "Deploy",
        "color": AZ_GOLD,
        "time": "1-2\nweeks",
        "lines": [
            ("Manual change requests", True, AZ_DARK),
            (" in ServiceNow.", False, AZ_DARK),
            ("Weekly CAB meetings gate all deployments.", False, AZ_DARK),
            ("20-page deployment runbooks executed by hand.", False, AZ_DARK),
            ("Post-deploy: stare at dashboards for an hour.", False, AZ_MEDIUM),
        ],
    },
]

for i, ph in enumerate(today_phases):
    y = TBAR_Y + Inches(0.05) + (PH + PG) * i

    # Phase box
    rrbox(sl, PX, y, PW, PH, ph["color"])
    add_text(sl, ph["name"], PX, y + Inches(0.3), PW, Inches(0.35), 12, True, WHITE, PP_ALIGN.CENTER)

    # Timeline
    add_text(sl, ph["time"], TX, y + Inches(0.25), TW, Inches(0.5), 9, False, AZ_DARK, PP_ALIGN.CENTER)

    # Arrow
    add_text(sl, "►", TX + TW - Inches(0.1), y + Inches(0.35), Inches(0.2), Inches(0.2), 8, True, AZ_GOLD, PP_ALIGN.CENTER)

    # Description
    add_multiline(sl, ph["lines"], DX, y + Inches(0.05), DW, PH - Inches(0.1), 8)

    # Connector line between phases
    if i < len(today_phases) - 1:
        rect(sl, PX + Inches(0.48), y + PH, Inches(0.02), PG, AZ_GOLD)

# Bottom note
rrbox(sl, LX + Inches(0.15), LY + LH - Inches(0.45), LW - Inches(0.3), Inches(0.35), LIGHT_MAROON, AZ_MAROON, 1)
add_text(sl, "Siloed Change Management & end user adoption achieved over a period",
         LX + Inches(0.25), LY + LH - Inches(0.42), LW - Inches(0.5), Inches(0.3), 8, True, AZ_DARK, PP_ALIGN.LEFT)


# ============================================================================
# CENTER — Arrow
# ============================================================================
arrow = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.2), Inches(3.5), Inches(0.45), Inches(0.7))
arrow.fill.solid()
arrow.fill.fore_color.rgb = AZ_MAROON
arrow.line.fill.background()


# ============================================================================
# RIGHT SIDE — "Tomorrow:"
# ============================================================================
RX = Inches(6.7)
RW = Inches(4.55)
RY = Inches(1.25)
RH = Inches(5.95)

# Container
rect(sl, RX, RY, RW, RH, CREAM, GRAY_LINE, 1)

# Header
add_text(sl, "Tomorrow:", RX + Inches(0.15), RY + Inches(0.05), Inches(2.5), Inches(0.35), 18, True, AZ_DARK)
add_rich_text(sl, [
    ("Agentic AI-powered delivery", 9, True, AZ_MAROON, True),
    (", accelerating and redefining future ways of delivering AI projects", 9, False, AZ_MAROON, True),
], RX + Inches(0.15), RY + Inches(0.38), RW - Inches(0.3), Inches(0.35))

# Timeline bar
TBAR2_X = RX + Inches(0.15)
TBAR2_Y = RY + Inches(0.85)
TBAR2_H = Inches(3.65)
rect(sl, TBAR2_X, TBAR2_Y, Inches(0.07), TBAR2_H, AZ_GREEN)
add_text(sl, "6-10\nweeks", TBAR2_X - Inches(0.55), TBAR2_Y + Inches(1.0), Inches(0.55), Inches(1.5), 9, True, AZ_GREEN, PP_ALIGN.CENTER)

# Tomorrow phases
RPX = RX + Inches(0.45)     # phase box x
RPW = Inches(0.85)          # phase box width
RTX = RPX + RPW + Inches(0.03)  # time x
RTW = Inches(0.5)           # time width
RDX = RTX + RTW + Inches(0.1)   # desc x
RDW = Inches(2.4)           # desc width
RPH = Inches(0.82)          # phase height
RPG = Inches(0.1)           # phase gap

tomorrow_phases = [
    {
        "name": "Design",
        "color": AZ_GREEN,
        "time": "1-2\nweeks",
        "lines": [
            ("Front-end built via Claude / Replit", True, AZ_DARK),
            (" —", False, AZ_DARK),
            ("type a prompt, get a working UI.", False, AZ_DARK),
            ("No PowerPoint mockups needed.", False, AZ_GREEN),
        ],
    },
    {
        "name": "Build",
        "color": AZ_BLUE,
        "time": "3-4\nweeks",
        "lines": [
            ("Claude Code + GitHub Copilot", True, AZ_DARK),
            (" generate", False, AZ_DARK),
            ("code. Agents on Bedrock + Agent Mesh.", False, AZ_DARK),
            ("PRs reviewed in hours, not days.", False, AZ_GREEN),
        ],
    },
    {
        "name": "Test",
        "color": AZ_GREEN,
        "time": "1-2\nweeks",
        "lines": [
            ("AI generates test cases from code", True, AZ_DARK),
            (" — no", False, AZ_DARK),
            ("more Excel. Auto-healing tests.", False, AZ_DARK),
            ("GxP evidence auto-generated.", False, AZ_GREEN),
        ],
    },
    {
        "name": "Deploy",
        "color": AZ_ORANGE,
        "time": "0-1\nweeks",
        "lines": [
            ("Automated via GitHub Actions.", True, AZ_DARK),
            ("", False, AZ_DARK),
            ("Datadog monitoring from day one.", False, AZ_DARK),
            ("Agent validates deployed solution.", False, AZ_GREEN),
        ],
    },
]

for i, ph in enumerate(tomorrow_phases):
    y = TBAR2_Y + Inches(0.05) + (RPH + RPG) * i

    # Phase box
    rrbox(sl, RPX, y, RPW, RPH, ph["color"])
    add_text(sl, ph["name"], RPX, y + Inches(0.22), RPW, Inches(0.3), 10, True, WHITE, PP_ALIGN.CENTER)

    # Timeline
    add_text(sl, ph["time"], RTX, y + Inches(0.2), RTW, Inches(0.4), 8, False, AZ_DARK, PP_ALIGN.CENTER)

    # Arrow
    add_text(sl, "►", RTX + RTW - Inches(0.1), y + Inches(0.28), Inches(0.2), Inches(0.2), 8, True, ph["color"], PP_ALIGN.CENTER)

    # Description
    add_multiline(sl, ph["lines"], RDX, y + Inches(0.05), RDW, RPH - Inches(0.1), 7)

    # Connector line
    if i < len(tomorrow_phases) - 1:
        rect(sl, RPX + Inches(0.4), y + RPH, Inches(0.02), RPG, AZ_GREEN)

# Enhancements row
enh_y = TBAR2_Y + Inches(3.75)
rrbox(sl, RPX, enh_y, RPW, Inches(0.45), AZ_LIGHT_GRAY, AZ_MEDIUM, 1)
add_text(sl, "Enhance", RPX + Inches(0.02), enh_y + Inches(0.07), RPW - Inches(0.04), Inches(0.3), 8, True, AZ_DARK, PP_ALIGN.CENTER)
add_text(sl, "Faster feature builds and change release cycles via Claude Code + continuous agent learning",
         RDX - Inches(0.5), enh_y + Inches(0.05), RDW + Inches(0.5), Inches(0.35), 7, False, AZ_DARK, PP_ALIGN.LEFT)

# Bottom note
rrbox(sl, RX + Inches(0.15), RY + RH - Inches(0.45), RW - Inches(0.3), Inches(0.35), AZ_GREEN)
add_text(sl, "AI-powered delivery with ready organization redesign — future ready from day 1",
         RX + Inches(0.25), RY + RH - Inches(0.42), RW - Inches(0.5), Inches(0.3), 8, True, WHITE, PP_ALIGN.CENTER)


# ============================================================================
# FAR RIGHT — "AI for AI Delivery" tools column
# ============================================================================
FX = Inches(11.35)
FW = Inches(1.85)
FY = Inches(1.25)
FH = Inches(5.95)

rect(sl, FX, FY, FW, FH, None, AZ_MAROON, 1.5)

# Header
add_text(sl, "AI for AI Delivery", FX + Inches(0.05), FY + Inches(0.08), FW - Inches(0.1), Inches(0.3), 11, True, AZ_MAROON, PP_ALIGN.CENTER)

# Subheader
add_text(sl, "Move away from manual tools\nand accelerate through AI",
         FX + Inches(0.05), FY + Inches(0.35), FW - Inches(0.1), Inches(0.4), 7, False, AZ_MEDIUM, PP_ALIGN.CENTER)

# Tool rows — each matched to a Tomorrow phase
# Dashed line separators between
tool_entries = [
    {
        "label": "AI tools",
        "tools": "Replit  |  Claude",
        "desc": "Prompt-to-UI. No more\nPowerPoint wireframes",
        "color": AZ_GREEN,
    },
    {
        "label": "AI tools",
        "tools": "GitHub Copilot\nClaude Code",
        "desc": "46% of code AI-generated.\nPRs in hours, not days",
        "color": AZ_BLUE,
    },
    {
        "label": "AI tools",
        "tools": "TestSigma\nTest Agent (self-test)",
        "desc": "No more Excel test cases.\nAI generates from Jira",
        "color": AZ_GREEN,
    },
    {
        "label": "AI tools",
        "tools": "GitHub Actions\nDatadog",
        "desc": "Auto-deploy, auto-monitor.\nNo manual runbooks",
        "color": AZ_ORANGE,
    },
    {
        "label": "AI tools",
        "tools": "Claude  |  Bedrock",
        "desc": "Continuous AI-powered\nenhancements & learning",
        "color": AZ_MAROON,
    },
]

tool_y_start = FY + Inches(0.8)
tool_h = Inches(0.92)
tool_gap = Inches(0.08)

for i, t in enumerate(tool_entries):
    y = tool_y_start + (tool_h + tool_gap) * i

    # Small color bar on left
    rect(sl, FX + Inches(0.08), y + Inches(0.05), Inches(0.04), tool_h - Inches(0.1), t["color"])

    # "AI tools" label
    add_text(sl, t["label"], FX + Inches(0.18), y, Inches(0.6), Inches(0.18), 6, False, AZ_MEDIUM, PP_ALIGN.LEFT)

    # Tool names (bold)
    add_text(sl, t["tools"], FX + Inches(0.18), y + Inches(0.18), FW - Inches(0.3), Inches(0.35), 8, True, AZ_MAROON, PP_ALIGN.LEFT)

    # Description (italic)
    add_text(sl, t["desc"], FX + Inches(0.18), y + Inches(0.55), FW - Inches(0.3), Inches(0.35), 6, False, AZ_MEDIUM, PP_ALIGN.LEFT, italic=True)

    # Separator line
    if i < len(tool_entries) - 1:
        rect(sl, FX + Inches(0.1), y + tool_h + Inches(0.02), FW - Inches(0.2), Inches(0.01), GRAY_LINE)


# ============================================================================
# SAVE
# ============================================================================
output_path = "/Users/devesh.b.sharma/Astra Zeneca/R&D IT/AI-for-AI-Test-Agent-Delivery.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
