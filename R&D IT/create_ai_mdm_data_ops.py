from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ============================================================================
# AZ BRAND COLORS
# ============================================================================
AZ_MAROON = RGBColor(0x83, 0x00, 0x51)
AZ_DARK_PURPLE = RGBColor(0x4A, 0x14, 0x8C)
AZ_PURPLE = RGBColor(0x5E, 0x35, 0xB1)
AZ_PINK = RGBColor(0xE9, 0x1E, 0x63)
AZ_GREEN = RGBColor(0x2E, 0x7D, 0x32)
AZ_LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
AZ_GOLD = RGBColor(0xC4, 0xA0, 0x00)
AZ_BLUE = RGBColor(0x1E, 0x88, 0xE5)
AZ_LIGHT_BLUE = RGBColor(0xE3, 0xF2, 0xFD)
AZ_ORANGE = RGBColor(0xFF, 0x98, 0x00)
AZ_RED = RGBColor(0xC6, 0x28, 0x28)
AZ_LIGHT_RED = RGBColor(0xFF, 0xEB, 0xEE)
AZ_DARK = RGBColor(0x2D, 0x2D, 0x2D)
AZ_MEDIUM = RGBColor(0x55, 0x55, 0x55)
AZ_LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xFA, 0xF9, 0xF7)
LIGHT_MAROON = RGBColor(0xF8, 0xE0, 0xEB)
GRAY_LINE = RGBColor(0xCC, 0xCC, 0xCC)
AZ_TEAL = RGBColor(0x00, 0x89, 0x7B)
LIGHT_TEAL = RGBColor(0xE0, 0xF2, 0xF1)
LIGHT_PURPLE = RGBColor(0xF3, 0xE5, 0xF5)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xEA, 0xF6)

# ============================================================================
# HELPERS
# ============================================================================
def add_text(slide, txt, left, top, width, height, size=12, bold=False, color=AZ_DARK, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return box

def add_rich_text(slide, runs_list, left, top, width, height, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for txt, sz, b, c, it in runs_list:
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = b
        r.font.color.rgb = c
        r.font.italic = it
    return box

def add_multiline(slide, lines, left, top, width, height, size=10, color=AZ_DARK, spacing=2):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        p.space_before = Pt(0)
        if isinstance(line, tuple):
            txt, b, c = line
            p.text = txt
            p.font.size = Pt(size)
            p.font.bold = b
            p.font.color.rgb = c
        else:
            p.text = line
            p.font.size = Pt(size)
            p.font.color.rgb = color
    return box

def rect(slide, left, top, width, height, fill=None, line=None, lw=1):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def rrbox(slide, left, top, width, height, fill, line=None, lw=1):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def chevron(slide, left, top, width, height, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def arrow_right(slide, left, top, width, height, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def set_cell(cell, text, size, bold, color, align, bg):
    cell.text = text
    cell.text_frame.paragraphs[0].font.size = Pt(size)
    cell.text_frame.paragraphs[0].font.bold = bold
    cell.text_frame.paragraphs[0].font.color.rgb = color
    cell.text_frame.paragraphs[0].alignment = align
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def pentagon(slide, left, top, width, height, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

# ============================================================================
# PRESENTATION SETUP
# ============================================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================================================
# SLIDE 1: TITLE — AGENTIC DATA MANAGEMENT
# ============================================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])

rect(s1, Inches(0), Inches(0), Inches(13.333), Inches(7.5), AZ_MAROON)
rect(s1, Inches(0), Inches(0), Inches(13.333), Inches(0.06), AZ_GOLD)

add_text(s1, "Agentic Data Management",
         Inches(0.8), Inches(1.0), Inches(11), Inches(0.8), 42, True, WHITE, PP_ALIGN.LEFT)

add_text(s1, "Transforming Data Operations for the Future",
         Inches(0.8), Inches(1.8), Inches(11), Inches(0.5), 22, False, RGBColor(0xFF, 0xCC, 0xDD), PP_ALIGN.LEFT, italic=True)

rect(s1, Inches(0.8), Inches(2.5), Inches(4.0), Inches(0.04), AZ_GOLD)

add_multiline(s1, [
    ("Data Strategy Office | AstraZeneca", True, WHITE),
    ("", False, WHITE),
    ("Autonomous AI agents that discover, trace, monitor,", False, RGBColor(0xFF, 0xCC, 0xDD)),
    ("diagnose, suggest, and fix data issues — with humans in the loop", False, RGBColor(0xFF, 0xCC, 0xDD)),
    ("", False, WHITE),
    ("From reactive, manual data management", False, WHITE),
    ("to proactive, self-healing data operations", False, WHITE),
    ("", False, WHITE),
    ("MDM  |  Data Governance  |  Data Quality  |  Data Operations", False, RGBColor(0xFF, 0xCC, 0xDD)),
], Inches(0.8), Inches(2.8), Inches(9), Inches(3.5), 14, WHITE)

# Five core characteristics strip at bottom
chars = [
    ("Autonomy", AZ_DARK_PURPLE),
    ("Goal-Driven", AZ_PURPLE),
    ("Collaboration", AZ_BLUE),
    ("Continuous Learning", AZ_TEAL),
    ("Embedded Guardrails", AZ_GREEN),
]
for i, (label, color) in enumerate(chars):
    x = Inches(0.8 + i * 2.45)
    rrbox(s1, x, Inches(5.8), Inches(2.2), Inches(0.55), color)
    add_text(s1, label, x, Inches(5.83), Inches(2.2), Inches(0.5), 11, True, WHITE, PP_ALIGN.CENTER)

rect(s1, Inches(0), Inches(6.6), Inches(13.333), Inches(0.9), RGBColor(0x5A, 0x00, 0x38))
add_text(s1, "CONFIDENTIAL", Inches(0.8), Inches(6.75), Inches(3), Inches(0.4), 10, True, AZ_GOLD, PP_ALIGN.LEFT)
add_text(s1, "February 2026", Inches(10.5), Inches(6.75), Inches(2.5), Inches(0.4), 10, False, RGBColor(0xFF, 0xCC, 0xDD), PP_ALIGN.RIGHT)

# ============================================================================
# SLIDE 2: MATURITY MODEL — 5 LEVELS
# ============================================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])

rect(s2, Inches(0), Inches(0), Inches(13.333), Inches(0.85), AZ_MAROON)
add_text(s2, "Agentic Data Management Maturity Model",
         Inches(0.4), Inches(0.15), Inches(12), Inches(0.55), 24, True, WHITE, PP_ALIGN.LEFT)

add_text(s2, "Where is AstraZeneca today — and where should we go?",
         Inches(0.4), Inches(0.95), Inches(12), Inches(0.3), 11, False, AZ_MEDIUM, PP_ALIGN.LEFT, italic=True)

# Five maturity levels as ascending blocks
levels = [
    ("Level 1", "MANUAL", "Entirely human-driven\ndata operations", [
        "Manual data entry & cleansing",
        "Spreadsheet-based tracking",
        "No automation of workflows",
        "Tribal knowledge dependency",
    ], RGBColor(0xB0, 0xB0, 0xB0), RGBColor(0xF0, 0xF0, 0xF0)),
    ("Level 2", "TOOL-ASSISTED", "Tools support humans\nbut humans drive", [
        "ETL tools automate movement",
        "Basic profiling & rule checks",
        "Cataloging in structured tools",
        "Human-initiated, tool-aided",
    ], AZ_ORANGE, RGBColor(0xFF, 0xF3, 0xE0)),
    ("Level 3", "AUGMENTED", "AI augments humans\nwith recommendations", [
        "ML-based matching suggestions",
        "AI-powered anomaly detection",
        "Predictive quality scoring",
        "Human approves AI suggestions",
    ], AZ_BLUE, AZ_LIGHT_BLUE),
    ("Level 4", "INTELLIGENT", "AI handles routine decisions\nhumans handle exceptions", [
        "Auto-resolution of low-risk issues",
        "Proactive quality monitoring",
        "AI-driven governance workflows",
        "Human oversight on escalations",
    ], AZ_DARK_PURPLE, LIGHT_PURPLE),
    ("Level 5", "AGENTIC", "Autonomous agents pursue\ngoals with human oversight", [
        "Self-healing data pipelines",
        "Autonomous goal-driven agents",
        "Continuous learning & adaptation",
        "Human-in-the-loop for strategy",
    ], AZ_MAROON, LIGHT_MAROON),
]

# Draw ascending staircase
for i, (lvl, name, desc, bullets, color, bg) in enumerate(levels):
    x = Inches(0.2 + i * 2.6)
    # Height increases with level
    box_height = Inches(1.1 + i * 0.25)
    y = Inches(4.3) - i * Inches(0.25)

    # Main box
    rrbox(s2, x, Inches(1.4), Inches(2.45), Inches(4.8), bg, color, 2)

    # Header
    rrbox(s2, x, Inches(1.4), Inches(2.45), Inches(0.9), color)
    add_text(s2, lvl, x + Inches(0.1), Inches(1.42), Inches(2.25), Inches(0.25), 9, False, WHITE, PP_ALIGN.LEFT)
    add_text(s2, name, x + Inches(0.1), Inches(1.65), Inches(2.25), Inches(0.3), 13, True, WHITE, PP_ALIGN.LEFT)
    add_text(s2, desc, x + Inches(0.1), Inches(1.95), Inches(2.25), Inches(0.35), 7, False, RGBColor(0xFF, 0xCC, 0xDD), PP_ALIGN.LEFT, italic=True)

    # Bullets
    for j, bullet in enumerate(bullets):
        add_text(s2, f"\u2022 {bullet}", x + Inches(0.1), Inches(2.45 + j * Inches(0.35)), Inches(2.25), Inches(0.3), 8, False, AZ_DARK)

    # Connector arrow between levels
    if i < 4:
        add_text(s2, "\u25B6", Inches(2.45 + i * 2.6), Inches(2.2), Inches(0.3), Inches(0.3), 12, True, color, PP_ALIGN.CENTER)

# AZ position marker
rrbox(s2, Inches(2.7), Inches(3.95), Inches(5.3), Inches(0.5), None, AZ_MAROON, 2)
add_rich_text(s2, [
    ("AZ TODAY: ", 10, True, AZ_MAROON, False),
    ("Most domains at Level 2-3. Some pockets of Level 4. ", 9, False, AZ_DARK, False),
    ("NORTH STAR: ", 10, True, AZ_GREEN, False),
    ("Level 5 Agentic across MDM, Governance & Quality", 9, False, AZ_DARK, False),
], Inches(2.85), Inches(3.98), Inches(5.0), Inches(0.4))

# Industry context strip
rect(s2, Inches(0), Inches(6.45), Inches(13.333), Inches(0.03), AZ_MAROON)
add_multiline(s2, [
    ("INDUSTRY CONTEXT:  Most pharma companies are at Level 2-3.  Eli Lilly and Novartis have reached Level 3-4 in select domains.  Level 5 (fully Agentic) is the emerging frontier — first movers gain competitive advantage in data-driven drug development.", False, AZ_MEDIUM),
], Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.7), 9, AZ_MEDIUM)

# ============================================================================
# SLIDE 3: 6-PHASE DATA MANAGEMENT LIFECYCLE
# ============================================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])

rect(s3, Inches(0), Inches(0), Inches(13.333), Inches(0.85), AZ_MAROON)
add_text(s3, "The 6-Phase Agentic Data Management Lifecycle",
         Inches(0.4), Inches(0.15), Inches(12), Inches(0.55), 24, True, WHITE, PP_ALIGN.LEFT)

add_text(s3, "Traditional human-led approach vs. Agentic AI with human-in-the-loop — across each lifecycle phase",
         Inches(0.4), Inches(0.95), Inches(12), Inches(0.3), 10, False, AZ_MEDIUM, PP_ALIGN.LEFT, italic=True)

# Phase chevron flow
phases_lifecycle = [
    ("DISCOVER", AZ_MAROON),
    ("TRACE", AZ_DARK_PURPLE),
    ("MONITOR", AZ_BLUE),
    ("DIAGNOSE", AZ_TEAL),
    ("SUGGEST", AZ_GREEN),
    ("FIX", AZ_GOLD),
]

for i, (label, color) in enumerate(phases_lifecycle):
    x = Inches(0.2 + i * 2.15)
    chevron(s3, x, Inches(1.4), Inches(2.05), Inches(0.6), color)
    add_text(s3, label, x + Inches(0.35), Inches(1.42), Inches(1.4), Inches(0.55), 12, True, WHITE, PP_ALIGN.CENTER)

# Table: Traditional vs Agentic for each phase
# Header row
trad_label_y = Inches(2.2)
agent_label_y = Inches(3.85)

rrbox(s3, Inches(0.2), trad_label_y, Inches(1.7), Inches(1.45), AZ_RED)
add_text(s3, "TRADITIONAL\nHUMAN-LED", Inches(0.25), trad_label_y + Inches(0.3), Inches(1.6), Inches(0.8), 11, True, WHITE, PP_ALIGN.CENTER)

rrbox(s3, Inches(0.2), agent_label_y, Inches(1.7), Inches(1.45), AZ_GREEN)
add_text(s3, "AGENTIC\nAI-DRIVEN", Inches(0.25), agent_label_y + Inches(0.3), Inches(1.6), Inches(0.8), 11, True, WHITE, PP_ALIGN.CENTER)

traditional = [
    ("Manual metadata\ncollection via\ninterviews and\ndocumentation", AZ_LIGHT_RED),
    ("Lineage documented\nmanually in\nspreadsheets\nand Visio diagrams", AZ_LIGHT_RED),
    ("Periodic quality\naudits (quarterly)\nwith manual\nscoring", AZ_LIGHT_RED),
    ("Reactive root-cause\nanalysis after\nbusiness impact\nalready felt", AZ_LIGHT_RED),
    ("Stewards manually\nresearch and\npropose fixes\nvia tickets", AZ_LIGHT_RED),
    ("Manual remediation\ncycles taking\ndays to weeks\nper issue", AZ_LIGHT_RED),
]

agentic = [
    ("AI agents auto-discover\nmetadata, schemas,\nrelationships across\nall platforms", AZ_LIGHT_GREEN),
    ("Automated lineage\ntracing across\npipelines, APIs,\nand transformations", AZ_LIGHT_GREEN),
    ("24/7 real-time\nquality monitoring\nwith predictive\nanomaly detection", AZ_LIGHT_GREEN),
    ("AI agents diagnose\nroot causes instantly\nusing pattern analysis\nand knowledge base", AZ_LIGHT_GREEN),
    ("Context-aware\nremediation suggestions\nranked by confidence\nand business impact", AZ_LIGHT_GREEN),
    ("Auto-fix low-risk\nissues, escalate\nhigh-risk with full\ncontext to stewards", AZ_LIGHT_GREEN),
]

for i, ((trad_text, trad_bg), (agent_text, agent_bg)) in enumerate(zip(traditional, agentic)):
    x = Inches(2.05 + i * 1.9)
    # Traditional row
    rrbox(s3, x, trad_label_y, Inches(1.75), Inches(1.45), trad_bg, AZ_RED, 1)
    add_text(s3, trad_text, x + Inches(0.08), trad_label_y + Inches(0.1), Inches(1.6), Inches(1.25), 8, False, AZ_DARK)
    # Agentic row
    rrbox(s3, x, agent_label_y, Inches(1.75), Inches(1.45), agent_bg, AZ_GREEN, 1)
    add_text(s3, agent_text, x + Inches(0.08), agent_label_y + Inches(0.1), Inches(1.6), Inches(1.25), 8, False, AZ_DARK)

# Down arrow between rows
add_text(s3, "\u25BC  Agentic AI Transformation  \u25BC", Inches(4.5), Inches(3.67), Inches(4.5), Inches(0.2), 9, True, AZ_MAROON, PP_ALIGN.CENTER)

# Bottom insight
rect(s3, Inches(0), Inches(5.5), Inches(13.333), Inches(0.03), AZ_MAROON)
add_text(s3, "How This Maps to AstraZeneca", Inches(0.4), Inches(5.65), Inches(6), Inches(0.3), 13, True, AZ_MAROON)

az_mapping = [
    ("DISCOVER", "Auto-discover metadata across Informatica MDM, SAP MDG, Collibra, Databricks, Veeva", AZ_MAROON),
    ("TRACE", "Automated lineage from SAP S/4HANA through MDM to downstream analytics and BIKG", AZ_DARK_PURPLE),
    ("MONITOR", "Real-time quality monitoring replacing quarterly Excel scorecards across all domains", AZ_BLUE),
    ("DIAGNOSE", "AI-powered root cause analysis on HCP matching failures, material master conflicts", AZ_TEAL),
    ("SUGGEST", "Context-aware remediation for data stewards — ranked by confidence and risk", AZ_GREEN),
    ("FIX", "Autonomous low-risk fixes (dedup, format, enrich); escalation for high-risk changes", AZ_GOLD),
]

for i, (phase, desc, color) in enumerate(az_mapping):
    col = i % 3
    row = i // 3
    x = Inches(0.3 + col * 4.35)
    y = Inches(6.05 + row * Inches(0.6))
    add_rich_text(s3, [
        (f"{phase}: ", 9, True, color, False),
        (desc, 8, False, AZ_DARK, False),
    ], x, y, Inches(4.2), Inches(0.55))

# ============================================================================
# SLIDE 4: AGENTIC PLATFORM ARCHITECTURE (CAPCO LAYERED MODEL)
# ============================================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])

rect(s4, Inches(0), Inches(0), Inches(13.333), Inches(0.85), AZ_MAROON)
add_text(s4, "Agentic Data Management Platform Architecture",
         Inches(0.4), Inches(0.15), Inches(12), Inches(0.55), 24, True, WHITE, PP_ALIGN.LEFT)

add_text(s4, "Layered architecture — agentic capabilities on top of existing AZ data investments (no rip-and-replace)",
         Inches(0.4), Inches(0.95), Inches(12), Inches(0.3), 10, False, AZ_MEDIUM, PP_ALIGN.LEFT, italic=True)

# ── LAYER 1: HUMAN INTERACTION CONSOLE ──
layer1_y = Inches(1.4)
rrbox(s4, Inches(0.25), layer1_y, Inches(2.3), Inches(0.9), AZ_MAROON)
add_text(s4, "HUMAN\nINTERACTION", Inches(0.3), layer1_y + Inches(0.15), Inches(2.2), Inches(0.7), 11, True, WHITE, PP_ALIGN.CENTER)

console_items = [
    ("Data Steward\nConsole", "Review, approve,\noverride agent decisions"),
    ("Governance\nDashboard", "Real-time quality\nscores & compliance"),
    ("Natural Language\nInterface", "Ask questions in plain\nEnglish, get answers"),
    ("Audit &\nExplainability", "Full trail of every\nagent decision & action"),
]
for j, (item_name, item_desc) in enumerate(console_items):
    x = Inches(2.8 + j * 2.65)
    rrbox(s4, x, layer1_y, Inches(2.45), Inches(0.9), LIGHT_MAROON, AZ_MAROON, 1)
    add_text(s4, item_name, x + Inches(0.1), layer1_y + Inches(0.02), Inches(2.25), Inches(0.4), 9, True, AZ_MAROON)
    add_text(s4, item_desc, x + Inches(0.1), layer1_y + Inches(0.42), Inches(2.25), Inches(0.45), 8, False, AZ_MEDIUM)

# Down arrows
add_text(s4, "\u2193  \u2193  \u2193  \u2193  \u2193", Inches(4.5), Inches(2.3), Inches(4.5), Inches(0.25), 12, True, AZ_MAROON, PP_ALIGN.CENTER)

# ── LAYER 2: POLICY, ORCHESTRATION & MONITORING AGENTS ──
layer2_y = Inches(2.6)
rrbox(s4, Inches(0.25), layer2_y, Inches(2.3), Inches(0.9), AZ_DARK_PURPLE)
add_text(s4, "POLICY &\nORCHESTRATION", Inches(0.3), layer2_y + Inches(0.15), Inches(2.2), Inches(0.7), 11, True, WHITE, PP_ALIGN.CENTER)

policy_items = [
    ("Policy Agent", "Enforces data governance\npolicies autonomously"),
    ("Dynamic Workflow\nAgent", "Routes tasks based on\nrisk, urgency, context"),
    ("Observability\nAgent", "Monitors all agent\nactions & performance"),
    ("Orchestration\nAgent", "Coordinates multi-agent\nworkflows end-to-end"),
]
for j, (item_name, item_desc) in enumerate(policy_items):
    x = Inches(2.8 + j * 2.65)
    rrbox(s4, x, layer2_y, Inches(2.45), Inches(0.9), LIGHT_PURPLE, AZ_DARK_PURPLE, 1)
    add_text(s4, item_name, x + Inches(0.1), layer2_y + Inches(0.02), Inches(2.25), Inches(0.4), 9, True, AZ_DARK_PURPLE)
    add_text(s4, item_desc, x + Inches(0.1), layer2_y + Inches(0.42), Inches(2.25), Inches(0.45), 8, False, AZ_MEDIUM)

# Down arrows
add_text(s4, "\u2193  \u2193  \u2193  \u2193  \u2193", Inches(4.5), Inches(3.5), Inches(4.5), Inches(0.25), 12, True, AZ_DARK_PURPLE, PP_ALIGN.CENTER)

# ── LAYER 3: DATA MANAGEMENT AGENTS ──
layer3_y = Inches(3.8)
rrbox(s4, Inches(0.25), layer3_y, Inches(2.3), Inches(0.9), AZ_BLUE)
add_text(s4, "DATA\nMANAGEMENT\nAGENTS", Inches(0.3), layer3_y + Inches(0.05), Inches(2.2), Inches(0.8), 11, True, WHITE, PP_ALIGN.CENTER)

dm_agents = [
    ("Metadata\nDiscovery Agent", "Auto-discovers schemas,\nrelationships, owners"),
    ("Lineage\nAgent", "Traces data flow across\nentire platform"),
    ("DQ Monitoring\nAgent", "Real-time quality checks,\nanomaly detection"),
    ("Root Cause\nAnalysis Agent", "Diagnoses data issues\nusing pattern analysis"),
]
for j, (item_name, item_desc) in enumerate(dm_agents):
    x = Inches(2.8 + j * 2.65)
    rrbox(s4, x, layer3_y, Inches(2.45), Inches(0.9), AZ_LIGHT_BLUE, AZ_BLUE, 1)
    add_text(s4, item_name, x + Inches(0.1), layer3_y + Inches(0.02), Inches(2.25), Inches(0.4), 9, True, AZ_BLUE)
    add_text(s4, item_desc, x + Inches(0.1), layer3_y + Inches(0.42), Inches(2.25), Inches(0.45), 8, False, AZ_MEDIUM)

# Additional agents row
dm_agents_2 = [
    ("Resolution\nAgent", "Applies automated\nfixes for known issues"),
    ("Escalation\nAgent", "Routes complex issues\nto human stewards"),
]
for j, (item_name, item_desc) in enumerate(dm_agents_2):
    x = Inches(2.8 + j * 2.65)
    rrbox(s4, x, layer3_y + Inches(0.95), Inches(2.45), Inches(0.65), AZ_LIGHT_BLUE, AZ_BLUE, 1)
    add_text(s4, item_name, x + Inches(0.1), layer3_y + Inches(0.97), Inches(2.25), Inches(0.35), 9, True, AZ_BLUE)
    add_text(s4, item_desc, x + Inches(0.1), layer3_y + Inches(1.3), Inches(2.25), Inches(0.25), 8, False, AZ_MEDIUM)

# Down arrows
add_text(s4, "\u2193  \u2193  \u2193  \u2193  \u2193", Inches(4.5), Inches(5.4), Inches(4.5), Inches(0.25), 12, True, AZ_BLUE, PP_ALIGN.CENTER)

# ── LAYER 4: KNOWLEDGE BASE ──
layer4_y = Inches(5.7)
rrbox(s4, Inches(0.25), layer4_y, Inches(2.3), Inches(0.65), AZ_TEAL)
add_text(s4, "KNOWLEDGE\nBASE", Inches(0.3), layer4_y + Inches(0.08), Inches(2.2), Inches(0.5), 11, True, WHITE, PP_ALIGN.CENTER)

kb_items = [
    ("Business Rules\n& Policies", "Data standards, validation\nrules, compliance policies"),
    ("Historical Patterns\n& Decisions", "Past resolutions, steward\ndecisions, outcomes"),
    ("Domain Knowledge\n& Ontologies", "Pharma domain models,\ndata dictionaries"),
]
for j, (item_name, item_desc) in enumerate(kb_items):
    x = Inches(2.8 + j * 2.65)
    rrbox(s4, x, layer4_y, Inches(2.45), Inches(0.65), LIGHT_TEAL, AZ_TEAL, 1)
    add_text(s4, item_name, x + Inches(0.1), layer4_y + Inches(0.02), Inches(2.25), Inches(0.3), 9, True, AZ_TEAL)
    add_text(s4, item_desc, x + Inches(0.1), layer4_y + Inches(0.3), Inches(2.25), Inches(0.3), 8, False, AZ_MEDIUM)

# ── LAYER 5: DATA PLATFORMS (AZ-specific) ──
layer5_y = Inches(6.55)
rrbox(s4, Inches(0.25), layer5_y, Inches(2.3), Inches(0.65), AZ_GREEN)
add_text(s4, "AZ DATA\nPLATFORMS", Inches(0.3), layer5_y + Inches(0.08), Inches(2.2), Inches(0.5), 11, True, WHITE, PP_ALIGN.CENTER)

platform_items = [
    ("Informatica MDM\n+ SAP MDG", "Master data hubs"),
    ("Collibra", "Data catalog & governance"),
    ("AWS / Databricks", "Cloud & analytics"),
    ("SAP S/4HANA\n+ Veeva", "ERP & CRM"),
]
for j, (item_name, item_desc) in enumerate(platform_items):
    x = Inches(2.8 + j * 2.65)
    rrbox(s4, x, layer5_y, Inches(2.45), Inches(0.65), AZ_LIGHT_GREEN, AZ_GREEN, 1)
    add_text(s4, item_name, x + Inches(0.1), layer5_y + Inches(0.02), Inches(2.25), Inches(0.3), 9, True, AZ_GREEN)
    add_text(s4, item_desc, x + Inches(0.1), layer5_y + Inches(0.32), Inches(2.25), Inches(0.25), 8, False, AZ_MEDIUM)

# Right-side label showing Capco source
add_text(s4, "Based on Capco/Wipro\nAgentic Data Management\nFramework", Inches(10.8), Inches(5.7), Inches(2.3), Inches(0.6), 8, False, AZ_MEDIUM, PP_ALIGN.RIGHT, italic=True)

# ============================================================================
# SLIDE 5: TODAY & TOMORROW — OPERATING MODEL SHIFT
# ============================================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])

rect(s5, Inches(0), Inches(0), Inches(13.333), Inches(0.85), AZ_MAROON)
add_text(s5, "The Operating Model Shift — Today & Tomorrow with Agentic AI",
         Inches(0.4), Inches(0.15), Inches(12), Inches(0.55), 24, True, WHITE, PP_ALIGN.LEFT)

# ── LEFT: TODAY ──
rect(s5, Inches(0.25), Inches(1.1), Inches(6.2), Inches(6.1), None, GRAY_LINE, 1)
add_text(s5, "TODAY", Inches(0.4), Inches(1.15), Inches(2), Inches(0.35), 18, True, AZ_RED, PP_ALIGN.LEFT)
add_text(s5, "Manual, Reactive, Resource-Intensive", Inches(1.6), Inches(1.2), Inches(4.5), Inches(0.3), 10, True, AZ_MEDIUM, PP_ALIGN.LEFT, italic=True)

today_items = [
    ("Master Data Management", AZ_RED, [
        "Manual golden record creation across 6+ domains",
        "HCP/Customer matching: high false positives",
        "6 data stewards for material master alone",
        "SAP MDG + Informatica MDM: siloed processes",
    ]),
    ("Data Governance", AZ_ORANGE, [
        "Policy documentation in static SharePoint docs",
        "Manual data classification and cataloging",
        "Lineage traced through interviews, not automation",
        "Compliance checks: quarterly, not real-time",
    ]),
    ("Data Quality", AZ_GOLD, [
        "Reactive issue detection post-incident",
        "Excel-based quality scorecards, lag 2-4 weeks",
        "Manual profiling and remediation cycles",
        "No predictive quality monitoring capability",
    ]),
    ("Data Operations", AZ_MEDIUM, [
        "High OSP dependency for routine tasks",
        "FTE-based staffing model, not outcome-driven",
        "Ticket-based data creation/maintenance workflows",
        "Manual data migration and cleansing efforts",
    ]),
]

for i, (title, color, bullets) in enumerate(today_items):
    y = Inches(1.7) + i * Inches(1.12)
    rrbox(s5, Inches(0.4), y, Inches(5.9), Inches(1.0), None, color, 1)
    add_text(s5, title, Inches(0.55), y + Inches(0.02), Inches(2.5), Inches(0.25), 10, True, color, PP_ALIGN.LEFT)
    for j, bullet in enumerate(bullets):
        add_text(s5, f"\u2022 {bullet}", Inches(0.55), y + Inches(0.25 + j * Inches(0.17)), Inches(5.6), Inches(0.18), 8, False, AZ_DARK, PP_ALIGN.LEFT)

# ── RIGHT: TOMORROW ──
rect(s5, Inches(6.8), Inches(1.1), Inches(6.25), Inches(6.1), None, AZ_GREEN, 2)
add_text(s5, "TOMORROW", Inches(6.95), Inches(1.15), Inches(3), Inches(0.35), 18, True, AZ_GREEN, PP_ALIGN.LEFT)
add_text(s5, "Agentic, Autonomous, Self-Healing", Inches(8.95), Inches(1.2), Inches(4), Inches(0.3), 10, True, AZ_GREEN, PP_ALIGN.LEFT, italic=True)

tomorrow_items = [
    ("Agentic MDM", AZ_GREEN, [
        "Discovery agents auto-find metadata across all systems",
        "Resolution agents auto-fix low-risk matching issues",
        "Stewards shift from data entry to strategy (6 to 1)",
        "Orchestration agent coordinates cross-domain harmonization",
    ]),
    ("Agentic Governance", AZ_TEAL, [
        "Policy agent enforces compliance continuously, not quarterly",
        "Lineage agent traces data flow across entire platform",
        "Automated classification — PII/PHI tagged in real-time",
        "Observability agent monitors all governance agent actions",
    ]),
    ("Agentic Data Quality", AZ_BLUE, [
        "DQ monitoring agent detects anomalies 24/7 proactively",
        "Root cause analysis agent diagnoses issues instantly",
        "Self-healing: auto-fix, escalate when confidence is low",
        "Continuous learning from steward corrections over time",
    ]),
    ("Agentic Operations", AZ_DARK_PURPLE, [
        "70% reduction in manual data stewardship effort",
        "Dynamic workflow agent routes tasks by risk and urgency",
        "Outcome-based OSP model, not FTE-based",
        "Human-in-the-loop for strategy; agents for execution",
    ]),
]

for i, (title, color, bullets) in enumerate(tomorrow_items):
    y = Inches(1.7) + i * Inches(1.12)
    rrbox(s5, Inches(6.95), y, Inches(5.95), Inches(1.0), None, color, 1)
    add_text(s5, title, Inches(7.1), y + Inches(0.02), Inches(2.5), Inches(0.25), 10, True, color, PP_ALIGN.LEFT)
    for j, bullet in enumerate(bullets):
        add_text(s5, f"\u2022 {bullet}", Inches(7.1), y + Inches(0.25 + j * Inches(0.17)), Inches(5.6), Inches(0.18), 8, False, AZ_DARK, PP_ALIGN.LEFT)

# VS circle
rrbox(s5, Inches(6.15), Inches(3.5), Inches(1.0), Inches(0.75), AZ_MAROON)
add_text(s5, "Agentic\nAI", Inches(6.15), Inches(3.52), Inches(1.0), Inches(0.7), 12, True, WHITE, PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 6: CASE STUDIES — ELI LILLY & NOVARTIS
# ============================================================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])

rect(s6, Inches(0), Inches(0), Inches(13.333), Inches(0.85), AZ_MAROON)
add_text(s6, "What Pharma Leaders Actually Did — Eli Lilly & Novartis",
         Inches(0.4), Inches(0.15), Inches(12), Inches(0.55), 24, True, WHITE, PP_ALIGN.LEFT)

# ── ELI LILLY ──
rrbox(s6, Inches(0.25), Inches(1.05), Inches(6.3), Inches(6.0), None, AZ_DARK_PURPLE, 2)
rrbox(s6, Inches(0.25), Inches(1.05), Inches(6.3), Inches(0.55), AZ_DARK_PURPLE)
add_text(s6, "Eli Lilly", Inches(0.5), Inches(1.1), Inches(5.8), Inches(0.45), 18, True, WHITE, PP_ALIGN.LEFT)

lilly_sections = [
    ("What They Built", [
        ("Global Manufacturing Data Fabric (GMDF) on Databricks", True, AZ_DARK),
        ("Serverless architecture, PySpark automation", False, AZ_DARK),
        ("Unity Catalog for governance, Auto Loader for streaming", False, AZ_DARK),
        ("Centralized analytics platform scaling AI enterprise-wide", False, AZ_MEDIUM),
    ]),
    ("What Changed in MDM", [
        ("Material Master: 100 fields put under global governance", True, AZ_DARK),
        ("Syniti platform automated data governance workflows", False, AZ_DARK),
        ("Global data stewards reduced from 6 to 1", True, AZ_GREEN),
        ("Changes auto-distributed to local stewards via workflow", False, AZ_MEDIUM),
    ]),
    ("Mapping to Agentic Lifecycle", [
        ("DISCOVER + TRACE: Auto-discovery via Unity Catalog & lineage", True, AZ_DARK_PURPLE),
        ("MONITOR: Automated quality checks on manufacturing data", False, AZ_DARK),
        ("SUGGEST + FIX: Workflows auto-route and auto-apply changes", False, AZ_DARK),
        ("10+ AI use cases running autonomously across the enterprise", False, AZ_MEDIUM),
    ]),
]

y_offset = Inches(1.75)
for section_title, items in lilly_sections:
    add_text(s6, section_title, Inches(0.5), y_offset, Inches(5.8), Inches(0.25), 11, True, AZ_DARK_PURPLE)
    y_offset += Inches(0.28)
    for txt, bold, color in items:
        add_text(s6, f"\u2022 {txt}", Inches(0.65), y_offset, Inches(5.7), Inches(0.2), 9, bold, color)
        y_offset += Inches(0.22)
    y_offset += Inches(0.12)

rrbox(s6, Inches(0.5), Inches(6.15), Inches(5.8), Inches(0.7), LIGHT_MAROON, AZ_DARK_PURPLE, 1)
add_rich_text(s6, [
    ("KEY RESULT: ", 10, True, AZ_DARK_PURPLE, False),
    ("6 stewards to 1. 100 material master fields under automated governance. Autonomous workflows replaced manual ticket-based processes.", 9, False, AZ_DARK, False),
], Inches(0.65), Inches(6.25), Inches(5.5), Inches(0.5))

# ── NOVARTIS ──
rrbox(s6, Inches(6.8), Inches(1.05), Inches(6.3), Inches(6.0), None, AZ_TEAL, 2)
rrbox(s6, Inches(6.8), Inches(1.05), Inches(6.3), Inches(0.55), AZ_TEAL)
add_text(s6, "Novartis", Inches(7.05), Inches(1.1), Inches(5.8), Inches(0.45), 18, True, WHITE, PP_ALIGN.LEFT)

novartis_sections = [
    ("What They Built", [
        ("Multi-cloud data platform (AWS primary, Azure secondary)", True, AZ_DARK),
        ("Accenture-built data operating model + governance", False, AZ_DARK),
        ("Centralized data catalogue, Dataiku for AI governance", False, AZ_DARK),
        ("Centralized multi-tenant data hub across countries", False, AZ_MEDIUM),
    ]),
    ("What Changed in Data Management", [
        ("Enterprise Data Management + Quality + Compliance unified", True, AZ_DARK),
        ("New use case setup: went from 2 weeks to 1 day", True, AZ_GREEN),
        ("Data & Digital Leadership made 1 of 5 strategic priorities", False, AZ_DARK),
        ("Partnerships with Microsoft, ConcertAI, Amazon for AI", False, AZ_MEDIUM),
    ]),
    ("Mapping to Agentic Lifecycle", [
        ("DISCOVER + TRACE: Centralized catalogue, automated lineage", True, AZ_TEAL),
        ("MONITOR + DIAGNOSE: 90% faster insights (21 days to 2)", False, AZ_DARK),
        ("SUGGEST + FIX: Modular AI components reused across use cases", False, AZ_DARK),
        ("AI governance ensures agents themselves are compliant", False, AZ_MEDIUM),
    ]),
]

y_offset = Inches(1.75)
for section_title, items in novartis_sections:
    add_text(s6, section_title, Inches(7.05), y_offset, Inches(5.8), Inches(0.25), 11, True, AZ_TEAL)
    y_offset += Inches(0.28)
    for txt, bold, color in items:
        add_text(s6, f"\u2022 {txt}", Inches(7.2), y_offset, Inches(5.7), Inches(0.2), 9, bold, color)
        y_offset += Inches(0.22)
    y_offset += Inches(0.12)

rrbox(s6, Inches(7.05), Inches(6.15), Inches(5.8), Inches(0.7), LIGHT_TEAL, AZ_TEAL, 1)
add_rich_text(s6, [
    ("KEY RESULT: ", 10, True, AZ_TEAL, False),
    ("90% faster insights. 2 weeks to 1 day for new use cases. Autonomous AI pipelines across the data platform — not manual, not prompted.", 9, False, AZ_DARK, False),
], Inches(7.2), Inches(6.25), Inches(5.5), Inches(0.5))

# ============================================================================
# SLIDE 7: BENEFITS FRAMEWORK — 6 PILLARS
# ============================================================================
s7 = prs.slides.add_slide(prs.slide_layouts[6])

rect(s7, Inches(0), Inches(0), Inches(13.333), Inches(0.85), AZ_MAROON)
add_text(s7, "Benefits of Agentic Data Management — Why Now",
         Inches(0.4), Inches(0.15), Inches(12), Inches(0.55), 24, True, WHITE, PP_ALIGN.LEFT)

add_text(s7, "Six pillars of value from adopting an agentic approach to data management at AstraZeneca",
         Inches(0.4), Inches(0.95), Inches(12), Inches(0.3), 10, False, AZ_MEDIUM, PP_ALIGN.LEFT, italic=True)

benefits = [
    ("User\nEmpowerment", AZ_MAROON, [
        "Self-service via natural language",
        "Stewards become strategists",
        "Reduced dependency on IT/OSP",
        "Faster time-to-answer",
    ]),
    ("Smart Metadata\nManagement", AZ_DARK_PURPLE, [
        "Auto-discovery across platforms",
        "Relationships mapped automatically",
        "Always-current data catalogue",
        "Reduced manual documentation",
    ]),
    ("Proactive Data\nQuality", AZ_BLUE, [
        "24/7 anomaly detection",
        "Predictive quality scoring",
        "Self-healing data pipelines",
        "Issues caught before impact",
    ]),
    ("Process\nAutomation", AZ_TEAL, [
        "70% stewardship automation",
        "Autonomous workflow routing",
        "Low-risk auto-resolution",
        "Outcome-based OSP model",
    ]),
    ("Future-Ready\nFoundation", AZ_GREEN, [
        "Scalable agent architecture",
        "Continuous learning system",
        "Ready for new AI capabilities",
        "Modular, extensible platform",
    ]),
    ("CDO\nEffectiveness", AZ_GOLD, [
        "Real-time governance visibility",
        "Data-driven decision making",
        "Reduced time-to-insight by 90%",
        "Measurable ROI on data ops",
    ]),
]

for i, (title, color, bullets) in enumerate(benefits):
    col = i % 3
    row = i // 3
    x = Inches(0.3 + col * 4.35)
    y = Inches(1.45 + row * 2.85)

    rrbox(s7, x, y, Inches(4.1), Inches(2.65), None, color, 2)
    rrbox(s7, x, y, Inches(4.1), Inches(0.75), color)
    add_text(s7, title, x + Inches(0.15), y + Inches(0.05), Inches(3.8), Inches(0.65), 13, True, WHITE, PP_ALIGN.LEFT)

    for j, bullet in enumerate(bullets):
        add_text(s7, f"\u2022 {bullet}", x + Inches(0.2), y + Inches(0.85 + j * Inches(0.4)), Inches(3.7), Inches(0.35), 10, False, AZ_DARK)

# Impact metrics bar at bottom
rect(s7, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5), AZ_LIGHT_GRAY)
metrics = [
    ("6 \u2192 1", "Steward Reduction\n(Eli Lilly)", AZ_MAROON),
    ("90%", "Faster Insights\n(Novartis)", AZ_GREEN),
    ("70%", "Stewardship Automation\n(McKinsey)", AZ_BLUE),
    ("40%", "Faster Time-to-Value\n(Gartner)", AZ_DARK_PURPLE),
    ("$3.70", "Return per $1\n(Industry Avg)", AZ_GOLD),
]
for i, (metric, label, color) in enumerate(metrics):
    x = Inches(0.3 + i * 2.6)
    add_rich_text(s7, [
        (metric, 14, True, color, False),
        (f"  {label}", 8, False, AZ_MEDIUM, False),
    ], x, Inches(7.02), Inches(2.4), Inches(0.45))

# ============================================================================
# SLIDE 8: IMPLEMENTATION ROADMAP — MAPPING A PATH FORWARD
# ============================================================================
s8 = prs.slides.add_slide(prs.slide_layouts[6])

rect(s8, Inches(0), Inches(0), Inches(13.333), Inches(0.85), AZ_MAROON)
add_text(s8, "Mapping a Path Forward — Phased Agentic AI Adoption",
         Inches(0.4), Inches(0.15), Inches(12), Inches(0.55), 24, True, WHITE, PP_ALIGN.LEFT)

phases = [
    ("PHASE 1", "Q1-Q2 2026", "Assess & Prove", AZ_MAROON, [
        "Maturity assessment across\nall data domains (Level 1-5)",
        "POC: agentic entity resolution\non HCP/Customer domain",
        "POC: DQ monitoring agent\non top 3 quality dimensions",
        "Pilot governance copilot\n(NL interface on Collibra)",
        "Baseline current metrics:\neffort, quality, cost, speed",
    ]),
    ("PHASE 2", "Q3-Q4 2026", "Build & Scale", AZ_DARK_PURPLE, [
        "Deploy policy + orchestration\nagent layer in production",
        "Agentic matching & dedup\nacross all MDM domains",
        "Autonomous stewardship for\nlow-risk data changes",
        "Real-time quality dashboards\nreplace Excel scorecards",
        "OSP model pilot: shift to\noutcome-based contracts",
    ]),
    ("PHASE 3", "H1 2027", "Optimize & Learn", AZ_BLUE, [
        "Full 6-phase lifecycle running\n(Discover to Fix)",
        "Self-healing data pipelines\nin production across domains",
        "Knowledge base learns from\nsteward decisions over time",
        "Continuous compliance\nmonitoring (real-time)",
        "Expand to Material +\nSupplier + Site domains",
    ]),
    ("PHASE 4", "H2 2027+", "Autonomous Ops", AZ_TEAL, [
        "Level 5 maturity across\ncore data domains",
        "Autonomous data operations\nwith human oversight",
        "Cross-domain intelligent\ndata harmonization",
        "Industry-leading data\nquality and governance",
        "Measurable ROI: $3.70+\nreturn per $1 invested",
    ]),
]

for i, (phase, timeline, label, color, items) in enumerate(phases):
    x = Inches(0.25 + i * 3.25)
    rrbox(s8, x, Inches(1.1), Inches(3.0), Inches(0.85), color)
    add_text(s8, phase, x + Inches(0.15), Inches(1.12), Inches(1.2), Inches(0.3), 14, True, WHITE)
    add_text(s8, timeline, x + Inches(1.4), Inches(1.12), Inches(1.4), Inches(0.3), 10, False, WHITE, PP_ALIGN.RIGHT)
    add_text(s8, label, x + Inches(0.15), Inches(1.45), Inches(2.7), Inches(0.3), 12, True, RGBColor(0xFF, 0xCC, 0xDD))
    rrbox(s8, x, Inches(2.05), Inches(3.0), Inches(3.6), None, color, 1)
    for j, item in enumerate(items):
        add_text(s8, f"\u2022 {item}", x + Inches(0.1), Inches(2.15 + j * Inches(0.65)), Inches(2.8), Inches(0.6), 9, False, AZ_DARK)
    if i < 3:
        add_text(s8, "\u25B6", Inches(3.05 + i * 3.25), Inches(1.35), Inches(0.3), Inches(0.3), 14, True, AZ_MAROON, PP_ALIGN.CENTER)

# Maturity progression bar
rect(s8, Inches(0), Inches(5.85), Inches(13.333), Inches(0.03), AZ_MAROON)
add_text(s8, "MATURITY PROGRESSION", Inches(0.4), Inches(5.95), Inches(4), Inches(0.3), 12, True, AZ_MAROON)

maturity_bar = [
    ("Level 2-3", "Current State", GRAY_LINE, AZ_MEDIUM),
    ("Level 3", "After Phase 1", AZ_ORANGE, WHITE),
    ("Level 3-4", "After Phase 2", AZ_BLUE, WHITE),
    ("Level 4-5", "After Phase 3", AZ_DARK_PURPLE, WHITE),
    ("Level 5", "Agentic (Target)", AZ_MAROON, WHITE),
]

for i, (level, label, color, text_color) in enumerate(maturity_bar):
    x = Inches(0.25 + i * 2.6)
    rrbox(s8, x, Inches(6.3), Inches(2.4), Inches(0.5), color)
    add_text(s8, f"{level}: {label}", x, Inches(6.32), Inches(2.4), Inches(0.45), 10, True, text_color, PP_ALIGN.CENTER)
    if i < 4:
        add_text(s8, "\u25B6", Inches(2.45 + i * 2.6), Inches(6.35), Inches(0.3), Inches(0.3), 10, True, AZ_MAROON, PP_ALIGN.CENTER)

# ── BOTTOM: PRINCIPLES ──
add_text(s8, "GUIDING PRINCIPLES", Inches(0.4), Inches(6.95), Inches(4), Inches(0.25), 11, True, AZ_MAROON)

principles = [
    ("Augment, Don't Replace", "Informatica, SAP MDG, Collibra\nstay. Agentic layer on top.\nHuman oversight always.", AZ_MAROON),
    ("Start Small, Prove Value", "Follow Lilly/Novartis: prove\nROI on one domain first,\nthen scale. No big bang.", AZ_GREEN),
    ("Govern the Agents", "Agentic AI must itself be\ngoverned. Audit trails,\nexplainability, compliance.", AZ_BLUE),
    ("People & Change First", "60% of success is change\nmanagement. Stewards become\nstrategists. OSPs become partners.", AZ_DARK_PURPLE),
]

for i, (title, desc, color) in enumerate(principles):
    x = Inches(0.3 + i * 3.25)
    rrbox(s8, x, Inches(7.2), Inches(3.0), Inches(0.04), color)
    add_rich_text(s8, [
        (f"{title}: ", 9, True, color, False),
        (desc.replace('\n', ' '), 8, False, AZ_DARK, False),
    ], x, Inches(7.25), Inches(3.0), Inches(0.22))

# ============================================================================
# SAVE
# ============================================================================
output_path = "/Users/devesh.b.sharma/Astra Zeneca/R&D IT/AI-MDM-DataOps-DataGovernance.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print("Slides created:")
print("  1. Title — Agentic Data Management")
print("  2. Maturity Model — 5 Levels (Manual to Agentic)")
print("  3. 6-Phase Lifecycle — Discover, Trace, Monitor, Diagnose, Suggest, Fix")
print("  4. Platform Architecture — Layered Agentic Architecture (Capco)")
print("  5. Today & Tomorrow — Operating Model Shift")
print("  6. Case Studies — Eli Lilly & Novartis")
print("  7. Benefits Framework — 6 Pillars of Value")
print("  8. Implementation Roadmap — Mapping a Path Forward")
