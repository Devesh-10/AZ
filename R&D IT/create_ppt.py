from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# AstraZeneca brand colors
AZ_MAROON = RGBColor(0x83, 0x00, 0x51)
AZ_GOLD = RGBColor(0xC4, 0xA0, 0x00)
AZ_GREEN = RGBColor(0x2E, 0x7D, 0x32)
TEXT_DARK = RGBColor(0x2D, 0x2D, 0x2D)
TEXT_MEDIUM = RGBColor(0x55, 0x55, 0x55)
TEXT_LIGHT = RGBColor(0x77, 0x77, 0x77)
BG_CREAM = RGBColor(0xFA, 0xF9, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def set_cell_fill(cell, color):
    cell.fill.solid()
    cell.fill.fore_color.rgb = color

def set_cell_text(cell, text, font_size=12, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    cell.text = text
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(font_size)
    para.font.bold = bold
    para.font.color.rgb = color
    para.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_title_text(slide, text, left, top, width, height, font_size=38, bold=False, italic=True, color=TEXT_DARK):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.italic = italic
    p.font.bold = bold
    p.font.color.rgb = color
    p.text = text
    return shape

def add_text_box(slide, text, left, top, width, height, font_size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.text = text
    return shape

def add_metric_box(slide, value, label, left, top, width=Inches(2.2), height=Inches(1.2)):
    # Background box
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CREAM
    shape.line.color.rgb = AZ_MAROON
    shape.line.width = Pt(3)

    # Value text
    val_box = slide.shapes.add_textbox(left, top + Inches(0.15), width, Inches(0.6))
    tf = val_box.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = AZ_MAROON
    p.alignment = PP_ALIGN.CENTER

    # Label text
    label_box = slide.shapes.add_textbox(left, top + Inches(0.75), width, Inches(0.4))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MEDIUM
    p.alignment = PP_ALIGN.CENTER

def add_slide_number(slide, num):
    shape = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(0.5), Inches(0.3))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_LIGHT

def add_az_logo(slide):
    # Simple gold/maroon triangle logo representation
    left = Inches(12.5)
    top = Inches(6.8)
    # Gold outer triangle
    shape1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, Inches(0.6), Inches(0.6))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = AZ_GOLD
    shape1.line.fill.background()
    # Maroon inner triangle
    shape2 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left + Inches(0.1), top + Inches(0.15), Inches(0.4), Inches(0.4))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = AZ_MAROON
    shape2.line.fill.background()

# Create presentation (16:9 widescreen)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============ SLIDE 1: Title ============
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Main title
title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(10), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
run1 = p.add_run()
run1.text = "Agentic AI "
run1.font.size = Pt(44)
run1.font.italic = True
run1.font.color.rgb = TEXT_DARK
run2 = p.add_run()
run2.text = "MVP Programme"
run2.font.size = Pt(44)
run2.font.italic = True
run2.font.bold = True
run2.font.color.rgb = AZ_MAROON

# Subtitle
add_text_box(slide1, "A 9-month plan to build 4 AI agents for R&D IT platforms, starting with 3DP and BIKG",
             Inches(0.8), Inches(1.6), Inches(11), Inches(0.5), font_size=16, color=TEXT_MEDIUM)

# Metrics row
metrics = [("4", "AI Agents"), ("2-3", "Primary Platforms"), ("6", "FTEs"), ("9", "Months")]
for i, (val, label) in enumerate(metrics):
    add_metric_box(slide1, val, label, Inches(0.8 + i * 2.5), Inches(2.4))

# Meta information boxes
meta_items = [
    ("Duration", "March – November 2026"),
    ("Primary Platforms", "3DP, BIKG"),
    ("Agents", "Test Intelligence, Knowledge Library, Ingestion, Auto-Lineage"),
    ("Extensibility", "Additional platforms tailored during discovery")
]

for i, (label, value) in enumerate(meta_items):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 5.5)
    top = Inches(4.2 + row * 0.9)

    # Gold left border
    border = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.05), Inches(0.6))
    border.fill.solid()
    border.fill.fore_color.rgb = AZ_GOLD
    border.line.fill.background()

    # Label
    label_box = slide1.shapes.add_textbox(left + Inches(0.15), top, Inches(4), Inches(0.25))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label.upper()
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_LIGHT

    # Value
    val_box = slide1.shapes.add_textbox(left + Inches(0.15), top + Inches(0.25), Inches(4.5), Inches(0.35))
    tf = val_box.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK

add_slide_number(slide1, 1)
add_az_logo(slide1)

# ============ SLIDE 2: Team Structure ============
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10), Inches(0.7))
tf = title_box.text_frame
p = tf.paragraphs[0]
run1 = p.add_run()
run1.text = "Team "
run1.font.size = Pt(32)
run1.font.italic = True
run1.font.color.rgb = TEXT_DARK
run2 = p.add_run()
run2.text = "Structure"
run2.font.size = Pt(32)
run2.font.italic = True
run2.font.bold = True
run2.font.color.rgb = AZ_MAROON

add_text_box(slide2, "Core team of 6 FTEs with part-time platform SME support across the 9-month programme.",
             Inches(0.5), Inches(1.0), Inches(11), Inches(0.4), font_size=13, color=TEXT_MEDIUM)

# Team table
rows, cols = 5, 3
table = slide2.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(8), Inches(2.0)).table

# Header
headers = ["Role", "Count", "Name / Location"]
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    set_cell_fill(cell, AZ_MAROON)
    set_cell_text(cell, h, font_size=11, bold=True, color=WHITE)

# Data
team_data = [
    ("Tech Lead / Architect", "1", "Devesh"),
    ("Senior AI/ML Engineer", "2", "Abdalla Abdakarim + 1 Offshore"),
    ("Backend Engineer", "2", "Offshore"),
    ("Frontend Engineer / Tester", "1", "Offshore")
]

for row_idx, (role, count, name) in enumerate(team_data, 1):
    set_cell_text(table.cell(row_idx, 0), role, font_size=11, bold=True, color=AZ_MAROON)
    set_cell_text(table.cell(row_idx, 1), count, font_size=11)
    set_cell_text(table.cell(row_idx, 2), name, font_size=11)
    if row_idx % 2 == 0:
        for c in range(cols):
            set_cell_fill(table.cell(row_idx, c), BG_CREAM)

# Note box
note_shape = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.6), Inches(12), Inches(0.5))
note_shape.fill.solid()
note_shape.fill.fore_color.rgb = BG_CREAM
note_shape.line.fill.background()
# Gold left border
border = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.6), Inches(0.06), Inches(0.5))
border.fill.solid()
border.fill.fore_color.rgb = AZ_GOLD
border.line.fill.background()
add_text_box(slide2, "Part-time SMEs: 3DP (20%), BIKG (20%), AI Ops (10%), Data Foundation (10%) — additional platform SMEs engaged during discovery phase as needed.",
             Inches(0.65), Inches(3.7), Inches(11.5), Inches(0.4), font_size=11, color=TEXT_MEDIUM)

# Team Allocation section title
section_box = slide2.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(4), Inches(0.3))
tf = section_box.text_frame
p = tf.paragraphs[0]
p.text = "Team Allocation by Phase"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = AZ_MAROON

# Phase allocation table
rows2, cols2 = 7, 4
table2 = slide2.shapes.add_table(rows2, cols2, Inches(0.5), Inches(4.7), Inches(12), Inches(2.2)).table

phase_headers = ["Role", "Phase 1 (Mar-May)", "Phase 2 (Jun-Aug)", "Phase 3 (Sep-Nov)"]
phase_colors = [AZ_MAROON, AZ_MAROON, AZ_GOLD, AZ_GREEN]
for i, h in enumerate(phase_headers):
    cell = table2.cell(0, i)
    set_cell_fill(cell, phase_colors[i])
    text_color = TEXT_DARK if i == 2 else WHITE
    set_cell_text(cell, h, font_size=10, bold=True, color=text_color)

allocation_data = [
    ("Devesh", "Architecture, AI Ops integration", "RAG + Ingestion architecture", "Graph DB, lineage design"),
    ("Abdalla", "LLM test generation, prompts", "RAG pipeline, Q&A", "Lineage extraction"),
    ("ML Engineer", "Synthetic data generation", "Schema inference", "Impact analysis"),
    ("Backend Eng 1", "Test generation APIs", "Knowledge search APIs", "Lineage storage APIs"),
    ("Backend Eng 2", "Test runner service", "Ingestion mapping APIs", "Lineage query APIs"),
    ("Frontend / Tester", "Test dashboard, QA", "Chat UI, Mapping UI, QA", "Lineage graph UI, QA")
]

for row_idx, row_data in enumerate(allocation_data, 1):
    for col_idx, val in enumerate(row_data):
        cell = table2.cell(row_idx, col_idx)
        bold = col_idx == 0
        color = AZ_MAROON if col_idx == 0 else TEXT_DARK
        set_cell_text(cell, val, font_size=9, bold=bold, color=color)
        if row_idx % 2 == 0:
            set_cell_fill(cell, BG_CREAM)

add_slide_number(slide2, 2)
add_az_logo(slide2)

# ============ SLIDE 3: 4 Agents Overview ============
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10), Inches(0.7))
tf = title_box.text_frame
p = tf.paragraphs[0]
run1 = p.add_run()
run1.text = "4 Agents for "
run1.font.size = Pt(32)
run1.font.italic = True
run1.font.color.rgb = TEXT_DARK
run2 = p.add_run()
run2.text = "R&D IT Platforms"
run2.font.size = Pt(32)
run2.font.italic = True
run2.font.bold = True
run2.font.color.rgb = AZ_MAROON

add_text_box(slide3, "A portfolio of Agentic AI opportunities selected for high impact on primary platforms (3DP, BIKG), with flexibility to extend to additional platforms during discovery.",
             Inches(0.5), Inches(1.0), Inches(11), Inches(0.4), font_size=12, color=TEXT_MEDIUM)

# Agents table
rows3, cols3 = 5, 5
table3 = slide3.shapes.add_table(rows3, cols3, Inches(0.5), Inches(1.5), Inches(12), Inches(1.8)).table

agent_headers = ["Rank", "Agent", "Score", "Primary Platforms", "Process Bottleneck Addressed"]
for i, h in enumerate(agent_headers):
    cell = table3.cell(0, i)
    set_cell_fill(cell, AZ_MAROON)
    set_cell_text(cell, h, font_size=10, bold=True, color=WHITE)

agents_data = [
    ("#1", "Test Intelligence Agent", "4.30", "3DP, BIKG", "Slow Validation and Release Cycles"),
    ("#5", "Knowledge Library Companion", "4.00", "3DP, BIKG", "Lack of Reusable & Standardised Pipeline Components"),
    ("#6", "Ingestion Companion", "3.80", "3DP, BIKG", "Time-Consuming Data Ingestion and Preparation"),
    ("#8", "Auto-Lineage Agent", "3.75", "3DP, BIKG", "Manual Data Management and Traceability")
]

for row_idx, row_data in enumerate(agents_data, 1):
    for col_idx, val in enumerate(row_data):
        cell = table3.cell(row_idx, col_idx)
        bold = col_idx in [0, 1]
        set_cell_text(cell, val, font_size=10, bold=bold, color=TEXT_DARK)
        if row_idx % 2 == 0:
            set_cell_fill(cell, BG_CREAM)

# Platform tags
add_text_box(slide3, "Platform Strategy", Inches(0.5), Inches(3.45), Inches(3), Inches(0.3), font_size=12, bold=True, color=AZ_MAROON)

# Tags
tags = [("3DP (Primary)", AZ_MAROON, WHITE), ("BIKG (Primary)", AZ_MAROON, WHITE),
        ("Additional platforms identified during discovery", BG_CREAM, TEXT_MEDIUM)]
left = Inches(0.5)
for text, bg, fg in tags:
    tag = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.75), Inches(len(text) * 0.08 + 0.3), Inches(0.28))
    tag.fill.solid()
    tag.fill.fore_color.rgb = bg
    tag.line.fill.background()
    tag_text = slide3.shapes.add_textbox(left + Inches(0.1), Inches(3.78), Inches(len(text) * 0.08 + 0.2), Inches(0.25))
    tf = tag_text.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = fg
    left += Inches(len(text) * 0.08 + 0.5)

# Note box
note_shape = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.15), Inches(12), Inches(0.5))
note_shape.fill.solid()
note_shape.fill.fore_color.rgb = BG_CREAM
note_shape.line.fill.background()
border = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.15), Inches(0.06), Inches(0.5))
border.fill.solid()
border.fill.fore_color.rgb = AZ_GOLD
border.line.fill.background()
add_text_box(slide3, "Extensibility: All agents are designed with a platform-agnostic architecture. During discovery workshops, we will identify 1-2 additional platforms where these agents can be tailored and deployed.",
             Inches(0.65), Inches(4.25), Inches(11.5), Inches(0.4), font_size=10, color=TEXT_MEDIUM)

# Target Outcomes
add_text_box(slide3, "Target Outcomes", Inches(0.5), Inches(4.8), Inches(3), Inches(0.3), font_size=12, bold=True, color=AZ_MAROON)

outcomes = [("25%", "Less Manual Testing"), ("50%", "Faster Answers"), ("25-30%", "Faster Ingestion"), ("15-25%", "Less Manual Traceability")]
for i, (val, label) in enumerate(outcomes):
    add_metric_box(slide3, val, label, Inches(0.5 + i * 3.0), Inches(5.2), width=Inches(2.5), height=Inches(1.0))

add_slide_number(slide3, 3)
add_az_logo(slide3)

# ============ SLIDE 4: Agent Benefits ============
slide4 = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10), Inches(0.7))
tf = title_box.text_frame
p = tf.paragraphs[0]
run1 = p.add_run()
run1.text = "Agent Benefits: "
run1.font.size = Pt(32)
run1.font.italic = True
run1.font.color.rgb = TEXT_DARK
run2 = p.add_run()
run2.text = "3DP & BIKG"
run2.font.size = Pt(32)
run2.font.italic = True
run2.font.bold = True
run2.font.color.rgb = AZ_MAROON

add_text_box(slide4, "Detailed value proposition for the two primary pilot platforms.",
             Inches(0.5), Inches(1.0), Inches(11), Inches(0.4), font_size=13, color=TEXT_MEDIUM)

# Benefits table
rows4, cols4 = 12, 3
table4 = slide4.shapes.add_table(rows4, cols4, Inches(0.5), Inches(1.5), Inches(12), Inches(5.2)).table

benefits_headers = ["Agent", "3DP Benefit", "BIKG Benefit"]
for i, h in enumerate(benefits_headers):
    cell = table4.cell(0, i)
    set_cell_fill(cell, AZ_MAROON)
    set_cell_text(cell, h, font_size=11, bold=True, color=WHITE)

benefits_data = [
    ("Test Intelligence Agent", "Generate regulatory validation test cases", "Generate KG query test cases"),
    ("", "Synthetic patient data for submission testing", "Synthetic biological entity data"),
    ("", "Find similar past submission test failures", "Find similar failed BIKG queries"),
    ("Knowledge Library Companion", '"How do I prepare a Module 3 submission?"', '"How do I query gene-disease relationships?"'),
    ("", "Search regulatory guidelines, SOPs", "Search biological ontologies, data dictionaries"),
    ("", "Summarize 500-page FDA guidance docs", "Summarize pathway documentation"),
    ("Ingestion Companion", "Auto-detect regulatory data schemas", "Auto-detect biological data schemas"),
    ("", "Map source data to submission format", "Map source data to KG ontologies"),
    ("", "Improve ingestion prep by 25-30%", "Reduce manual data mapping"),
    ("Auto-Lineage Agent", "Regulatory audit trail (FDA/EMA requirement)", "Gene → Protein → Drug lineage"),
    ("", "FAIR compliance evidence", "Data provenance tracking")
]

for row_idx, row_data in enumerate(benefits_data, 1):
    for col_idx, val in enumerate(row_data):
        cell = table4.cell(row_idx, col_idx)
        if col_idx == 0:
            set_cell_fill(cell, BG_CREAM)
            set_cell_text(cell, val, font_size=10, bold=True, color=AZ_MAROON)
        else:
            set_cell_text(cell, val, font_size=10, color=TEXT_DARK)
            if row_idx % 2 == 0:
                set_cell_fill(cell, BG_CREAM)

add_slide_number(slide4, 4)
add_az_logo(slide4)

# ============ SLIDE 5: Three Phase Plan ============
slide5 = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
tf = title_box.text_frame
p = tf.paragraphs[0]
run1 = p.add_run()
run1.text = "Main steps to "
run1.font.size = Pt(30)
run1.font.italic = True
run1.font.color.rgb = TEXT_DARK
run2 = p.add_run()
run2.text = "mobilise and build"
run2.font.size = Pt(30)
run2.font.italic = True
run2.font.bold = True
run2.font.color.rgb = AZ_MAROON
run3 = p.add_run()
run3.text = " Agentic AI"
run3.font.size = Pt(30)
run3.font.italic = True
run3.font.color.rgb = TEXT_DARK

add_text_box(slide5, "A structured plan to build 4 AI agents, co-design Human+AI workflows, and enable change across R&D IT.",
             Inches(0.5), Inches(0.9), Inches(12), Inches(0.35), font_size=12, color=TEXT_MEDIUM)

# Three phase cards
phases = [
    ("March – May 2026", "Phase 1: Test Intelligence Agent", AZ_MAROON, WHITE,
     "Build foundational testing automation for 3DP and BIKG, validating the agentic approach with high-impact MVP.",
     ["Conduct discovery workshops with 3DP, BIKG teams", "Build LLM-powered test generation, synthetic data", "Integrate with 3DP and BIKG platforms"],
     ["Test Intelligence Agent deployed", "25% reduction in manual test effort", "Reusable prompts & templates"]),
    ("June – August 2026", "Phase 2: Knowledge Library + Ingestion", AZ_GOLD, TEXT_DARK,
     "Scale to knowledge management and data ingestion, enabling self-service docs and automated data onboarding.",
     ["Build RAG pipeline for document search and Q&A", "Build schema inference and auto-mapping", "Design human-in-the-loop review workflows"],
     ["Knowledge Library with docs indexed", "Ingestion Companion with auto-schema", "50% faster answers, 25-30% faster ingestion"]),
    ("September – November 2026", "Phase 3: Auto-Lineage Agent", AZ_GREEN, WHITE,
     "Complete the data lifecycle with lineage tracking, enabling FAIR compliance and regulatory audit readiness.",
     ["Define lineage requirements for regulatory and biological data", "Build graph-based lineage extraction", "Implement FAIR evidence generation"],
     ["Auto-Lineage Agent deployed", "FAIR compliance evidence for 3DP", "Training materials for scaling"])
]

card_width = Inches(3.9)
card_height = Inches(5.5)
card_gap = Inches(0.2)

for i, (dates, title, header_color, header_text_color, objective, activities, outputs) in enumerate(phases):
    left = Inches(0.5) + i * (card_width + card_gap)
    top = Inches(1.35)

    # Card border
    card = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, card_width, card_height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)

    # Header
    header = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, card_width, Inches(0.4))
    header.fill.solid()
    header.fill.fore_color.rgb = header_color
    header.line.fill.background()
    add_text_box(slide5, dates, left + Inches(0.1), top + Inches(0.08), card_width - Inches(0.2), Inches(0.3),
                 font_size=11, bold=True, color=header_text_color)

    # Subheader
    subheader = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.4), card_width, Inches(0.35))
    subheader.fill.solid()
    subheader.fill.fore_color.rgb = BG_CREAM
    subheader.line.fill.background()
    add_text_box(slide5, title, left + Inches(0.1), top + Inches(0.45), card_width - Inches(0.2), Inches(0.3),
                 font_size=11, bold=True, color=TEXT_DARK)

    # Content
    content_top = top + Inches(0.85)
    add_text_box(slide5, "OBJECTIVE", left + Inches(0.15), content_top, card_width - Inches(0.3), Inches(0.2),
                 font_size=9, bold=True, color=AZ_MAROON)
    add_text_box(slide5, objective, left + Inches(0.15), content_top + Inches(0.2), card_width - Inches(0.3), Inches(0.7),
                 font_size=9, color=TEXT_MEDIUM)

    add_text_box(slide5, "ACTIVITIES", left + Inches(0.15), content_top + Inches(0.95), card_width - Inches(0.3), Inches(0.2),
                 font_size=9, bold=True, color=AZ_MAROON)
    activities_text = "\n".join([f"• {a}" for a in activities])
    add_text_box(slide5, activities_text, left + Inches(0.15), content_top + Inches(1.15), card_width - Inches(0.3), Inches(1.1),
                 font_size=9, color=TEXT_MEDIUM)

    add_text_box(slide5, "OUTPUTS", left + Inches(0.15), content_top + Inches(2.4), card_width - Inches(0.3), Inches(0.2),
                 font_size=9, bold=True, color=AZ_MAROON)
    outputs_text = "\n".join([f"{j+1}. {o}" for j, o in enumerate(outputs)])
    add_text_box(slide5, outputs_text, left + Inches(0.15), content_top + Inches(2.6), card_width - Inches(0.3), Inches(1.0),
                 font_size=9, color=TEXT_DARK)

add_slide_number(slide5, 5)
add_az_logo(slide5)

# ============ SLIDE 6: Timeline ============
slide6 = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
tf = title_box.text_frame
p = tf.paragraphs[0]
run1 = p.add_run()
run1.text = "Proposed "
run1.font.size = Pt(30)
run1.font.italic = True
run1.font.color.rgb = TEXT_DARK
run2 = p.add_run()
run2.text = "timelines"
run2.font.size = Pt(30)
run2.font.italic = True
run2.font.bold = True
run2.font.color.rgb = AZ_MAROON
run3 = p.add_run()
run3.text = " for 9-month mobilisation"
run3.font.size = Pt(30)
run3.font.italic = True
run3.font.color.rgb = TEXT_DARK

# Timeline table
rows6, cols6 = 13, 10
table6 = slide6.shapes.add_table(rows6, cols6, Inches(0.3), Inches(0.9), Inches(12.7), Inches(5.3)).table

# Headers
timeline_headers = ["Workstream", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov"]
header_colors = [AZ_MAROON, AZ_MAROON, AZ_MAROON, AZ_MAROON, AZ_GOLD, AZ_GOLD, AZ_GOLD, AZ_GREEN, AZ_GREEN, AZ_GREEN]
for i, h in enumerate(timeline_headers):
    cell = table6.cell(0, i)
    set_cell_fill(cell, header_colors[i])
    text_color = TEXT_DARK if header_colors[i] == AZ_GOLD else WHITE
    set_cell_text(cell, h, font_size=9, bold=True, color=text_color, align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

# Timeline data - using █ for bars and ◆ for milestones
timeline_data = [
    ("Opportunity Selection & Discovery", ["█", "█◆", "", "", "", "", "", "", ""]),
    ("Phase 1: Test Intelligence Agent", ["█", "█", "█◆", "", "", "", "", "", ""]),
    ("  Discovery & Architecture", ["█", "", "", "", "", "", "", "", ""]),
    ("  LLM Test Gen & Synthetic Data", ["", "█", "", "", "", "", "", "", ""]),
    ("  Platform Integration & Pilot", ["", "", "█", "", "", "", "", "", ""]),
    ("Phase 2: Knowledge Library", ["", "", "", "█", "█", "█◆", "", "", ""]),
    ("Phase 2: Ingestion Companion", ["", "", "", "█", "█", "█◆", "", "", ""]),
    ("Phase 3: Auto-Lineage Agent", ["", "", "", "", "", "", "█", "█", "█◆"]),
    ("  Graph DB & Lineage Design", ["", "", "", "", "", "", "█", "", ""]),
    ("  Extraction & Visualization", ["", "", "", "", "", "", "", "█", ""]),
    ("  FAIR Compliance & Deployment", ["", "", "", "", "", "", "", "", "█"]),
    ("Training & Enablement", ["", "", "", "", "", "", "", "█", "█◆"])
]

bar_colors = [AZ_MAROON, AZ_MAROON, AZ_MAROON, AZ_MAROON, AZ_MAROON, AZ_GOLD, AZ_GOLD, AZ_GREEN, AZ_GREEN, AZ_GREEN, AZ_GREEN, AZ_GREEN]

for row_idx, (workstream, months) in enumerate(timeline_data, 1):
    # Workstream name
    cell = table6.cell(row_idx, 0)
    is_sub = workstream.startswith("  ")
    is_phase = workstream.startswith("Phase") or workstream == "Training & Enablement"
    set_cell_text(cell, workstream, font_size=8 if is_sub else 9, bold=is_phase,
                  color=TEXT_MEDIUM if is_sub else TEXT_DARK)
    if row_idx % 2 == 0:
        set_cell_fill(cell, BG_CREAM)

    # Month cells
    for col_idx, val in enumerate(months):
        cell = table6.cell(row_idx, col_idx + 1)
        if row_idx % 2 == 0:
            set_cell_fill(cell, BG_CREAM)
        if val:
            bar_color = bar_colors[row_idx - 1]
            display_val = val.replace("█", "■")
            set_cell_text(cell, display_val, font_size=12, bold=True, color=bar_color, align=PP_ALIGN.CENTER)

# Legend
legend_y = Inches(6.35)
add_text_box(slide6, "◆ Milestone", Inches(0.5), legend_y, Inches(1.2), Inches(0.25), font_size=9, color=AZ_GOLD)
add_text_box(slide6, "■ Phase 1: Test Intelligence", Inches(2.0), legend_y, Inches(2.2), Inches(0.25), font_size=9, color=AZ_MAROON)
add_text_box(slide6, "■ Phase 2: Knowledge + Ingestion", Inches(4.5), legend_y, Inches(2.6), Inches(0.25), font_size=9, color=AZ_GOLD)
add_text_box(slide6, "■ Phase 3: Auto-Lineage", Inches(7.5), legend_y, Inches(2.0), Inches(0.25), font_size=9, color=AZ_GREEN)

add_slide_number(slide6, 6)
add_az_logo(slide6)

# ============ SLIDE 7: Success Metrics ============
slide7 = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide7.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10), Inches(0.7))
tf = title_box.text_frame
p = tf.paragraphs[0]
run1 = p.add_run()
run1.text = "Success "
run1.font.size = Pt(32)
run1.font.italic = True
run1.font.color.rgb = TEXT_DARK
run2 = p.add_run()
run2.text = "Metrics"
run2.font.size = Pt(32)
run2.font.italic = True
run2.font.bold = True
run2.font.color.rgb = AZ_MAROON

add_text_box(slide7, "Key performance indicators for each agent MVP.",
             Inches(0.5), Inches(1.0), Inches(11), Inches(0.4), font_size=13, color=TEXT_MEDIUM)

# Metrics table
rows7, cols7 = 13, 4
table7 = slide7.shapes.add_table(rows7, cols7, Inches(0.5), Inches(1.5), Inches(12), Inches(5.2)).table

metrics_headers = ["Phase", "Agent", "Metric", "Target"]
for i, h in enumerate(metrics_headers):
    cell = table7.cell(0, i)
    set_cell_fill(cell, AZ_MAROON)
    set_cell_text(cell, h, font_size=11, bold=True, color=WHITE)

metrics_data = [
    ("P1", "Test Intelligence Agent", "Manual test effort reduction", "25%"),
    ("", "", "Test coverage increase", "15%"),
    ("", "", "Platforms integrated", "3DP, BIKG + discovery"),
    ("P2", "Knowledge Library Companion", "Time-to-answer reduction", "50%"),
    ("", "", "Documents indexed", "1000+"),
    ("", "", "Query success rate", "80%"),
    ("P2", "Ingestion Companion", "Ingestion prep time reduction", "25-30%"),
    ("", "", "Schema inference accuracy", "85%"),
    ("", "", "Platforms supported", "3DP, BIKG + discovery"),
    ("P3", "Auto-Lineage Agent", "Manual traceability reduction", "15-25%"),
    ("", "", "Lineage coverage", "3DP, BIKG + discovery"),
    ("", "", "FAIR compliance", "Audit-ready")
]

phase_badge_colors = {"P1": AZ_MAROON, "P2": AZ_GOLD, "P3": AZ_GREEN}
phase_badge_text = {"P1": WHITE, "P2": TEXT_DARK, "P3": WHITE}

for row_idx, (phase, agent, metric, target) in enumerate(metrics_data, 1):
    # Phase
    cell = table7.cell(row_idx, 0)
    if phase:
        set_cell_fill(cell, phase_badge_colors[phase])
        set_cell_text(cell, f"Phase {phase[1]}", font_size=9, bold=True, color=phase_badge_text[phase], align=PP_ALIGN.CENTER)
    elif row_idx % 2 == 0:
        set_cell_fill(cell, BG_CREAM)

    # Agent
    cell = table7.cell(row_idx, 1)
    set_cell_text(cell, agent, font_size=10, bold=bool(agent), color=TEXT_DARK)
    if row_idx % 2 == 0:
        set_cell_fill(cell, BG_CREAM)

    # Metric
    cell = table7.cell(row_idx, 2)
    set_cell_text(cell, metric, font_size=10, color=TEXT_DARK)
    if row_idx % 2 == 0:
        set_cell_fill(cell, BG_CREAM)

    # Target
    cell = table7.cell(row_idx, 3)
    set_cell_text(cell, target, font_size=10, bold=True, color=TEXT_DARK)
    if row_idx % 2 == 0:
        set_cell_fill(cell, BG_CREAM)

add_slide_number(slide7, 7)
add_az_logo(slide7)

# ============ SLIDE 8: Key Milestones ============
slide8 = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide8.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10), Inches(0.7))
tf = title_box.text_frame
p = tf.paragraphs[0]
run1 = p.add_run()
run1.text = "Key "
run1.font.size = Pt(32)
run1.font.italic = True
run1.font.color.rgb = TEXT_DARK
run2 = p.add_run()
run2.text = "Milestones"
run2.font.size = Pt(32)
run2.font.italic = True
run2.font.bold = True
run2.font.color.rgb = AZ_MAROON

add_text_box(slide8, "9-month delivery roadmap with clear deliverables for each phase.",
             Inches(0.5), Inches(1.0), Inches(11), Inches(0.4), font_size=13, color=TEXT_MEDIUM)

# Milestones table
rows8, cols8 = 10, 4
table8 = slide8.shapes.add_table(rows8, cols8, Inches(0.5), Inches(1.5), Inches(12), Inches(5.0)).table

milestones_headers = ["Month", "Phase", "Milestone", "Key Deliverable"]
for i, h in enumerate(milestones_headers):
    cell = table8.cell(0, i)
    set_cell_fill(cell, AZ_MAROON)
    set_cell_text(cell, h, font_size=11, bold=True, color=WHITE)

milestones_data = [
    ("M1 (Mar)", "P1", "Discovery Complete", "Requirements from 3DP, BIKG + additional platforms identified"),
    ("M2 (Apr)", "P1", "Test Gen API Live", "LLM test generation working with Claude"),
    ("M3 (May)", "P1", "◆ Agent 1 Pilot", "3DP + BIKG using Test Intelligence Agent"),
    ("M4 (Jun)", "P2", "RAG + Schema Ready", "Docs crawled, schema detection working"),
    ("M5 (Jul)", "P2", "Q&A + Mapping Live", "Knowledge search + ingestion mapping APIs"),
    ("M6 (Aug)", "P2", "◆ Agents 2 & 3 Pilot", "Knowledge Library + Ingestion Companion live"),
    ("M7 (Sep)", "P3", "Graph DB Ready", "Lineage schema deployed to Neptune/Neo4j"),
    ("M8 (Oct)", "P3", "Lineage Extraction", "Metadata captured across platforms"),
    ("M9 (Nov)", "P3", "◆ All 4 Agents Live", "FAIR compliance, full deployment, training complete")
]

highlight_rows = [3, 6, 9]  # M3, M6, M9

for row_idx, (month, phase, milestone, deliverable) in enumerate(milestones_data, 1):
    is_highlight = row_idx in highlight_rows
    bg_color = RGBColor(0xFF, 0xF8, 0xE1) if is_highlight else (BG_CREAM if row_idx % 2 == 0 else None)

    # Month
    cell = table8.cell(row_idx, 0)
    set_cell_text(cell, month, font_size=10, bold=True, color=TEXT_DARK)
    if bg_color:
        set_cell_fill(cell, bg_color)

    # Phase badge
    cell = table8.cell(row_idx, 1)
    set_cell_fill(cell, phase_badge_colors[phase])
    set_cell_text(cell, f"Phase {phase[1]}", font_size=9, bold=True, color=phase_badge_text[phase], align=PP_ALIGN.CENTER)

    # Milestone
    cell = table8.cell(row_idx, 2)
    set_cell_text(cell, milestone, font_size=10, bold=is_highlight, color=AZ_GOLD if "◆" in milestone else TEXT_DARK)
    if bg_color and not is_highlight:
        set_cell_fill(cell, bg_color)
    elif is_highlight:
        set_cell_fill(cell, RGBColor(0xFF, 0xF8, 0xE1))

    # Deliverable
    cell = table8.cell(row_idx, 3)
    set_cell_text(cell, deliverable, font_size=10, color=TEXT_DARK)
    if bg_color:
        set_cell_fill(cell, bg_color)

add_slide_number(slide8, 8)
add_az_logo(slide8)

# Save presentation
output_path = "/Users/devesh.b.sharma/Astra Zeneca/R&D IT/Agentic-AI-MVP-Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
